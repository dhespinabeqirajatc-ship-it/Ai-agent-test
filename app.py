import asyncio
import json
import os
import re
from copy import deepcopy
from datetime import datetime
from io import BytesIO
from xml.sax.saxutils import escape

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

from copilot import CopilotClient
from copilot.session import PermissionHandler
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from workiva import (
    discover_available_quarters,
    find_quarter_source_smart,
    get_sheet_data,
    workiva_sheet_to_dataframe,
)


load_dotenv()


# ============================================================
# PAGE
# ============================================================

st.set_page_config(
    page_title="Workiva AI Agent",
    layout="wide",
)

APP_NAME = "Workiva AI Agent"


# ============================================================
# BASIC HELPERS
# ============================================================

def now_text():

    return (
        datetime.now()
        .astimezone()
        .strftime(
            "%Y-%m-%d %H:%M:%S %Z"
        )
    )


def friendly_error(error):

    text = str(error)

    lower = text.lower()

    if "401" in text:

        return (
            "Workiva authentication failed. "
            "Check the Workiva credentials, "
            "region and API version."
        )

    if "403" in text:

        return (
            "The Workiva connection does not "
            "have permission to read this resource."
        )

    if "404" in text:

        return (
            "The requested Workiva spreadsheet "
            "or sheet could not be found."
        )

    if "429" in text:

        return (
            "Workiva is temporarily rate-limiting "
            "requests. Wait briefly and try again."
        )

    if (
        "timeout" in lower
        or "timed out" in lower
    ):

        return (
            "The Workiva request took too long. "
            "Try again in a moment."
        )

    if "github_token" in lower:

        return (
            "The AI connection is not configured. "
            "Check GITHUB_TOKEN in .env or "
            "Streamlit Secrets."
        )

    return (
        "The operation could not be completed. "
        f"Technical detail: {text}"
    )


# ============================================================
# SESSION STATE
# ============================================================

DEFAULTS = {

    "app_open":
        False,

    "agent_plan":
        None,

    "agent_request":
        (
            "Build an executive dashboard highlighting the most important "
            "KPIs, trends, risks, and business drivers in this dataset."
        ),

    "agent_request_input":
        (
            "Build an executive dashboard highlighting the most important "
            "KPIs, trends, risks, and business drivers in this dataset."
        ),

    "chat_messages":
        [],

    "comparison_result":
        None,

    "comparison_ai_commentary":
        None,

    "audit_log":
        [],

    "last_refreshed":
        None,

    "current_source_key":
        None,

    "pdf_bytes":
        None,

    "pptx_bytes":
        None,

    "management_summary":
        None,

    "management_summary_key":
        None,

    "forecast_result":
        None,

    "cross_quarter_quality":
        None,

    "cross_quarter_quality_key":
        None,

    "quality_ai_summary":
        None,

    "quality_ai_summary_key":
        None,

    "forecast_ai_commentary":
        None,

    "forecast_ai_commentary_key":
        None,

    "trends_ai_summary":
        None,

    "trends_ai_summary_key":
        None,

    "unified_management_brief":
        None,

    "unified_management_brief_key":
        None,

    "management_watchlist":
        [],
}


for key, value in DEFAULTS.items():

    if key not in st.session_state:

        st.session_state[
            key
        ] = value


def log_event(
    action,
    detail="",
):

    st.session_state[
        "audit_log"
    ].append(
        {
            "time":
                now_text(),

            "action":
                action,

            "detail":
                detail,
        }
    )


# ============================================================
# DATA CLEANING
# ============================================================

def clean_dataframe(data):

    data = (
        data.copy()
        .replace(
            "",
            pd.NA,
        )
    )

    data = (
        data
        .dropna(
            axis=0,
            how="all",
        )
        .dropna(
            axis=1,
            how="all",
        )
    )

    return data


# ============================================================
# DATA PROFILING
# ============================================================

def _parse_quarter_value(value, default_year=None, allow_bare_quarter=False):
    """Convert common quarter labels to the first day of the quarter.

    ``default_year`` lets standalone values such as Q1/Q2 become real
    datetimes when the surrounding Workiva sheet already supplies the year.
    """
    if value is None or pd.isna(value):
        return pd.NaT

    text = str(value).strip().upper()
    if not text:
        return pd.NaT

    # Normalize common finance variants first.
    text = re.sub(r"\bFISCAL\s+YEAR\b", "FY", text)
    text = re.sub(r"\bFISCAL\s+QUARTER\b", "Q", text)
    text = re.sub(r"\bQUARTER\b|\bQTR\b", "Q", text)
    text = re.sub(r"\s+", " ", text).strip()

    patterns = (
        r"^Q([1-4])[\s\-_/]*(?:FY)?(20\d{2})$",   # Q1 2026 / Q1 FY2026
        r"^(?:FY)?(20\d{2})[\s\-_/]*Q([1-4])$",   # 2026 Q1 / FY2026-Q1
        r"^([1-4])Q[\s\-_/]*(?:FY)?(20\d{2})$",   # 1Q 2026
    )

    for index, pattern in enumerate(patterns):
        match = re.match(pattern, text)
        if not match:
            continue

        if index == 1:
            year = int(match.group(1))
            quarter = int(match.group(2))
        else:
            quarter = int(match.group(1))
            year = int(match.group(2))

        month = (quarter - 1) * 3 + 1
        return pd.Timestamp(year=year, month=month, day=1)

    # Standalone Q1/Q2/Q3/Q4 can still be made chronological when the
    # caller knows the sheet/reporting year.
    standalone = re.fullmatch(r"Q([1-4])", text)
    if standalone is None and allow_bare_quarter:
        standalone = re.fullmatch(r"([1-4])", text)
    if standalone and default_year is not None:
        quarter = int(standalone.group(1))
        month = (quarter - 1) * 3 + 1
        return pd.Timestamp(year=int(default_year), month=month, day=1)

    return pd.NaT


def infer_date_series(series, default_year=None, allow_bare_quarter=False):
    """Infer dates from values rather than relying on the column heading.

    Returns a parsed datetime Series and a confidence ratio.  The parser
    recognises pandas datetime values, common quarter labels, year/month
    labels and ordinary date strings.  Pure numeric measures are deliberately
    protected from accidental timestamp conversion.
    """
    source = series.copy()
    parsed = pd.Series(pd.NaT, index=source.index, dtype="datetime64[ns]")
    non_null = source.dropna()

    if non_null.empty:
        return parsed, 0.0

    if pd.api.types.is_datetime64_any_dtype(source):
        parsed = pd.to_datetime(source, errors="coerce")
        return parsed, float(parsed.loc[non_null.index].notna().mean())

    text = source.astype("string").str.strip()
    present = source.notna() & text.ne("")

    # Quarter labels such as Q1 2026, 2026-Q1 and 1Q2026.
    quarter_values = source[present].map(
        lambda value: _parse_quarter_value(
            value,
            default_year=default_year,
            allow_bare_quarter=allow_bare_quarter,
        )
    )
    quarter_mask = quarter_values.notna()
    if quarter_mask.any():
        parsed.loc[quarter_values.index[quarter_mask]] = quarter_values[quarter_mask]

    # Do not let plain numbers (amounts, IDs, counts) become nanosecond dates.
    remaining = present & parsed.isna()
    remaining_text = text[remaining]
    numeric_text = remaining_text.str.fullmatch(r"[+-]?(?:\d+(?:\.\d*)?|\.\d+)")
    date_candidates = remaining_text[~numeric_text.fillna(False)]

    if not date_candidates.empty:
        try:
            generic_dates = pd.to_datetime(
                date_candidates,
                errors="coerce",
                format="mixed",
            )
        except (TypeError, ValueError):
            # Compatibility fallback for older pandas versions.
            generic_dates = pd.to_datetime(
                date_candidates,
                errors="coerce",
            )
        parsed.loc[generic_dates.index] = generic_dates

    # Four-digit years are date-like when the column values consistently look
    # like plausible calendar years.  This remains value-driven and therefore
    # works even for headings such as "Period" or "Fiscal".
    remaining = present & parsed.isna()
    year_text = text[remaining]
    year_mask = year_text.str.fullmatch(r"(?:19|20|21)\d{2}")
    if year_mask.any():
        year_candidates = year_text[year_mask]
        year_ratio = len(year_candidates) / max(int(present.sum()), 1)
        if year_ratio >= 0.8:
            parsed.loc[year_candidates.index] = pd.to_datetime(
                year_candidates + "-01-01",
                errors="coerce",
            )

    confidence = float(parsed[present].notna().mean()) if present.any() else 0.0
    return parsed, confidence


def profile_dataframe(data, default_year=None):

    data = (
        clean_dataframe(
            data
        )
    )

    profile = {

        "rows":
            len(data),

        "columns":
            len(
                data.columns
            ),

        "numeric":
            [],

        "categories":
            [],

        "dates":
            [],

        "identifiers":
            [],

        "details":
            {},
    }


    for column in data.columns:

        series = (
            data[column]
            .dropna()
        )

        if series.empty:

            continue


        name = str(
            column
        )

        lower = (
            name
            .strip()
            .lower()
        )


        looks_like_id = (

            lower == "id"

            or lower.endswith(
                "_id"
            )

            or lower.endswith(
                " id"
            )

            or lower.endswith(
                " code"
            )
        )


        numeric_version = (
            pd.to_numeric(
                series,
                errors="coerce",
            )
        )

        numeric_ratio = (
            numeric_version
            .notna()
            .mean()
        )

        _, date_ratio = infer_date_series(series, default_year=default_year)

        normalized_name = re.sub(r"[\s_\-]+", " ", lower).strip()
        is_quarter_heading = normalized_name in {
            "quarter",
            "fiscal quarter",
            "financial quarter",
        }
        quarter_text = (
            series.astype("string")
            .str.strip()
            .str.upper()
        )
        quarter_only_ratio = float(
            quarter_text
            .str.fullmatch(r"(?:Q|FQ|QTR\s*)[1-4]")
            .fillna(False)
            .mean()
        )
        # Quarter-shaped values are temporal even when the heading is called
        # something generic such as Period or Reporting Bucket. The heading
        # remains only a fallback, not the primary signal.
        bare_quarter_ratio = float(
            quarter_text
            .str.fullmatch(r"[1-4]")
            .fillna(False)
            .mean()
        )
        is_quarter_date = (
            quarter_only_ratio >= 0.8
            or (is_quarter_heading and bare_quarter_ratio >= 0.8)
            or (is_quarter_heading and date_ratio >= 0.5)
        )

        # A date interpretation takes precedence when the values themselves
        # strongly support it. The small heading fallback covers standalone
        # Q1/Q2/Q3/Q4 columns, which do not contain enough information to map
        # to an absolute timestamp. All other date detection is value-driven.
        if looks_like_id:

            kind = (
                "identifier"
            )

            profile[
                "identifiers"
            ].append(
                name
            )


        elif date_ratio >= 0.8 or is_quarter_date:

            kind = (
                "date"
            )

            profile[
                "dates"
            ].append(
                name
            )


        elif numeric_ratio >= 0.8:

            kind = (
                "numeric"
            )

            profile[
                "numeric"
            ].append(
                name
            )


        else:

            kind = (
                "category"
            )

            profile[
                "categories"
            ].append(
                name
            )


        profile[
            "details"
        ][name] = {

            "type":
                kind,

            "unique":
                int(
                    series.nunique()
                ),

            "missing":
                int(
                    data[column]
                    .isna()
                    .sum()
                ),

            "date_confidence":
                round(float(max(date_ratio, quarter_only_ratio if is_quarter_date else 0.0)), 3),
        }


    return profile


def extend_profile_with_derived_numeric(data, profile, derived_measures):
    """Add deterministic derived numeric columns to an existing profile.

    Derived finance measures are created by vectorized arithmetic, so they do not
    need a second full pass through date/type inference. Avoiding that re-profile
    materially reduces quarter-load CPU cost on wide Workiva sheets.
    """
    if not derived_measures:
        profile["rows"] = len(data)
        profile["columns"] = len(data.columns)
        return profile

    updated = deepcopy(profile)
    updated["rows"] = len(data)
    updated["columns"] = len(data.columns)

    for column in derived_measures:
        if column not in data.columns:
            continue
        name = str(column)
        if name not in updated["numeric"]:
            updated["numeric"].append(name)
        for bucket in ("categories", "dates", "identifiers"):
            if name in updated[bucket]:
                updated[bucket].remove(name)
        series = data[column]
        updated["details"][name] = {
            "type": "numeric",
            "unique": int(series.dropna().nunique()),
            "missing": int(series.isna().sum()),
            "date_confidence": 0.0,
        }

    return updated


def prepare_dataframe(
    data,
    profile,
    default_year=None,
):

    data = (
        clean_dataframe(
            data
        )
    )


    for column in profile[
        "numeric"
    ]:

        data[column] = (
            pd.to_numeric(
                data[column],
                errors="coerce",
            )
        )


    for column in profile[
        "dates"
    ]:

        try:
            parsed, confidence = infer_date_series(
                data[column],
                default_year=default_year,
                allow_bare_quarter=(
                    re.sub(r"[\s_\-]+", " ", str(column).strip().lower()).strip()
                    in {"quarter", "fiscal quarter", "financial quarter"}
                ),
            )
            if confidence >= 0.8:
                data[column] = parsed

        except Exception:

            pass


    return data


# ============================================================
# BUSINESS COLUMN DETECTION
# ============================================================

def find_column(
    columns,
    words,
):

    normalized = [

        (
            column,
            str(column)
            .strip()
            .lower(),
        )

        for column in columns
    ]


    for column, lower in normalized:

        if lower in words:

            return column


    for column, lower in normalized:

        if any(
            word in lower
            for word in words
        ):

            return column


    return None


def _find_financial_column(columns, aliases):
    """Prefer exact semantic matches, then safe substring matches."""
    normalized = [
        (column, str(column).strip().lower().replace("_", " ").replace("-", " "))
        for column in columns
    ]
    normalized_aliases = [
        str(alias).strip().lower().replace("_", " ").replace("-", " ")
        for alias in aliases
    ]

    for column, lower in normalized:
        if lower in normalized_aliases:
            return column

    # Avoid unsafe tiny substring aliases such as PY / LY, and do not let a
    # generated measure such as "Actual vs Budget Variance" become the
    # source Actual field merely because it contains the word "actual".
    safe_aliases = [alias for alias in normalized_aliases if len(alias) >= 4]
    unsafe_generic_substrings = {"actual", "actuals", "plan", "target", "income"}
    safe_aliases = [a for a in safe_aliases if a not in unsafe_generic_substrings]
    for column, lower in normalized:
        if any(alias in lower for alias in safe_aliases):
            return column

    return None


def detect_business_columns(data):
    """Detect a broader finance/operating vocabulary without using AI."""
    columns = list(data.columns)

    concepts = {
        "actual": ["actual", "actuals", "current actual"],
        "revenue": ["revenue", "sales", "net sales", "turnover", "income"],
        "budget": ["budget", "plan", "target", "budget amount"],
        "forecast": ["forecast", "latest estimate", "latest forecast", "outlook"],
        "costs": ["cost", "costs", "expense", "expenses", "total cost", "total costs"],
        "profit": ["profit", "gross profit", "net profit", "operating profit", "earnings"],
        "ebitda": ["ebitda", "adjusted ebitda"],
        "opex": ["opex", "operating expense", "operating expenses"],
        "capex": ["capex", "capital expenditure", "capital expenditures"],
        "headcount": ["headcount", "fte", "full time equivalent", "employees", "employee count"],
        "volume": ["volume", "units", "unit volume", "quantity", "qty"],
        "price": ["price", "unit price", "average price", "avg price", "asp"],
        "prior_year": [
            "prior year", "previous year", "last year", "prior year actual",
            "previous year actual", "prior year revenue", "ly actual", "py actual"
        ],
        "margin": ["margin", "gross margin", "operating margin", "net margin", "margin %"],
        "dimension": [
            "region", "country", "market", "department", "business unit",
            "segment", "division", "product", "product line", "customer group"
        ],
    }

    result = {
        key: _find_financial_column(columns, aliases)
        for key, aliases in concepts.items()
    }

    # Preserve the existing fallback behavior for common finance sheets.
    if result["budget"] is None and result["forecast"] is not None:
        result["budget"] = result["forecast"]
    if result["revenue"] is None and result["actual"] is not None:
        result["revenue"] = result["actual"]

    return result


def add_financial_intelligence(data, business):
    """Add deterministic derived finance measures. No Copilot call is used."""
    enriched = data.copy()
    derived = {}

    def numeric(column):
        if not column or column not in enriched.columns:
            return None
        return pd.to_numeric(enriched[column], errors="coerce")

    def add_column(name, values, formula):
        if values is None:
            return
        final_name = name
        suffix = 2
        while final_name in enriched.columns:
            final_name = f"{name} {suffix}"
            suffix += 1
        enriched[final_name] = values
        derived[final_name] = formula

    actual_col = business.get("actual") or business.get("revenue")
    budget_col = business.get("budget")
    revenue_col = business.get("revenue")
    costs_col = business.get("costs")
    profit_col = business.get("profit")
    ebitda_col = business.get("ebitda")
    opex_col = business.get("opex")
    capex_col = business.get("capex")
    prior_year_col = business.get("prior_year")

    actual = numeric(actual_col)
    budget = numeric(budget_col)
    revenue = numeric(revenue_col)
    costs = numeric(costs_col)
    profit = numeric(profit_col)
    ebitda = numeric(ebitda_col)
    opex = numeric(opex_col)
    capex = numeric(capex_col)
    prior_year = numeric(prior_year_col)

    if actual is not None and budget is not None and actual_col != budget_col:
        variance = actual - budget
        add_column(
            "Actual vs Budget Variance",
            variance,
            f"{actual_col} - {budget_col}",
        )
        denominator = budget.replace(0, pd.NA)
        add_column(
            "Actual vs Budget Variance %",
            variance / denominator * 100,
            f"({actual_col} - {budget_col}) / {budget_col} × 100",
        )

    if revenue is not None and profit is not None and revenue_col != profit_col:
        add_column(
            "Margin %",
            profit / revenue.replace(0, pd.NA) * 100,
            f"{profit_col} / {revenue_col} × 100",
        )

    if revenue is not None and costs is not None and revenue_col != costs_col:
        add_column(
            "Cost Ratio %",
            costs / revenue.replace(0, pd.NA) * 100,
            f"{costs_col} / {revenue_col} × 100",
        )

    if revenue is not None and ebitda is not None and revenue_col != ebitda_col:
        add_column(
            "EBITDA Margin %",
            ebitda / revenue.replace(0, pd.NA) * 100,
            f"{ebitda_col} / {revenue_col} × 100",
        )

    if revenue is not None and opex is not None and revenue_col != opex_col:
        add_column(
            "Opex Ratio %",
            opex / revenue.replace(0, pd.NA) * 100,
            f"{opex_col} / {revenue_col} × 100",
        )

    if revenue is not None and capex is not None and revenue_col != capex_col:
        add_column(
            "Capex Ratio %",
            capex / revenue.replace(0, pd.NA) * 100,
            f"{capex_col} / {revenue_col} × 100",
        )

    if actual is not None and prior_year is not None and actual_col != prior_year_col:
        yoy = actual - prior_year
        add_column(
            "Prior Year Variance",
            yoy,
            f"{actual_col} - {prior_year_col}",
        )
        add_column(
            "Prior Year Variance %",
            yoy / prior_year.replace(0, pd.NA) * 100,
            f"({actual_col} - {prior_year_col}) / {prior_year_col} × 100",
        )

    return enriched, derived


# ============================================================
# ANALYTICAL CONTEXT
# ============================================================

def _is_additive_driver_metric(metric):
    """Return True for measures where category contributions can be summed safely."""
    lower = str(metric).strip().lower()
    non_additive_tokens = (
        "%", "ratio", "margin", "rate", "average", "avg", "per ", "share",
    )
    return not any(token in lower for token in non_additive_tokens)


def build_driver_analysis(data, profile, business):
    """Deterministically explain what categories drove change between two periods.

    The analysis intentionally uses arithmetic only.  It compares the earliest and
    latest observed dates, decomposes additive metric changes by low-cardinality
    dimensions, and ranks contributors by absolute impact.
    """
    if not profile.get("dates") or not profile.get("numeric"):
        return {}

    date_col = profile["dates"][0]
    working_dates = pd.to_datetime(data[date_col], errors="coerce")
    valid_dates = working_dates.dropna().sort_values().unique()
    if len(valid_dates) < 2:
        return {}

    first_date = pd.Timestamp(valid_dates[0])
    last_date = pd.Timestamp(valid_dates[-1])

    # Put the detected management dimension first, then other usable categories.
    category_candidates = []
    preferred = business.get("dimension")
    if preferred in data.columns:
        category_candidates.append(preferred)
    for category in profile.get("categories", []):
        if category not in category_candidates and category in data.columns:
            category_candidates.append(category)

    category_candidates = [
        category for category in category_candidates
        if 2 <= data[category].nunique(dropna=True) <= 40
    ][:4]

    metric_candidates = [
        metric for metric in profile.get("numeric", [])
        if metric in data.columns and _is_additive_driver_metric(metric)
    ][:8]

    if not category_candidates or not metric_candidates:
        return {}

    result = {
        "date_column": str(date_col),
        "first_period": str(first_date),
        "last_period": str(last_date),
        "metrics": {},
    }

    for metric in metric_candidates:
        metric_values = pd.to_numeric(data[metric], errors="coerce")
        metric_result = {}

        for category in category_candidates:
            frame = pd.DataFrame({
                "__date": working_dates,
                "__category": data[category].astype("string"),
                "__value": metric_values,
            }).dropna(subset=["__date", "__category", "__value"])

            frame = frame[frame["__date"].isin([first_date, last_date])]
            if frame.empty:
                continue

            grouped = (
                frame.groupby(["__date", "__category"], dropna=False)["__value"]
                .sum()
                .unstack(fill_value=0.0)
            )
            if first_date not in grouped.index or last_date not in grouped.index:
                continue

            first_values = grouped.loc[first_date]
            last_values = grouped.loc[last_date]
            all_categories = first_values.index.union(last_values.index)
            first_values = first_values.reindex(all_categories, fill_value=0.0)
            last_values = last_values.reindex(all_categories, fill_value=0.0)
            changes = (last_values - first_values).sort_values(
                key=lambda x: x.abs(), ascending=False
            )

            total_first = float(first_values.sum())
            total_last = float(last_values.sum())
            total_change = float(total_last - total_first)
            if changes.empty or abs(total_change) < 1e-12:
                continue

            drivers = []
            for category_value, change in changes.head(8).items():
                change = float(change)
                drivers.append({
                    "category": str(category_value),
                    "first_value": float(first_values.get(category_value, 0.0)),
                    "last_value": float(last_values.get(category_value, 0.0)),
                    "change": change,
                    "contribution_pct_of_net_change": float(change / total_change * 100),
                    "direction": "increase" if change > 0 else "decrease",
                })

            metric_result[str(category)] = {
                "total_first": total_first,
                "total_last": total_last,
                "total_change": total_change,
                "total_change_pct": (
                    float(total_change / total_first * 100) if total_first else None
                ),
                "drivers": drivers,
            }

        if metric_result:
            result["metrics"][str(metric)] = metric_result

    return result if result["metrics"] else {}


def build_analytical_context(data, profile, business):
    """Pre-calculate trustworthy facts for the AI summary and chat."""

    context = {
        "metric_totals": {},
        "top_breakdowns": {},
        "trends": {},
        "correlations": [],
        "anomalies": {},
        "driver_analysis": {},
        "financial_variance": {},
        "data_quality": {},
    }

    row_count = max(len(data), 1)
    missing_by_column = data.isna().sum().sort_values(ascending=False)
    context["data_quality"] = {
        "missing_cells": int(data.isna().sum().sum()),
        "missing_rate_pct": float(data.isna().sum().sum() / max(data.size, 1) * 100),
        "most_missing_columns": [
            {"column": str(col), "missing": int(value), "missing_pct": float(value / row_count * 100)}
            for col, value in missing_by_column.head(5).items()
            if value > 0
        ],
    }

    for metric in profile["numeric"][:12]:
        values = pd.to_numeric(data[metric], errors="coerce").dropna()
        if values.empty:
            continue
        context["metric_totals"][metric] = {
            "sum": float(values.sum()),
            "average": float(values.mean()),
            "median": float(values.median()),
            "minimum": float(values.min()),
            "maximum": float(values.max()),
            "non_null_count": int(values.count()),
        }

        # Deterministic IQR outlier screening gives AI an evidence-backed
        # anomaly signal without asking the model to invent thresholds.
        if len(values) >= 4:
            q1 = float(values.quantile(0.25))
            q3 = float(values.quantile(0.75))
            iqr = q3 - q1
            if iqr > 0:
                lower_bound = q1 - 1.5 * iqr
                upper_bound = q3 + 1.5 * iqr
                outliers = values[(values < lower_bound) | (values > upper_bound)]
                if not outliers.empty:
                    context["anomalies"][metric] = {
                        "method": "IQR 1.5x",
                        "count": int(outliers.count()),
                        "lower_bound": float(lower_bound),
                        "upper_bound": float(upper_bound),
                        "minimum_outlier": float(outliers.min()),
                        "maximum_outlier": float(outliers.max()),
                    }

    for category in profile["categories"][:6]:
        unique = data[category].nunique(dropna=True)
        if unique < 2 or unique > 40:
            continue
        for metric in profile["numeric"][:6]:
            grouped = grouped_data(data, category, metric, "sum").dropna()
            if grouped.empty:
                continue
            grouped = grouped.sort_values(metric, ascending=False)
            key = f"{metric} by {category}"
            context["top_breakdowns"][key] = {
                "top": [
                    {"category": str(row[category]), "value": float(row[metric])}
                    for _, row in grouped.head(5).iterrows()
                ],
                "bottom": [
                    {"category": str(row[category]), "value": float(row[metric])}
                    for _, row in grouped.tail(3).iterrows()
                ],
            }

    if profile["dates"]:
        date_col = profile["dates"][0]
        for metric in profile["numeric"][:8]:
            working = data[[date_col, metric]].copy()
            working[date_col] = pd.to_datetime(working[date_col], errors="coerce")
            working[metric] = pd.to_numeric(working[metric], errors="coerce")
            working = working.dropna()
            if working.empty:
                continue
            grouped = working.groupby(date_col)[metric].sum().sort_index()
            if len(grouped) < 2:
                continue
            first = float(grouped.iloc[0])
            last = float(grouped.iloc[-1])
            change = last - first
            differences = grouped.diff().dropna()
            nonzero_differences = differences[differences.abs() > 1e-12]
            if nonzero_differences.empty:
                direction_consistency_pct = 100.0
            else:
                positive_share = float((nonzero_differences > 0).mean())
                negative_share = float((nonzero_differences < 0).mean())
                direction_consistency_pct = max(positive_share, negative_share) * 100

            context["trends"][metric] = {
                "date_column": str(date_col),
                "first_date": str(grouped.index[0]),
                "last_date": str(grouped.index[-1]),
                "first_value": first,
                "last_value": last,
                "change": float(change),
                "change_pct": float(change / first * 100) if first else None,
                "peak_date": str(grouped.idxmax()),
                "peak_value": float(grouped.max()),
                "low_date": str(grouped.idxmin()),
                "low_value": float(grouped.min()),
                "observation_count": int(len(grouped)),
                "direction_consistency_pct": float(direction_consistency_pct),
            }

    numeric_frame = data[profile["numeric"][:10]].apply(pd.to_numeric, errors="coerce")
    if numeric_frame.shape[1] >= 2:
        corr = numeric_frame.corr()
        pairs = []
        columns = list(corr.columns)
        for i, first in enumerate(columns):
            for second in columns[i + 1:]:
                value = corr.loc[first, second]
                if pd.notna(value):
                    pairs.append((abs(float(value)), first, second, float(value)))
        for _, first, second, value in sorted(pairs, reverse=True)[:5]:
            context["correlations"].append({
                "metric_1": str(first),
                "metric_2": str(second),
                "correlation": value,
            })

    revenue = business.get("revenue")
    budget = business.get("budget")
    costs = business.get("costs")
    if revenue and budget:
        actual = pd.to_numeric(data[revenue], errors="coerce").sum()
        plan = pd.to_numeric(data[budget], errors="coerce").sum()
        variance = actual - plan
        context["financial_variance"]["actual_vs_budget"] = {
            "actual_metric": str(revenue),
            "budget_metric": str(budget),
            "actual": float(actual),
            "budget": float(plan),
            "variance": float(variance),
            "variance_pct": float(variance / plan * 100) if plan else None,
        }
    if revenue and costs:
        rev = pd.to_numeric(data[revenue], errors="coerce").sum()
        cost = pd.to_numeric(data[costs], errors="coerce").sum()
        context["financial_variance"]["revenue_vs_costs"] = {
            "revenue": float(rev),
            "costs": float(cost),
            "spread": float(rev - cost),
            "spread_pct_of_revenue": float((rev - cost) / rev * 100) if rev else None,
        }

    context["driver_analysis"] = build_driver_analysis(
        data,
        profile,
        business,
    )

    return context


def build_question_context(question, data, profile, business):
    """Calculate extra facts specifically relevant to the user's question."""

    lower = question.lower()
    metric_candidates = [
        c for c in profile["numeric"]
        if str(c).lower() in lower
    ]
    category_candidates = [
        c for c in profile["categories"]
        if str(c).lower() in lower
    ]

    keyword_map = {
        "revenue": business.get("revenue"),
        "sales": business.get("revenue"),
        "budget": business.get("budget"),
        "forecast": business.get("budget"),
        "cost": business.get("costs"),
        "expense": business.get("costs"),
        "profit": business.get("profit"),
        "ebitda": business.get("ebitda"),
        "opex": business.get("opex"),
        "capex": business.get("capex"),
        "headcount": business.get("headcount"),
        "fte": business.get("headcount"),
        "volume": business.get("volume"),
        "price": business.get("price"),
        "margin": business.get("margin"),
        "prior year": business.get("prior_year"),
    }
    for word, column in keyword_map.items():
        if word in lower and column and column not in metric_candidates:
            metric_candidates.append(column)

    if not metric_candidates:
        metric_candidates = profile["numeric"][:3]
    if not category_candidates and business.get("dimension"):
        category_candidates = [business["dimension"]]

    result = {"metrics": {}, "breakdowns": {}, "sample_rows": []}
    for metric in metric_candidates[:4]:
        values = pd.to_numeric(data[metric], errors="coerce").dropna()
        if not values.empty:
            result["metrics"][metric] = {
                "sum": float(values.sum()),
                "average": float(values.mean()),
                "median": float(values.median()),
                "minimum": float(values.min()),
                "maximum": float(values.max()),
            }

    for category in category_candidates[:2]:
        if category not in data.columns or data[category].nunique(dropna=True) > 50:
            continue
        for metric in metric_candidates[:3]:
            grouped = grouped_data(data, category, metric, "sum").sort_values(metric, ascending=False)
            result["breakdowns"][f"{metric} by {category}"] = grouped.head(12).to_dict(orient="records")

    matched_columns = [c for c in data.columns if str(c).lower() in lower]
    if matched_columns:
        result["sample_rows"] = data[matched_columns[:6]].head(8).astype(str).to_dict(orient="records")

    return result


# ============================================================
# RICH DATA SUMMARY
# ============================================================

def build_data_summary(
    data,
    profile,
    quarter,
):

    summary = {

        "quarter":
            quarter,

        "rows":
            len(data),

        "columns":
            list(
                data.columns
            ),

        "numeric_columns":
            profile[
                "numeric"
            ],

        "category_columns":
            profile[
                "categories"
            ],

        "date_columns":
            profile[
                "dates"
            ],

        "numeric_stats":
            {},

        "categories":
            {},

        "date_ranges":
            {},

        "missing_cells":
            int(
                data
                .isna()
                .sum()
                .sum()
            ),
    }


    for column in profile[
        "numeric"
    ][:20]:

        values = (
            pd.to_numeric(
                data[column],
                errors="coerce",
            )
            .dropna()
        )


        if values.empty:

            continue


        summary[
            "numeric_stats"
        ][column] = {

            "sum":
                float(
                    values.sum()
                ),

            "average":
                float(
                    values.mean()
                ),

            "minimum":
                float(
                    values.min()
                ),

            "maximum":
                float(
                    values.max()
                ),

            "count":
                int(
                    values.count()
                ),
        }


    for column in profile[
        "categories"
    ][:10]:

        values = (
            data[column]
            .dropna()
            .astype(str)
        )


        summary[
            "categories"
        ][column] = {

            "unique":
                int(
                    values.nunique()
                ),

            "top_values":
                (
                    values
                    .value_counts()
                    .head(10)
                    .to_dict()
                ),
        }


    for column in profile[
        "dates"
    ]:

        values = (
            data[column]
            .dropna()
        )


        if not values.empty:

            summary[
                "date_ranges"
            ][column] = {

                "minimum":
                    str(
                        values.min()
                    ),

                "maximum":
                    str(
                        values.max()
                    ),
            }


    return summary


# ============================================================
# COPILOT
# ============================================================

async def ask_copilot(prompt):

    token = (
        os.environ.get(
            "GITHUB_TOKEN"
        )
    )


    if not token:

        raise ValueError(
            "GITHUB_TOKEN is missing."
        )


    client = CopilotClient(

        github_token=token,

        use_logged_in_user=False,
    )


    await client.start()


    try:

        session = (
            await client.create_session(

                model="auto",

                on_permission_request=(
                    PermissionHandler
                    .approve_all
                ),
            )
        )


        response = (
            await session
            .send_and_wait(
                prompt
            )
        )


        if response:

            return (
                response
                .data
                .content
            )


        return ""


    finally:

        await client.stop()


def run_copilot(prompt):

    return asyncio.run(
        ask_copilot(
            prompt
        )
    )


# ============================================================
# GROUNDED AI SAFETY
# ============================================================

def _numeric_tokens(text):
    """Extract user-visible numeric claims while ignoring list numbering."""
    text = str(text or "")
    pattern = re.compile(r"(?<![A-Za-z])[-+]?\d[\d,]*(?:\.\d+)?(?:%|x)?")
    values = []
    for match in pattern.finditer(text):
        token = match.group(0)
        start, end = match.span()
        before = text[max(0, start - 2):start]
        after = text[end:end + 1]
        # Ignore markdown/section list markers such as "1." at line starts.
        line_start = text.rfind("\n", 0, start) + 1
        prefix = text[line_start:start].strip()
        if prefix == "" and after == "." and token.lstrip("+-").isdigit():
            continue
        cleaned = token.rstrip("%x").replace(",", "")
        try:
            values.append((token, float(cleaned)))
        except ValueError:
            pass
    return values


def _allowed_evidence_numbers(evidence):
    """Build a tolerant set of numbers that are actually present in evidence."""
    serialized = json.dumps(evidence, default=str, ensure_ascii=False)
    raw = [value for _, value in _numeric_tokens(serialized)]
    allowed = set()
    for value in raw:
        # Permit normal presentation rounding, but not newly invented values.
        for decimals in range(0, 7):
            allowed.add(round(float(value), decimals))
    return allowed


def _unsupported_ai_numbers(text, evidence=None, allowed=None):
    # grounded_ai_response precomputes the allowed numeric universe once so a
    # repair pass does not repeatedly serialize and scan the same evidence.
    if allowed is None:
        allowed = _allowed_evidence_numbers(evidence)
    unsupported = []
    for token, value in _numeric_tokens(text):
        if not any(abs(value - candidate) <= 1e-9 for candidate in allowed):
            unsupported.append(token)
    return list(dict.fromkeys(unsupported))


@st.cache_data(ttl=900, show_spinner=False)
def grounded_ai_response(prompt, evidence, purpose="analysis"):
    """Use Copilot for interpretation while Python remains the source of truth.

    The model receives a bounded evidence package. Its numeric claims are checked
    after generation; if unsupported numbers appear, it gets one constrained
    rewrite. If the rewrite still introduces unsupported figures, the response is
    withheld rather than displaying an ungrounded numerical statement.
    """
    # Compact JSON reduces AI request size/latency without removing evidence.
    evidence_text = json.dumps(
        evidence,
        default=str,
        ensure_ascii=False,
        separators=(",", ":"),
    )
    allowed_numbers = _allowed_evidence_numbers(evidence)
    grounding_rules = f"""

GROUND-TRUTH EVIDENCE (calculated/retrieved by the application):
{evidence_text}

MANDATORY GROUNDING RULES:
- Treat the evidence above as the complete factual universe for this {purpose}.
- Python/Workiva values are authoritative. Never calculate or invent an official value yourself.
- Do not introduce a number unless that exact value, or an ordinary rounded display of it, is present in the evidence.
- Never invent causes, events, targets, benchmarks, seasonality, probabilities, confidence levels, or external drivers.
- You may interpret patterns that are directly supported by the evidence, but label uncertain interpretations as such.
- If the evidence cannot answer something, explicitly say the data does not establish it.
- Recommendations must be framed as management questions/actions to investigate, not as claims about facts not in evidence.
- Do not use numbered section headings; use plain markdown headings or bullets.
"""
    full_prompt = prompt + grounding_rules
    response = run_copilot(full_prompt)
    unsupported = _unsupported_ai_numbers(response, allowed=allowed_numbers)

    if unsupported:
        repair_prompt = f"""
{full_prompt}

Your previous draft contained unsupported numeric claims: {unsupported}
Previous draft:
{response}

Rewrite the answer now. Remove every unsupported number and any claim that depends on it.
Use only evidence-supported facts. Return only the corrected answer.
"""
        response = run_copilot(repair_prompt)
        unsupported = _unsupported_ai_numbers(response, allowed=allowed_numbers)

    if unsupported:
        return (
            "AI commentary was withheld because its generated text still contained "
            "numeric claims that could not be matched to the calculated dataset. "
            "The dashboard values and deterministic analysis remain available above."
        )

    return response




# ============================================================
# AI EVIDENCE / CONFIDENCE SURFACING
# ============================================================

def build_ai_evidence_indicator(summary, purpose="analysis", forecast_result=None, comparison_result=None):
    """Return a deterministic evidence-strength label for an AI interpretation.

    This is deliberately separate from the grounding validator. Grounding checks
    whether numerical claims are present in the supplied evidence; this indicator
    communicates how much analytical support exists for the interpretation itself.
    It is an application heuristic, not a statistical confidence interval.
    """
    quality = (summary or {}).get("data_quality_intelligence", {}) or {}
    quality_score = quality.get("score")
    try:
        score = float(quality_score) if quality_score is not None else 65.0
    except (TypeError, ValueError):
        score = 65.0

    basis = [
        "Workiva/source values are treated as the source of record.",
        "Python performs the calculations before AI interpretation.",
        "AI numeric claims are checked against the supplied evidence and unsupported figures are rejected.",
    ]

    purpose_lower = str(purpose or "analysis").lower()
    if "forecast" in purpose_lower:
        backtests = (forecast_result or {}).get("backtests", {}) or {}
        selected_metrics = (forecast_result or {}).get("metrics", []) or []
        labels = [
            str((backtests.get(metric, {}) or {}).get("confidence"))
            for metric in selected_metrics
            if (backtests.get(metric, {}) or {}).get("confidence")
        ]
        if labels:
            # Forecast confidence is intentionally capped elsewhere because the
            # history is short. Reflect the weakest selected metric here.
            if "Limited" in labels:
                score = min(score, 45.0)
                basis.append("At least one selected forecast metric has Limited walk-forward confidence.")
            else:
                score = min(score, 75.0)
                basis.append("Selected forecast metrics have Moderate walk-forward confidence; short history prevents a High label.")
        else:
            score = min(score, 45.0)
            basis.append("No usable forecast backtest confidence is available for the selected forecast state.")

    if comparison_result:
        basis.append("Quarter-comparison statements use calculated common-column and variance evidence from the selected periods.")

    score = max(0.0, min(100.0, score))
    if score >= 85:
        label = "Strong evidence"
    elif score >= 65:
        label = "Supported interpretation"
    else:
        label = "Limited evidence"

    return {
        "label": label,
        "score": round(score),
        "basis": basis,
        "quality_status": quality.get("status"),
        "quality_score": quality_score,
        "purpose": purpose,
    }


def render_ai_evidence_indicator(summary, purpose="analysis", forecast_result=None, comparison_result=None, key=None, compact=False):
    """Render a small trust indicator without crowding the AI narrative."""
    indicator = build_ai_evidence_indicator(
        summary,
        purpose=purpose,
        forecast_result=forecast_result,
        comparison_result=comparison_result,
    )
    st.caption(
        f"Evidence: **{indicator['label']}** · {indicator['score']}/100 · "
        "Grounded AI interpretation"
    )
    if not compact:
        expander_label = "Why this evidence rating?"
        with st.expander(expander_label, expanded=False):
            st.caption(
                "This rating describes analytical support for the AI interpretation. "
                "It is not a statistical confidence interval, audit opinion, or accounting materiality assessment."
            )
            for item in indicator["basis"]:
                st.markdown(f"- {item}")
            if indicator.get("quality_status"):
                quality_text = f"Underlying analysis readiness: {indicator['quality_status']}"
                if indicator.get("quality_score") is not None:
                    quality_text += f" · {float(indicator['quality_score']):.0f}/100"
                st.markdown(f"- {quality_text}")
    return indicator


# ============================================================
# JSON EXTRACTION
# ============================================================

def extract_json(text):

    start = (
        text.find(
            "{"
        )
    )

    end = (
        text.rfind(
            "}"
        )
    )


    if (
        start == -1
        or end == -1
        or end < start
    ):

        raise ValueError(
            "AI did not return "
            "a valid dashboard plan."
        )


    return json.loads(
        text[
            start:
            end + 1
        ]
    )


# ============================================================
# COMPACT AI CONTEXT
# ============================================================

def compact_ai_context(summary):
    """Return a smaller, decision-useful context for Copilot prompts.

    The full analytical context remains available in Python for the UI and
    exports.  Copilot only receives the most material subset, which reduces
    prompt size and model latency without adding any extra API calls.
    """
    context = summary.get("analytical_context", {}) or {}

    compact = {
        "financial_variance": context.get("financial_variance", {}),
        "data_quality": context.get("data_quality", {}),
        "quality_intelligence": context.get("quality_intelligence", {}),
        "trends": {},
        "top_breakdowns": {},
        "correlations": (context.get("correlations", []) or [])[:3],
        "anomalies": dict(list((context.get("anomalies", {}) or {}).items())[:5]),
        "driver_analysis": {},
        "materiality_ranking": [],
        "management_watchlist": list(summary.get("management_watchlist", []) or []),
        "watchlist_evidence": {},
    }

    # Keep only the strongest few trend signals.
    trend_rows = []
    for metric, details in (context.get("trends", {}) or {}).items():
        pct = details.get("change_pct")
        score = abs(float(pct)) if pct is not None else 0.0
        trend_rows.append((score, metric, details))
    for _, metric, details in sorted(trend_rows, reverse=True)[:5]:
        compact["trends"][metric] = details

    # Keep driver evidence compact: strongest few metrics/dimensions and top contributors.
    drivers = context.get("driver_analysis", {}) or {}
    if drivers.get("metrics"):
        compact["driver_analysis"] = {
            "date_column": drivers.get("date_column"),
            "first_period": drivers.get("first_period"),
            "last_period": drivers.get("last_period"),
            "metrics": {},
        }
        for metric, dimensions in list(drivers.get("metrics", {}).items())[:4]:
            compact["driver_analysis"]["metrics"][metric] = {}
            for dimension, details in list(dimensions.items())[:2]:
                compact["driver_analysis"]["metrics"][metric][dimension] = {
                    "total_first": details.get("total_first"),
                    "total_last": details.get("total_last"),
                    "total_change": details.get("total_change"),
                    "total_change_pct": details.get("total_change_pct"),
                    "drivers": (details.get("drivers", []) or [])[:5],
                }

    # Management watchlist is preference context, never a source of factual values.
    # Attach only already-calculated evidence for bookmarked metrics.
    watchlist = list(summary.get("management_watchlist", []) or [])
    numeric_stats = summary.get("numeric_stats", {}) or {}
    for metric in watchlist[:5]:
        compact["watchlist_evidence"][metric] = {
            "numeric_statistics": numeric_stats.get(metric),
            "trend": (context.get("trends", {}) or {}).get(metric),
            "anomalies": (context.get("anomalies", {}) or {}).get(metric),
        }

    # Deterministic priority ranking lets AI focus on the most decision-relevant
    # evidence without asking the model to decide materiality itself.
    compact["materiality_ranking"] = build_materiality_ranking(summary, limit=6)

    # Keep only a few category breakdowns and fewer rows within each.
    for key, details in list((context.get("top_breakdowns", {}) or {}).items())[:6]:
        compact["top_breakdowns"][key] = {
            "top": (details.get("top", []) or [])[:3],
            "bottom": (details.get("bottom", []) or [])[:2],
        }

    return compact




def compact_prompt_summary(summary):
    """Small prompt payload for faster Copilot responses."""
    numeric_stats = summary.get("numeric_stats", {}) or {}
    categories = summary.get("categories", {}) or {}
    return {
        "numeric_stats": dict(list(numeric_stats.items())[:10]),
        "categories": dict(list(categories.items())[:6]),
        "date_ranges": summary.get("date_ranges", {}),
    }

# ============================================================
# AI DASHBOARD PLANNER
# ============================================================

@st.cache_data(
    ttl=600,
    show_spinner=False,
)
def create_dashboard_plan(
    request,
    summary,
):

    prompt_summary = compact_prompt_summary(summary)

    prompt = f"""
You are designing an executive management dashboard.

USER REQUEST

{request}

DATASET

Columns:
{summary["columns"]}

Numeric columns:
{summary["numeric_columns"]}

Category columns:
{summary["category_columns"]}

Date columns:
{summary["date_columns"]}

Numeric statistics:
{prompt_summary["numeric_stats"]}

Date ranges:
{prompt_summary["date_ranges"]}

Pre-calculated analytical context:
{compact_ai_context(summary)}


Return JSON only using this schema:

{{
  "title": "Dashboard title",
  "reason": "Short explanation",
  "kpis": [
    {{
      "column": "exact numeric column",
      "aggregation": "sum",
      "label": "KPI label"
    }}
  ],
  "charts": [
    {{
      "type": "bar",
      "title": "Chart title",
      "x": "exact column or null",
      "y": "exact numeric column or null",
      "aggregation": "sum"
    }}
  ]
}}

Allowed KPI aggregations:

sum
average
minimum
maximum
count

Allowed chart types:

bar
horizontal_bar
line
area
donut
scatter
histogram
box
treemap

Rules:

- Use only exact supplied column names.
- Never invent a column.
- Build between 3 and 6 useful KPIs.
- Build between 3 and 6 useful charts when the data supports them.
- Prefer management usefulness over visual variety.
- Use donut only for low-cardinality categories.
- Use line or area when a real date column exists.
- Use scatter only when appropriate numeric fields exist.
- Use bars for categorical comparisons.
- Use horizontal_bar for ranked categories with long labels.
- Use histogram to show a numeric distribution.
- Use box to show spread/outliers across a category.
- Use treemap for hierarchical-looking category shares when cardinality is manageable.
- Choose charts that answer distinct management questions rather than repeating the same view.
- When management_watchlist evidence is supplied, prefer relevant bookmarked measures when useful, but do not ignore stronger contrary evidence.
- Do not calculate official financial values yourself.
"""


    return (
        extract_json(
            run_copilot(
                prompt
            )
        )
    )


def validate_plan(
    plan,
    profile,
):

    numeric = set(
        profile[
            "numeric"
        ]
    )

    categories = set(
        profile[
            "categories"
        ]
    )

    dates = set(
        profile[
            "dates"
        ]
    )

    all_columns = (
        numeric
        | categories
        | dates
    )


    allowed_aggs = {

        "sum",
        "average",
        "minimum",
        "maximum",
        "count",
    }


    allowed_charts = {

        "bar",
        "horizontal_bar",
        "line",
        "area",
        "donut",
        "scatter",
        "histogram",
        "box",
        "treemap",
    }


    clean = {

        "title":
            str(
                plan.get(
                    "title",
                    "AI Designed Dashboard",
                )
            ),

        "reason":
            str(
                plan.get(
                    "reason",
                    "",
                )
            ),

        "kpis":
            [],

        "charts":
            [],
    }


    for item in plan.get(
        "kpis",
        [],
    )[:6]:

        column = (
            item.get(
                "column"
            )
        )

        aggregation = (
            item.get(
                "aggregation",
                "sum",
            )
        )


        if column not in numeric:

            continue


        if aggregation not in allowed_aggs:

            aggregation = (
                "sum"
            )


        clean[
            "kpis"
        ].append(
            {

                "column":
                    column,

                "aggregation":
                    aggregation,

                "label":
                    str(
                        item.get(
                            "label",
                            column,
                        )
                    ),
            }
        )


    for item in plan.get(
        "charts",
        [],
    )[:6]:

        chart_type = (
            item.get(
                "type"
            )
        )

        x = (
            item.get(
                "x"
            )
        )

        y = (
            item.get(
                "y"
            )
        )

        aggregation = (
            item.get(
                "aggregation",
                "sum",
            )
        )


        if chart_type not in allowed_charts:

            continue


        if (
            x is not None
            and x not in all_columns
        ):

            continue


        if (
            y is not None
            and y not in numeric
        ):

            continue


        if aggregation not in allowed_aggs:

            aggregation = (
                "sum"
            )


        clean[
            "charts"
        ].append(
            {

                "type":
                    chart_type,

                "title":
                    str(
                        item.get(
                            "title",
                            chart_type.title(),
                        )
                    ),

                "x":
                    x,

                "y":
                    y,

                "aggregation":
                    aggregation,
            }
        )


    return clean


# ============================================================
# AI MANAGEMENT SUMMARY
# ============================================================

@st.cache_data(
    ttl=600,
    show_spinner=False,
)
def generate_management_summary(
    summary,
    plan,
):

    prompt_summary = compact_prompt_summary(summary)

    prompt = f"""
You are an executive management analyst.

Python has already calculated the dataset statistics.

Do not invent numbers.

Quarter:
{summary["quarter"]}

Dashboard:
{plan["title"]}

Dashboard reason:
{plan["reason"]}

Numeric statistics:
{prompt_summary["numeric_stats"]}

Categories:
{prompt_summary["categories"]}

Date ranges:
{prompt_summary["date_ranges"]}

Pre-calculated analytical context:
{compact_ai_context(summary)}

Write an executive-grade management summary with clear sections:
1. Executive takeaway — 2 to 4 sentences on what matters most.
2. Performance signals — the strongest supported movements, concentration, trend or variance signals.
3. Risks / watch items — only risks supported by the supplied facts, with evidence.
4. Decisions / actions — 2 to 4 practical management actions or questions.

Prioritize materiality and business relevance. Quantify observations whenever a supplied value supports them.
Distinguish facts from interpretation. Never invent causes.
If the available statistics do not support a conclusion, say so clearly.
"""


    evidence = {
        "quarter": summary.get("quarter"),
        "dashboard": {
            "title": plan.get("title"),
            "reason": plan.get("reason"),
        },
        "numeric_statistics": prompt_summary.get("numeric_stats", {}),
        "categories": prompt_summary.get("categories", {}),
        "date_ranges": prompt_summary.get("date_ranges", {}),
        "analytical_context": compact_ai_context(summary),
    }

    return grounded_ai_response(
        prompt,
        evidence,
        purpose="management summary",
    )


def generate_trends_ai_summary(summary):
    """Generate a grounded executive summary specifically for the Trends board."""
    context = compact_ai_context(summary)
    evidence = {
        "quarter": summary.get("quarter"),
        "date_ranges": summary.get("date_ranges", {}),
        "trends": context.get("trends", {}),
        "anomalies": context.get("anomalies", {}),
        "correlations": context.get("correlations", []),
        "driver_analysis": context.get("driver_analysis", {}),
        "financial_variance": context.get("financial_variance", {}),
        "materiality_ranking": context.get("materiality_ranking", []),
    }

    prompt = f"""
You are an executive analyst reviewing the Trends board of a Workiva dashboard.

Grounded trend evidence:
{evidence}

Rules:
- Use only the supplied evidence. Never invent dates, values, causes, seasonality, or business events.
- Describe direction, magnitude, peaks/lows, anomalies, relationships, and calculated category drivers only when directly supported.
- When driver analysis is available, distinguish mathematical contribution from causal explanation. A contributor explains where the numeric change occurred; it does not prove why it occurred.
- Treat IQR anomaly flags as screening signals, not proof of an error or root cause.
- Distinguish observed facts from interpretation.
- If there is not enough time history to support a strong trend conclusion, say so clearly.
- Follow the supplied deterministic materiality ranking when prioritizing signals; do not invent a different ranking.
- Treat the score as a decision-priority heuristic, not an accounting materiality threshold.
- Keep the tone concise, executive, and decision-oriented.

Write sections:
1. Trend takeaway
2. Strongest movements
3. Key drivers of change
4. Anomalies / watch items
5. Management questions
"""
    return grounded_ai_response(
        prompt,
        evidence,
        purpose="trends summary",
    )


# ============================================================
# CHAT
# ============================================================

def chat_answer(
    question,
    summary,
    plan,
    data,
    profile,
    business,
    chat_history=None,
    forecast_result=None,
    comparison_result=None,
):

    question_context = build_question_context(
        question,
        data,
        profile,
        business,
    )

    # Keep enough conversational continuity without repeatedly sending a long
    # transcript to Copilot. This materially reduces prompt size/latency.
    recent_history = (chat_history or [])[-4:]
    cross_quality_context = st.session_state.get("cross_quarter_quality")

    compact_context = compact_ai_context(summary)
    compact_summary = compact_prompt_summary(summary)

    # Give Copilot the exact dashboard the user is currently looking at.
    # This is local context construction only; it does not trigger another
    # Workiva request or AI call.
    visible_kpis = [
        {
            "label": item.get("label", item.get("column")),
            "column": item.get("column"),
            "aggregation": item.get("aggregation", "sum"),
        }
        for item in (plan or {}).get("kpis", [])
    ]

    visible_chart_items = [
        item for item in (plan or {}).get("charts", [])
        if item.get("visible", True)
    ]
    visible_charts = [
        {
            "position": (
                "primary" if index == 0 else "supporting"
            ),
            "number": index + 1,
            "title": item.get("title"),
            "type": item.get("type"),
            "x": item.get("x"),
            "y": item.get("y"),
            "aggregation": item.get("aggregation", "sum"),
            "top_n": item.get("top_n", "All"),
            "sort_order": item.get("sort_order", "Descending"),
            "show_percentage": item.get("show_percentage", False),
        }
        for index, item in enumerate(visible_chart_items)
    ]

    forecast_context = None
    if forecast_result and forecast_result.get("forecasts"):
        forecast_context = {
            "historical_quarters": forecast_result.get("quarters", []),
            "forecast_horizon_quarters": forecast_result.get("horizon", 0),
            "scenario_sensitivity_pct": forecast_result.get("scenario_sensitivity_pct", 0),
            "metrics": {},
        }
        for metric, frame in forecast_result.get("forecasts", {}).items():
            backtest = (forecast_result.get("backtests", {}) or {}).get(metric, {}) or {}
            selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
            scenario_frame = (forecast_result.get("scenarios", {}) or {}).get(metric)
            forecast_context["metrics"][metric] = {
                "selected_model": {
                    "model_name": selection.get("model_name"),
                    "reason": selection.get("reason"),
                },
                "series": (
                    frame[["Quarter", "Value", "Series"]]
                    .to_dict(orient="records")
                ),
                "scenarios": (
                    scenario_frame[["Quarter", "Value", "Scenario"]].to_dict(orient="records")
                    if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty
                    else []
                ),
                "backtest": {
                    "history_points": backtest.get("history_points"),
                    "backtest_points": backtest.get("backtest_points"),
                    "mae": backtest.get("mae"),
                    "rmse": backtest.get("rmse"),
                    "mape_pct": backtest.get("mape_pct"),
                    "wape_pct": backtest.get("wape_pct"),
                    "smape_pct": backtest.get("smape_pct"),
                    "confidence": backtest.get("confidence"),
                    "confidence_reason": backtest.get("confidence_reason"),
                },
            }

    trends_context = {
        "trends": compact_context.get("trends", {}),
        "anomalies": compact_context.get("anomalies", {}),
        "correlations": compact_context.get("correlations", []),
        "driver_analysis": compact_context.get("driver_analysis", {}),
    }

    comparison_context = _comparison_management_brief_payload(comparison_result)

    dashboard_context = {
        "title": (plan or {}).get("title"),
        "reason": (plan or {}).get("reason"),
        "layout": {
            "kpi_row": visible_kpis,
            "primary_chart": (
                visible_charts[0] if visible_charts else None
            ),
            "supporting_charts": (
                visible_charts[1:] if len(visible_charts) > 1 else []
            ),
            "management_focus_panels": [
                "Variance",
                "Risk / watch",
                "Management questions",
            ],
        },
        "charts": visible_charts,
    }

    prompt = f"""
You are a senior conversational analyst inside a Workiva management dashboard.

Quarter:
{summary["quarter"]}

Dashboard:
{plan.get("title", "Current dashboard")}

Visible dashboard configuration:
{dashboard_context}

Columns:
{summary["columns"]}

Relevant numeric statistics:
{compact_summary["numeric_stats"]}

Relevant category information:
{compact_summary["categories"]}

Pre-calculated analytical context:
{compact_context}

Trends board context:
{trends_context}

Current forecasting board context (if built):
{forecast_context}

Current deterministic data-quality intelligence:
{summary.get("data_quality_intelligence", {})}

Cross-quarter quality scan (only if the user has run it):
{cross_quality_context}

Current quarter-comparison / variance-bridge context (if built):
{comparison_context}

Management watchlist (user-prioritized metrics, not a factual source):
{summary.get("management_watchlist", [])}

Question-specific calculations:
{question_context}

Recent conversation:
{recent_history}

User question:
{question}

Rules:
- Never invent numbers, causes, definitions or data that are not supplied.
- Answer the question first, then give the most useful supporting evidence.
- Treat the visible dashboard configuration as authoritative for what the user can currently see.
- If the user refers to a chart by position or number (for example, "the first chart" or "the second chart"), resolve it from the visible dashboard configuration.
- If the user asks why a visible KPI or chart matters, connect that visual to the supplied analytical facts without inventing causes.
- Use exact figures when available and explain the calculation basis (sum, average, ranking, etc.).
- For ranking questions, identify leaders and laggards from the supplied breakdowns.
- Treat the management watchlist as a user preference for attention, never as proof that a metric is objectively more material.
- For broad management questions, mention watchlisted metrics when evidence is relevant, but do not suppress stronger non-watchlisted signals.
- For trend questions, use the Trends board context and describe direction, magnitude, peaks/lows, anomaly screening signals, correlations, and calculated drivers only when supported.
- For driver questions, explain which category values mathematically contributed to the change. Never convert contribution into an unsupported causal claim.
- For variance-bridge questions, use the current comparison bridge when available. Explain starting total, ranked category contributions, residual Other (if any), and ending total; confirm that arithmetic contribution does not establish causation.
- For forecasting questions, use the current forecasting board context when available. Clearly distinguish historical actuals from projected values and identify the selected deterministic forecast model when relevant.
- When forecast backtest evidence is available, use it to explain historical forecast error and the supplied confidence label. Never upgrade the supplied confidence label or imply that historical accuracy guarantees future performance.
- If the user asks about a forecast but no forecasting board has been built, say that clearly and direct them to build one rather than inventing a projection.
- For data-quality or reliability questions, use the deterministic quality-intelligence evidence and any cross-quarter scan that has actually been run. Treat the readiness score as an analytical-use heuristic, not an audit opinion.
- If quality status is Watch or Limited, explicitly qualify trend/forecast conclusions rather than hiding the limitation.
- If the question is ambiguous, state the most reasonable interpretation and what would change the answer.
- If the data is insufficient, say exactly what is missing and suggest the next useful question.
- Keep the tone concise, executive and analytical, but allow enough detail to explain the evidence.
"""

    evidence = {
        "quarter": summary.get("quarter"),
        "dashboard": dashboard_context,
        "columns": summary.get("columns", []),
        "numeric_statistics": compact_summary.get("numeric_stats", {}),
        "category_information": compact_summary.get("categories", {}),
        "analytical_context": compact_context,
        "trends_board_context": trends_context,
        "forecasting_board_context": forecast_context,
        "data_quality_intelligence": summary.get("data_quality_intelligence", {}),
        "cross_quarter_quality": cross_quality_context,
        "quarter_comparison_context": comparison_context,
        "management_watchlist": summary.get("management_watchlist", []),
        "question_specific_calculations": question_context,
        "user_question": question,
    }

    return grounded_ai_response(
        prompt,
        evidence,
        purpose="dashboard question answering",
    )


# ============================================================
# KPI / CHART HELPERS
# ============================================================

def aggregate_value(
    series,
    method,
):

    values = (
        pd.to_numeric(
            series,
            errors="coerce",
        )
    )


    if method == "average":

        return (
            values.mean()
        )


    if method == "minimum":

        return (
            values.min()
        )


    if method == "maximum":

        return (
            values.max()
        )


    if method == "count":

        return (
            values.count()
        )


    return (
        values.sum()
    )


def grouped_data(
    data,
    category,
    metric,
    aggregation,
):

    working = (
        data[
            [
                category,
                metric,
            ]
        ].copy()
    )


    working[
        metric
    ] = (
        pd.to_numeric(
            working[
                metric
            ],
            errors="coerce",
        )
    )


    working = (
        working.dropna()
    )


    group = (
        working.groupby(
            category,
            dropna=False,
        )[metric]
    )


    if aggregation == "average":

        result = (
            group.mean()
        )


    elif aggregation == "minimum":

        result = (
            group.min()
        )


    elif aggregation == "maximum":

        result = (
            group.max()
        )


    elif aggregation == "count":

        result = (
            group.count()
        )


    else:

        result = (
            group.sum()
        )


    return (
        result
        .reset_index()
    )


def render_kpis(
    data,
    kpis,
):

    if not kpis:

        return


    columns = (
        st.columns(
            len(kpis)
        )
    )


    for index, item in enumerate(
        kpis
    ):

        value = (
            aggregate_value(
                data[
                    item[
                        "column"
                    ]
                ],
                item[
                    "aggregation"
                ],
            )
        )


        columns[
            index
        ].metric(
            item[
                "label"
            ],
            f"{value:,.2f}",
        )


@st.cache_data(ttl=1800, show_spinner=False, max_entries=96)
def build_chart_figure(data, chart):
    if not chart.get("visible", True):
        return None

    chart_type = chart.get("type", "bar")
    x = chart.get("x")
    y = chart.get("y")
    title = chart.get("title", chart_type.title())
    aggregation = chart.get("aggregation", "sum")
    sort_order = chart.get("sort_order", "Descending")
    top_n = chart.get("top_n")
    show_percentage = bool(chart.get("show_percentage", False))

    if chart_type == "histogram":
        metric = y or x
        if metric not in data.columns:
            return None
        return px.histogram(data, x=metric, title=title)

    if chart_type == "scatter":
        if x not in data.columns or y not in data.columns:
            return None
        return px.scatter(data, x=x, y=y, title=title)

    if chart_type == "box":
        if y not in data.columns:
            return None
        if x and x in data.columns:
            return px.box(data, x=x, y=y, title=title)
        return px.box(data, y=y, title=title)

    if x is None or y is None or x not in data.columns or y not in data.columns:
        return None

    chart_data = grouped_data(data, x, y, aggregation)
    if chart_data.empty:
        return None

    is_date = pd.api.types.is_datetime64_any_dtype(data[x])

    if is_date:
        chart_data = chart_data.sort_values(x)
    elif sort_order == "Ascending":
        chart_data = chart_data.sort_values(y, ascending=True)
    elif sort_order == "Descending":
        chart_data = chart_data.sort_values(y, ascending=False)

    if not is_date and top_n not in (None, "All"):
        try:
            limit = int(top_n)
            if limit > 0:
                chart_data = chart_data.head(limit)
        except (TypeError, ValueError):
            pass

    display_y = y
    if show_percentage and not is_date and chart_type in {
        "bar", "horizontal_bar", "donut", "treemap"
    }:
        total = pd.to_numeric(chart_data[y], errors="coerce").sum()
        if pd.notna(total) and total != 0:
            display_y = "Share %"
            chart_data = chart_data.copy()
            chart_data[display_y] = (
                pd.to_numeric(chart_data[y], errors="coerce") / total * 100
            )
            title = f"{title} (%)"

    if chart_type == "bar":
        figure = px.bar(chart_data, x=x, y=display_y, title=title)
    elif chart_type == "horizontal_bar":
        figure = px.bar(
            chart_data,
            x=display_y,
            y=x,
            orientation="h",
            title=title,
        )
    elif chart_type == "line":
        figure = px.line(chart_data, x=x, y=y, markers=True, title=title)
    elif chart_type == "area":
        figure = px.area(chart_data, x=x, y=y, title=title)
    elif chart_type == "donut":
        figure = px.pie(
            chart_data,
            names=x,
            values=display_y,
            hole=0.55,
            title=title,
        )
    elif chart_type == "treemap":
        figure = px.treemap(chart_data, path=[x], values=display_y, title=title)
    else:
        return None

    if display_y == "Share %" and chart_type in {"bar", "horizontal_bar"}:
        figure.update_layout(yaxis_title="Share %") if chart_type == "bar" else figure.update_layout(xaxis_title="Share %")

    return figure



def build_visual_ai_evidence(data, chart):
    """Build a compact deterministic evidence package for one dashboard visual."""
    chart_type = str(chart.get("type", "bar"))
    x = chart.get("x")
    y = chart.get("y")
    aggregation = chart.get("aggregation", "sum")
    evidence = {
        "title": str(chart.get("title") or chart_type.title()),
        "chart_type": chart_type,
        "x": str(x) if x is not None else None,
        "y": str(y) if y is not None else None,
        "aggregation": aggregation,
        "source_rows": int(len(data)),
    }

    def numeric_stats(series):
        values = pd.to_numeric(series, errors="coerce").dropna()
        if values.empty:
            return {}
        return {
            "count": int(values.count()),
            "sum": float(values.sum()),
            "average": float(values.mean()),
            "median": float(values.median()),
            "minimum": float(values.min()),
            "maximum": float(values.max()),
        }

    if chart_type == "histogram":
        metric = y or x
        if metric in data.columns:
            evidence["metric_stats"] = numeric_stats(data[metric])
        return evidence

    if chart_type == "scatter":
        if x in data.columns and y in data.columns:
            working = data[[x, y]].copy()
            working[x] = pd.to_numeric(working[x], errors="coerce")
            working[y] = pd.to_numeric(working[y], errors="coerce")
            working = working.dropna()
            evidence["x_stats"] = numeric_stats(working[x])
            evidence["y_stats"] = numeric_stats(working[y])
            if len(working) >= 2:
                corr = working[x].corr(working[y])
                evidence["correlation"] = float(corr) if pd.notna(corr) else None
        return evidence

    if chart_type == "box":
        if y in data.columns:
            evidence["metric_stats"] = numeric_stats(data[y])
            if x in data.columns:
                grouped = []
                for name, group in data[[x, y]].dropna().groupby(x):
                    stats = numeric_stats(group[y])
                    if stats:
                        grouped.append({"category": str(name), **stats})
                grouped.sort(key=lambda row: abs(row.get("median", 0)), reverse=True)
                evidence["group_stats"] = grouped[:10]
        return evidence

    if x in data.columns and y in data.columns:
        grouped = grouped_data(data, x, y, aggregation).dropna()
        if not grouped.empty:
            is_date = pd.api.types.is_datetime64_any_dtype(data[x])
            if is_date:
                grouped = grouped.sort_values(x)
                evidence["ordered_values"] = [
                    {"x": str(row[x]), "value": float(row[y])}
                    for _, row in grouped.head(12).iterrows()
                ]
                first_value = float(grouped[y].iloc[0])
                last_value = float(grouped[y].iloc[-1])
                evidence["trend"] = {
                    "first": {"x": str(grouped[x].iloc[0]), "value": first_value},
                    "last": {"x": str(grouped[x].iloc[-1]), "value": last_value},
                    "change": float(last_value - first_value),
                    "change_pct": float((last_value - first_value) / abs(first_value) * 100) if first_value else None,
                    "peak": {"x": str(grouped.loc[grouped[y].idxmax(), x]), "value": float(grouped[y].max())},
                    "low": {"x": str(grouped.loc[grouped[y].idxmin(), x]), "value": float(grouped[y].min())},
                }
            else:
                ranked = grouped.sort_values(y, ascending=False)
                evidence["top_values"] = [
                    {"x": str(row[x]), "value": float(row[y])}
                    for _, row in ranked.head(7).iterrows()
                ]
                evidence["bottom_values"] = [
                    {"x": str(row[x]), "value": float(row[y])}
                    for _, row in ranked.tail(3).iterrows()
                ]
                evidence["total"] = float(pd.to_numeric(grouped[y], errors="coerce").sum())
                evidence["category_count"] = int(len(grouped))
    return evidence


def render_visual_ai_action(evidence, key, label="✦ Explain this visual"):
    """Render an on-demand grounded AI explanation without adding background AI calls."""
    if "visual_ai_explanations" not in st.session_state:
        st.session_state["visual_ai_explanations"] = {}

    state_key = str(key)
    if st.button(label, key=f"{state_key}_ai_explain", help="Explain this chart using only the calculated evidence shown by the application."):
        prompt = """
You are explaining one management dashboard visual.

Write a concise, executive-quality explanation with exactly these parts:
- **What it shows:** the main pattern in the visual.
- **Why it matters:** the decision-relevant implication supported by the evidence.
- **What to inspect next:** one useful follow-up question or investigation.

Rules:
- Explain only what this visual's evidence establishes.
- Distinguish observation from interpretation.
- Do not invent root causes, targets, benchmarks, probabilities, external events, or missing context.
- If the evidence is weak or too sparse, say so plainly.
- Keep the response compact; this is an inline chart explanation, not a full management summary.
"""
        with st.spinner("Explaining visual..."):
            st.session_state["visual_ai_explanations"][state_key] = grounded_ai_response(
                prompt,
                evidence,
                purpose="visual explanation",
            )

    explanation = st.session_state["visual_ai_explanations"].get(state_key)
    if explanation:
        with st.expander("AI visual explanation", expanded=True):
            st.markdown(explanation)
            st.caption("Grounded in the calculated evidence for this visual; AI does not determine the underlying values.")


def render_chart(data, chart, key):
    """Render one Plotly chart with an explicit Streamlit element key.

    Streamlit derives an internal element ID from a chart's parameters when no
    key is supplied. The same figure rendered in more than one place during a
    single rerun can therefore collide. Requiring a key here prevents that
    class of DuplicateElementId errors throughout the dashboard.

    Accessible reading also increases chart typography and interaction targets.
    This keeps the underlying data and chart choice unchanged while making the
    visual easier to read for low-vision and keyboard users.
    """
    figure = build_chart_figure(data, chart)
    if figure is not None:
        if st.session_state.get("display_mode") == "Accessible reading":
            figure.update_layout(
                font={"size": 16},
                title={"font": {"size": 22}},
                legend={"font": {"size": 15}},
                margin={"l": 70, "r": 35, "t": 75, "b": 70},
                hoverlabel={"font_size": 16},
            )
            figure.update_xaxes(
                title_font={"size": 17},
                tickfont={"size": 15},
                automargin=True,
                showgrid=True,
            )
            figure.update_yaxes(
                title_font={"size": 17},
                tickfont={"size": 15},
                automargin=True,
                showgrid=True,
            )
            figure.update_traces(
                marker_line_width=1.2,
                selector=dict(type="bar"),
            )

        st.plotly_chart(
            figure,
            use_container_width=True,
            key=key,
            config={
                "displaylogo": False,
                "responsive": True,
            },
        )

        # Accessible reading provides a concise textual equivalent of the chart
        # structure without duplicating the entire underlying table.
        if st.session_state.get("display_mode") == "Accessible reading":
            title = chart.get("title") or "Chart"
            chart_type = str(chart.get("type") or "chart").replace("_", " ")
            x = chart.get("x")
            y = chart.get("y")
            aggregation = chart.get("aggregation", "sum")
            parts = [f"{title} is a {chart_type} chart"]
            if y:
                parts.append(f"showing {aggregation} of {y}")
            if x:
                parts.append(f"by {x}")
            st.caption("Accessible chart description: " + " ".join(parts) + ".")

        render_visual_ai_action(
            build_visual_ai_evidence(data, chart),
            key=f"{key}_visual",
        )


def customize_dashboard_plan(plan, data, profile, key_prefix="dashboard"):
    """Return a locally customized copy of the dashboard plan.

    All controls operate on in-memory data only. They do not call Workiva or Copilot.
    """
    if not plan:
        return plan

    customized = deepcopy(plan)
    numeric_options = list(profile.get("numeric", []))
    dimension_options = list(dict.fromkeys(
        profile.get("categories", []) + profile.get("dates", []) + numeric_options
    ))

    fingerprint = abs(hash(json.dumps(plan, sort_keys=True, default=str)))
    base_key = f"{key_prefix}_{fingerprint}"

    with st.expander("Customize dashboard", expanded=False):
        st.caption(
            "These controls update the visible charts locally. "
            "They do not send another request to Copilot or Workiva. "
            "Changes are applied together to avoid rerunning the app for every control."
        )

        with st.form(f"{base_key}_customizer_form", clear_on_submit=False):
            for index, chart in enumerate(customized.get("charts", [])):
                st.markdown(f"**Chart {index + 1}: {chart.get('title', 'Dashboard chart')}**")
    
                row1 = st.columns([1, 2, 2, 2])
                chart["visible"] = row1[0].checkbox(
                    "Visible",
                    value=chart.get("visible", True),
                    key=f"{base_key}_visible_{index}",
                )
    
                chart_types = [
                    "bar", "horizontal_bar", "line", "area", "donut",
                    "scatter", "histogram", "box", "treemap",
                ]
                current_type = chart.get("type", "bar")
                chart["type"] = row1[1].selectbox(
                    "Chart type",
                    chart_types,
                    index=(chart_types.index(current_type) if current_type in chart_types else 0),
                    key=f"{base_key}_type_{index}",
                )
    
                x_options = [None] + dimension_options
                current_x = chart.get("x")
                chart["x"] = row1[2].selectbox(
                    "Dimension / X",
                    x_options,
                    index=(x_options.index(current_x) if current_x in x_options else 0),
                    key=f"{base_key}_x_{index}",
                    format_func=lambda value: "None" if value is None else str(value),
                )
    
                y_options = [None] + numeric_options
                current_y = chart.get("y")
                chart["y"] = row1[3].selectbox(
                    "Metric / Y",
                    y_options,
                    index=(y_options.index(current_y) if current_y in y_options else 0),
                    key=f"{base_key}_y_{index}",
                    format_func=lambda value: "None" if value is None else str(value),
                )
    
                row2 = st.columns(4)
                aggregations = ["sum", "average", "minimum", "maximum", "count"]
                current_agg = chart.get("aggregation", "sum")
                chart["aggregation"] = row2[0].selectbox(
                    "Aggregation",
                    aggregations,
                    index=(aggregations.index(current_agg) if current_agg in aggregations else 0),
                    key=f"{base_key}_agg_{index}",
                )
    
                top_options = ["All", 5, 10, 15, 20]
                current_top = chart.get("top_n", "All")
                if current_top not in top_options:
                    current_top = "All"
                chart["top_n"] = row2[1].selectbox(
                    "Top N",
                    top_options,
                    index=top_options.index(current_top),
                    key=f"{base_key}_top_{index}",
                )
    
                sort_options = ["Descending", "Ascending", "Original"]
                current_sort = chart.get("sort_order", "Descending")
                chart["sort_order"] = row2[2].selectbox(
                    "Sort",
                    sort_options,
                    index=(sort_options.index(current_sort) if current_sort in sort_options else 0),
                    key=f"{base_key}_sort_{index}",
                )
    
                chart["show_percentage"] = row2[3].checkbox(
                    "Show %",
                    value=bool(chart.get("show_percentage", False)),
                    key=f"{base_key}_pct_{index}",
                    help="For categorical bar, horizontal bar, donut and treemap charts, show each category as a share of the displayed total.",
                )
    
                if index < len(customized.get("charts", [])) - 1:
                    st.divider()

            st.form_submit_button("Apply dashboard changes", use_container_width=True)

    return customized


# ============================================================
# DATA-QUALITY INTELLIGENCE
# ============================================================

@st.cache_data(ttl=1800, show_spinner=False, max_entries=48)
def assess_data_quality_intelligence(data, profile, business, quarter):
    """Build a deterministic, explainable quality gate for analytical use.

    The score is an application-readiness heuristic, not an audit opinion.  It
    deliberately avoids AI so analysis quality can be assessed before any model
    commentary is generated.
    """
    rows = int(len(data))
    cells = max(int(data.size), 1)
    missing_cells = int(data.isna().sum().sum())
    missing_rate = float(missing_cells / cells * 100.0)
    duplicate_rows = int(data.duplicated().sum()) if rows else 0
    duplicate_rate = float(duplicate_rows / max(rows, 1) * 100.0)

    checks = []
    deductions = 0.0

    def add_check(name, status, detail, deduction=0.0, evidence=None):
        nonlocal deductions
        deductions += max(float(deduction), 0.0)
        checks.append({
            "check": name,
            "status": status,
            "detail": detail,
            "deduction": round(max(float(deduction), 0.0), 1),
            "evidence": evidence or {},
        })

    # Overall completeness.
    if missing_rate >= 20:
        add_check("Completeness", "Critical", f"{missing_rate:.1f}% of cells are missing.", 25, {"missing_rate_pct": missing_rate})
    elif missing_rate >= 10:
        add_check("Completeness", "Watch", f"{missing_rate:.1f}% of cells are missing.", 14, {"missing_rate_pct": missing_rate})
    elif missing_rate >= 5:
        add_check("Completeness", "Watch", f"{missing_rate:.1f}% of cells are missing.", 7, {"missing_rate_pct": missing_rate})
    else:
        add_check("Completeness", "Good", f"Overall missingness is {missing_rate:.1f}%.", 0, {"missing_rate_pct": missing_rate})

    # Exact duplicate rows are a safe, non-semantic duplicate check.
    if duplicate_rate >= 10:
        add_check("Duplicate rows", "Critical", f"{duplicate_rows:,} exact duplicate rows ({duplicate_rate:.1f}%).", 20, {"duplicate_rows": duplicate_rows, "duplicate_rate_pct": duplicate_rate})
    elif duplicate_rate >= 2:
        add_check("Duplicate rows", "Watch", f"{duplicate_rows:,} exact duplicate rows ({duplicate_rate:.1f}%).", 10, {"duplicate_rows": duplicate_rows, "duplicate_rate_pct": duplicate_rate})
    elif duplicate_rows:
        add_check("Duplicate rows", "Review", f"{duplicate_rows:,} exact duplicate rows ({duplicate_rate:.1f}%).", 3, {"duplicate_rows": duplicate_rows, "duplicate_rate_pct": duplicate_rate})
    else:
        add_check("Duplicate rows", "Good", "No exact duplicate rows detected.", 0, {"duplicate_rows": 0})

    # Concentrated missingness can hide behind a low dataset-wide average.
    missing_by_col = []
    for col in data.columns:
        pct = float(data[col].isna().mean() * 100.0)
        if pct > 0:
            missing_by_col.append((pct, str(col)))
    missing_by_col.sort(reverse=True)
    severe_cols = [(name, pct) for pct, name in missing_by_col if pct >= 30]
    watch_cols = [(name, pct) for pct, name in missing_by_col if 10 <= pct < 30]
    if severe_cols:
        preview = ", ".join(f"{name} ({pct:.0f}%)" for name, pct in severe_cols[:4])
        add_check("Column completeness", "Critical", f"High missingness in: {preview}.", min(20, 5 + 3 * len(severe_cols)), {"columns": severe_cols[:8]})
    elif watch_cols:
        preview = ", ".join(f"{name} ({pct:.0f}%)" for name, pct in watch_cols[:4])
        add_check("Column completeness", "Watch", f"Moderate missingness in: {preview}.", min(10, 2 + 2 * len(watch_cols)), {"columns": watch_cols[:8]})
    else:
        add_check("Column completeness", "Good", "No column has 10% or more missing values.", 0)

    # Period integrity: a quarter-labelled dataset should not silently contain
    # several different quarter periods after parsing.
    date_cols = list(profile.get("dates", []) or [])
    quarter_like = [c for c in date_cols if "quarter" in str(c).lower() or "period" in str(c).lower()]
    if quarter_like:
        period_col = quarter_like[0]
        values = pd.to_datetime(data[period_col], errors="coerce").dropna()
        periods = sorted({str(v.to_period("Q")) for v in values})
        if len(periods) > 1:
            add_check("Period integrity", "Watch", f"{period_col} contains {len(periods)} distinct quarter periods: {', '.join(periods[:6])}.", min(15, 5 + 2 * (len(periods)-1)), {"column": str(period_col), "periods": periods})
        elif len(periods) == 1:
            add_check("Period integrity", "Good", f"{period_col} resolves consistently to {periods[0]}.", 0, {"column": str(period_col), "periods": periods})

    # Analytical coverage. Missing business concepts are not data errors, so the
    # deduction is intentionally small and framed as analytical limitation.
    numeric_cols = list(profile.get("numeric", []) or [])
    if not numeric_cols:
        add_check("Analytical coverage", "Critical", "No numeric measures are available for quantitative analysis.", 30)
    elif len(numeric_cols) == 1:
        add_check("Analytical coverage", "Review", "Only one numeric measure is available; comparative analysis is limited.", 5)
    else:
        add_check("Analytical coverage", "Good", f"{len(numeric_cols)} numeric measures are available for analysis.", 0)

    # Constant numeric fields cannot support trends/variance even though they may
    # be valid data. Treat this as a usability warning, not a quality failure.
    constant_numeric = []
    for col in numeric_cols:
        vals = pd.to_numeric(data[col], errors="coerce").dropna()
        if len(vals) > 1 and vals.nunique() <= 1:
            constant_numeric.append(str(col))
    if constant_numeric:
        add_check("Metric variability", "Review", "Constant numeric fields: " + ", ".join(constant_numeric[:6]) + ".", min(6, len(constant_numeric)), {"columns": constant_numeric[:10]})
    else:
        add_check("Metric variability", "Good", "No constant numeric measures detected.", 0)

    score = round(max(0.0, 100.0 - deductions), 1)
    if score >= 90:
        status = "Strong"
    elif score >= 75:
        status = "Good"
    elif score >= 55:
        status = "Watch"
    else:
        status = "Limited"

    critical = [item for item in checks if item["status"] == "Critical"]
    watch = [item for item in checks if item["status"] in {"Watch", "Review"}]
    return {
        "quarter": str(quarter),
        "score": score,
        "status": status,
        "checks": checks,
        "critical_count": len(critical),
        "watch_count": len(watch),
        "analysis_guidance": (
            "Proceed with normal analytical interpretation." if status in {"Strong", "Good"}
            else "Interpret trend and forecast conclusions with additional caution." if status == "Watch"
            else "Resolve material data-quality issues before relying on trend or forecast conclusions."
        ),
    }


@st.cache_data(ttl=1800, show_spinner=False, max_entries=24)
def compare_quarter_quality(bundles):
    """Compare up to four already-requested quarter bundles for structural drift."""
    bundles = list(bundles or [])[:4]
    if not bundles:
        return {"quarters": [], "status": "Unavailable", "checks": [], "score": None}

    rows = {str(b["quarter"]): int(len(b["data"])) for b in bundles}
    schemas = {str(b["quarter"]): {str(c).strip().lower(): str(c) for c in b["data"].columns} for b in bundles}
    type_maps = {
        str(b["quarter"]): {
            str(col).strip().lower(): (b["profile"].get("details", {}).get(col, {}) or {}).get("type", "unknown")
            for col in b["data"].columns
        }
        for b in bundles
    }
    checks = []
    deductions = 0.0

    def add(name, status, detail, deduction=0.0, evidence=None):
        nonlocal deductions
        deductions += max(float(deduction), 0.0)
        checks.append({"check": name, "status": status, "detail": detail, "deduction": round(float(deduction), 1), "evidence": evidence or {}})

    # Row-count drift, using adjacent selected quarters.
    row_changes = []
    ordered = [str(b["quarter"]) for b in bundles]
    for previous, current in zip(ordered, ordered[1:]):
        prev_rows = rows[previous]
        curr_rows = rows[current]
        pct = ((curr_rows - prev_rows) / prev_rows * 100.0) if prev_rows else None
        row_changes.append({"from": previous, "to": current, "change_pct": pct, "rows_from": prev_rows, "rows_to": curr_rows})
    large = [x for x in row_changes if x["change_pct"] is not None and abs(x["change_pct"]) >= 30]
    if large:
        worst = max(large, key=lambda x: abs(x["change_pct"]))
        add("Row-count consistency", "Watch", f"Largest adjacent row-count change is {worst['change_pct']:+.1f}% ({worst['from']} → {worst['to']}).", min(18, 8 + 3 * len(large)), {"changes": row_changes})
    else:
        add("Row-count consistency", "Good", "No adjacent selected quarter changes row count by 30% or more.", 0, {"changes": row_changes})

    # Schema drift relative to the first selected quarter.
    base_q = ordered[0]
    base_cols = set(schemas[base_q])
    schema_drift = []
    for q in ordered[1:]:
        cols = set(schemas[q])
        missing = sorted(base_cols - cols)
        added = sorted(cols - base_cols)
        if missing or added:
            schema_drift.append({"quarter": q, "missing_columns": missing, "new_columns": added})
    if schema_drift:
        add("Schema consistency", "Watch", f"Column structure changes in {len(schema_drift)} selected quarter(s).", min(20, 6 + 4 * len(schema_drift)), {"drift": schema_drift})
    else:
        add("Schema consistency", "Good", "Column structure is consistent across selected quarters.", 0)

    # Type drift for columns that exist in all selected quarters.
    common = set.intersection(*(set(schemas[q]) for q in ordered)) if ordered else set()
    type_drift = []
    for key in sorted(common):
        observed = {q: type_maps[q].get(key, "unknown") for q in ordered}
        if len(set(observed.values())) > 1:
            type_drift.append({"column": schemas[base_q].get(key, key), "types": observed})
    if type_drift:
        add("Type consistency", "Watch", f"Detected type changes for {len(type_drift)} common column(s).", min(16, 4 + 3 * len(type_drift)), {"columns": type_drift[:10]})
    else:
        add("Type consistency", "Good", "Detected column types are consistent across selected quarters.", 0)

    # Category drift for low-cardinality common dimensions.
    category_drift = []
    base_bundle = bundles[0]
    for col in (base_bundle["profile"].get("categories", []) or [])[:8]:
        key = str(col).strip().lower()
        if key not in common:
            continue
        sets = {}
        too_large = False
        for b in bundles:
            q = str(b["quarter"])
            actual = schemas[q][key]
            vals = set(b["data"][actual].dropna().astype(str).str.strip().unique().tolist())
            if len(vals) > 50:
                too_large = True
                break
            sets[q] = vals
        if too_large or not sets:
            continue
        union = set().union(*sets.values())
        intersection = set.intersection(*sets.values()) if sets else set()
        if union and union != intersection:
            category_drift.append({"column": str(col), "shared": len(intersection), "total_distinct": len(union), "by_quarter": {q: len(v) for q, v in sets.items()}})
    if category_drift:
        add("Category consistency", "Review", f"Category membership changes in {len(category_drift)} common dimension(s).", min(10, 2 + 2 * len(category_drift)), {"dimensions": category_drift[:8]})
    else:
        add("Category consistency", "Good", "No material low-cardinality category drift detected.", 0)

    score = round(max(0.0, 100.0 - deductions), 1)
    status = "Strong" if score >= 90 else "Good" if score >= 75 else "Watch" if score >= 55 else "Limited"
    return {"quarters": ordered, "rows": rows, "score": score, "status": status, "checks": checks}


def compact_quality_evidence(quality):
    if not quality:
        return None
    return {
        "status": quality.get("status"),
        "score": quality.get("score"),
        "analysis_guidance": quality.get("analysis_guidance"),
        "checks": [
            {"check": x.get("check"), "status": x.get("status"), "detail": x.get("detail")}
            for x in (quality.get("checks", []) or [])[:8]
        ],
    }

# ============================================================
# SMART DASHBOARD LAYOUT HELPERS
# ============================================================

def _compact_value(value):
    """Format management values without changing the underlying calculation."""
    if value is None or pd.isna(value):
        return "n/a"
    value = float(value)
    absolute = abs(value)
    if absolute >= 1_000_000_000:
        return f"{value / 1_000_000_000:,.1f}B"
    if absolute >= 1_000_000:
        return f"{value / 1_000_000:,.1f}M"
    if absolute >= 1_000:
        return f"{value / 1_000:,.1f}K"
    return f"{value:,.2f}"


def build_key_management_signals(summary, comparison=None):
    """Return concise management signals using Python-calculated facts only."""
    context = summary.get("analytical_context", {}) or {}
    financial = context.get("financial_variance", {}) or {}
    trends = context.get("trends", {}) or {}
    breakdowns = context.get("top_breakdowns", {}) or {}
    quality = context.get("data_quality", {}) or {}
    signals = []

    actual_budget = financial.get("actual_vs_budget")
    if actual_budget:
        variance = actual_budget.get("variance")
        pct = actual_budget.get("variance_pct")
        actual_name = actual_budget.get("actual_metric", "Actual")
        budget_name = actual_budget.get("budget_metric", "Budget")
        if variance is not None:
            direction = "above" if variance >= 0 else "below"
            pct_text = f" ({abs(pct):.1f}%)" if pct is not None else ""
            signals.append(
                f"**{actual_name} vs {budget_name}:** {_compact_value(abs(variance))} "
                f"{direction} plan{pct_text}."
            )

    revenue_costs = financial.get("revenue_vs_costs")
    if revenue_costs:
        spread = revenue_costs.get("spread")
        spread_pct = revenue_costs.get("spread_pct_of_revenue")
        if spread is not None:
            signals.append(
                "**Revenue less costs:** "
                f"{_compact_value(spread)}"
                + (f" ({spread_pct:.1f}% of revenue)." if spread_pct is not None else ".")
            )

    # Largest absolute trend movement.
    trend_candidates = []
    for metric, details in trends.items():
        pct = details.get("change_pct")
        if pct is not None:
            trend_candidates.append((abs(pct), metric, pct))
    if trend_candidates:
        _, metric, pct = max(trend_candidates)
        direction = "increased" if pct >= 0 else "decreased"
        signals.append(f"**Largest trend movement:** {metric} {direction} {abs(pct):.1f}% across the observed period.")

    # Concentration signal from the strongest available category breakdown.
    for key, details in breakdowns.items():
        top = details.get("top", [])
        if len(top) >= 2:
            values = [float(item.get("value", 0) or 0) for item in top]
            positive_total = sum(v for v in values if v > 0)
            if positive_total > 0:
                top_two = sum(sorted((v for v in values if v > 0), reverse=True)[:2])
                share = top_two / positive_total * 100
                signals.append(f"**Contribution concentration:** the top two visible contributors in {key} represent {share:.1f}% of the top-five positive contribution.")
            break

    if comparison is not None:
        metrics = comparison.get("metrics")
        if metrics is not None and not metrics.empty and "Change %" in metrics.columns:
            comparable = metrics.dropna(subset=["Change %"]).copy()
            if not comparable.empty:
                idx = comparable["Change %"].abs().idxmax()
                row = comparable.loc[idx]
                signals.append(
                    f"**Quarter comparison:** {row['Metric']} moved {row['Change %']:+.1f}% "
                    f"from {comparison['first_quarter']} to {comparison['second_quarter']}."
                )

    missing_rate = quality.get("missing_rate_pct")
    if missing_rate is not None and missing_rate > 0:
        signals.append(f"**Data quality:** {missing_rate:.1f}% of cells are missing; interpret affected measures with care.")

    if not signals:
        signals.append("No material variance, trend, concentration or data-quality signal was detected from the available calculated facts.")

    return signals[:5]



def _metric_business_importance(metric, watchlist=None):
    """Heuristic business relevance score used only for prioritisation.

    This is intentionally not an accounting materiality threshold. It helps the
    interface rank competing analytical signals without asking AI to decide which
    metric is important.
    """
    lower = str(metric or "").strip().lower()
    tiers = (
        (100.0, ("revenue", "sales", "profit", "ebitda", "earnings", "cash", "actual")),
        (92.0, ("budget", "forecast", "cost", "expense", "opex", "capex", "margin")),
        (82.0, ("headcount", "fte", "volume", "price", "units")),
        (72.0, ("variance", "ratio", "rate", "%")),
    )
    base_score = 60.0
    for score, tokens in tiers:
        if any(token in lower for token in tokens):
            base_score = score
            break

    # A watchlist is a user preference, not an accounting materiality rule.
    # Give it a modest bounded boost so genuinely stronger evidence can still outrank it.
    normalized_watchlist = {str(item).strip().lower() for item in (watchlist or [])}
    if lower in normalized_watchlist:
        base_score = min(100.0, base_score + 10.0)
    return base_score


def _priority_score(magnitude, importance, consistency):
    """Blend bounded evidence components into a transparent 0-100 priority score."""
    magnitude = max(0.0, min(float(magnitude or 0.0), 100.0))
    importance = max(0.0, min(float(importance or 0.0), 100.0))
    consistency = max(0.0, min(float(consistency or 0.0), 100.0))
    return round(0.50 * magnitude + 0.30 * importance + 0.20 * consistency, 1)


@st.cache_data(ttl=900, show_spinner=False, max_entries=48)
def build_materiality_ranking(summary, forecast_result=None, comparison=None, limit=8):
    """Rank decision signals using deterministic evidence only.

    Score = 50% magnitude + 30% business importance + 20% consistency/reliability.
    The score is a decision-priority heuristic, not a statutory/accounting
    materiality threshold and not an AI-generated judgement.
    """
    context = summary.get("analytical_context", {}) or {}
    financial = context.get("financial_variance", {}) or {}
    trends = context.get("trends", {}) or {}
    anomalies = context.get("anomalies", {}) or {}
    totals = context.get("metric_totals", {}) or {}
    drivers = (context.get("driver_analysis", {}) or {}).get("metrics", {}) or {}
    quality = context.get("data_quality", {}) or {}
    watchlist = list(summary.get("management_watchlist", []) or [])
    normalized_watchlist = {str(item).strip().lower() for item in watchlist}
    signals = []

    def add(signal_type, title, detail, metric, magnitude, consistency, source):
        importance = _metric_business_importance(metric, watchlist=watchlist)
        score = _priority_score(magnitude, importance, consistency)
        signals.append({
            "type": signal_type,
            "title": str(title),
            "detail": str(detail),
            "metric": str(metric) if metric is not None else None,
            "score": score,
            "magnitude_component": round(float(magnitude or 0.0), 1),
            "importance_component": round(float(importance), 1),
            "consistency_component": round(float(consistency or 0.0), 1),
            "source": source,
            "watchlisted": str(metric).strip().lower() in normalized_watchlist if metric is not None else False,
        })

    actual_budget = financial.get("actual_vs_budget") or {}
    if actual_budget.get("variance_pct") is not None:
        pct = float(actual_budget["variance_pct"])
        metric = actual_budget.get("actual_metric", "Actual")
        add(
            "Financial variance",
            f"{metric} vs {actual_budget.get('budget_metric', 'Budget')}",
            f"{pct:+.1f}% versus plan",
            metric,
            min(abs(pct) * 2.0, 100.0),
            90.0,
            "calculated financial variance",
        )

    for metric, details in trends.items():
        pct = details.get("change_pct")
        if pct is None:
            continue
        pct = float(pct)
        consistency = float(details.get("direction_consistency_pct", 50.0) or 50.0)
        observations = int(details.get("observation_count", 0) or 0)
        # Short histories are useful but should not outrank equally large moves
        # supported by more repeated observations.
        history_factor = min(1.0, max(observations - 1, 1) / 3.0)
        consistency = consistency * (0.65 + 0.35 * history_factor)
        add(
            "Trend",
            f"{metric} trend",
            f"{pct:+.1f}% across the observed period; direction consistency {details.get('direction_consistency_pct', 0):.0f}%",
            metric,
            min(abs(pct) * 2.0, 100.0),
            consistency,
            "calculated trend",
        )

    for metric, details in anomalies.items():
        count = int(details.get("count", 0) or 0)
        total = int((totals.get(metric, {}) or {}).get("non_null_count", 0) or 0)
        rate = (count / total * 100.0) if total else 0.0
        add(
            "Anomaly screen",
            f"{metric} anomaly candidates",
            f"{count} IQR-screened value{'s' if count != 1 else ''} ({rate:.1f}% of non-null observations)",
            metric,
            min(rate * 5.0, 100.0),
            70.0,
            "deterministic IQR screening",
        )

    # One strongest decomposition signal per metric avoids flooding the ranking
    # with several dimensions describing the same underlying movement.
    for metric, dimensions in drivers.items():
        best = None
        for dimension, details in (dimensions or {}).items():
            pct = details.get("total_change_pct")
            driver_rows = details.get("drivers", []) or []
            if pct is None or not driver_rows:
                continue
            total_abs = sum(abs(float(item.get("change", 0) or 0)) for item in driver_rows)
            top_abs = abs(float(driver_rows[0].get("change", 0) or 0))
            concentration = (top_abs / total_abs * 100.0) if total_abs else 0.0
            candidate = (abs(float(pct)), concentration, dimension, details)
            if best is None or candidate[:2] > best[:2]:
                best = candidate
        if best:
            _, concentration, dimension, details = best
            pct = float(details.get("total_change_pct") or 0.0)
            top_driver = (details.get("drivers", []) or [{}])[0]
            add(
                "Driver concentration",
                f"{metric} by {dimension}",
                f"{pct:+.1f}% net movement; largest contributor {top_driver.get('category')} represents {concentration:.1f}% of absolute displayed movement",
                metric,
                min(abs(pct) * 1.5, 100.0),
                min(50.0 + concentration / 2.0, 100.0),
                "calculated driver decomposition",
            )

    missing_rate = quality.get("missing_rate_pct")
    if missing_rate is not None and float(missing_rate) > 0:
        missing_rate = float(missing_rate)
        add(
            "Data quality",
            "Missing data",
            f"{missing_rate:.1f}% of cells are missing",
            "Data quality",
            min(missing_rate * 3.0, 100.0),
            85.0,
            "calculated data-quality check",
        )

    if comparison is not None:
        metrics = comparison.get("metrics")
        if metrics is not None and not metrics.empty and "Change %" in metrics.columns:
            for _, row in metrics.dropna(subset=["Change %"]).iterrows():
                pct = float(row["Change %"])
                metric = row.get("Metric", "Metric")
                add(
                    "Quarter comparison",
                    f"{metric} quarter-on-quarter",
                    f"{pct:+.1f}% from {comparison.get('first_quarter')} to {comparison.get('second_quarter')}",
                    metric,
                    min(abs(pct) * 2.0, 100.0),
                    85.0,
                    "calculated quarter comparison",
                )

    if forecast_result and forecast_result.get("forecasts"):
        backtests = forecast_result.get("backtests", {}) or {}
        for metric, frame in (forecast_result.get("forecasts", {}) or {}).items():
            actual = frame[frame["Series"] == "Actual"]
            projected = frame[frame["Series"] == "Forecast"]
            if actual.empty or projected.empty:
                continue
            last_actual = float(actual.iloc[-1]["Value"])
            end_forecast = float(projected.iloc[-1]["Value"])
            pct = ((end_forecast - last_actual) / last_actual * 100.0) if last_actual else None
            if pct is None:
                continue
            backtest = backtests.get(metric, {}) or {}
            confidence = str(backtest.get("confidence", "Limited"))
            wape = backtest.get("wape_pct")
            reliability = 65.0 if confidence == "Moderate" else 35.0
            if wape is not None:
                reliability = min(reliability, max(15.0, 100.0 - float(wape) * 2.0))
            add(
                "Forecast outlook",
                f"{metric} forecast",
                f"{pct:+.1f}% from latest actual to end of horizon; confidence {confidence}",
                metric,
                min(abs(pct) * 2.0, 100.0) * reliability / 100.0,
                reliability,
                "calculated forecast and walk-forward validation",
            )

    # Keep related signals but prevent exact duplicate titles from occupying slots.
    deduped = {}
    for item in signals:
        key = (item["type"], item["title"])
        if key not in deduped or item["score"] > deduped[key]["score"]:
            deduped[key] = item
    ranked = sorted(deduped.values(), key=lambda item: (item["score"], item["magnitude_component"]), reverse=True)
    for index, item in enumerate(ranked[:limit], start=1):
        item["rank"] = index
    return ranked[:limit]


def build_management_pulse(summary, forecast_result=None, materiality_ranking=None):
    """Return a compact, deterministic management pulse for the current state.

    No AI call is made here. The pulse only prioritizes already-calculated facts,
    so it can be rendered on every rerun without latency or hallucination risk.
    """
    context = summary.get("analytical_context", {}) or {}
    trends = context.get("trends", {}) or {}
    quality = context.get("data_quality", {}) or {}
    financial = context.get("financial_variance", {}) or {}

    priority_signals = (
        list(materiality_ranking)[:3]
        if materiality_ranking is not None
        else build_materiality_ranking(summary, forecast_result=forecast_result, limit=3)
    )
    if priority_signals:
        top_signal = priority_signals[0]
        trend_text = f"#{top_signal['rank']} {top_signal['title']} — {top_signal['detail']}"
        trend_question = (
            "Explain why the highest-ranked decision signal matters, using the supplied "
            "materiality ranking and underlying calculated evidence only."
        )
    else:
        trend_text = "No material decision signal is available from the current calculations."
        trend_question = "What are the most important management signals in the current data?"

    actual_budget = financial.get("actual_vs_budget") or {}
    if actual_budget and actual_budget.get("variance_pct") is not None:
        variance_text = (
            f"{actual_budget.get('actual_metric', 'Actual')} vs "
            f"{actual_budget.get('budget_metric', 'Budget')}: "
            f"{float(actual_budget['variance_pct']):+.1f}%"
        )
    else:
        quality_intel = summary.get("data_quality_intelligence", {}) or {}
        if quality_intel.get("score") is not None:
            variance_text = (
                f"Analysis confidence: {quality_intel.get('status', 'n/a')} · "
                f"{float(quality_intel['score']):.0f}/100"
            )
        else:
            missing = quality.get("missing_rate_pct")
            variance_text = (
                f"Data completeness: {100 - float(missing):.1f}%"
                if missing is not None
                else "No material variance signal calculated."
            )

    forecast_text = "Forecast board not built yet."
    forecast_question = None
    if forecast_result:
        horizon = forecast_result.get("horizon")
        quarters = forecast_result.get("quarters", []) or []
        metrics = forecast_result.get("metrics", []) or []
        confidence_labels = [
            str(details.get("confidence"))
            for details in (forecast_result.get("backtests", {}) or {}).values()
            if details and details.get("confidence")
        ]
        confidence_text = ""
        if confidence_labels:
            confidence_text = f" Confidence: {', '.join(sorted(set(confidence_labels)))}."
        forecast_text = (
            f"Forecast ready: {len(quarters)} historical quarter"
            f"{'s' if len(quarters) != 1 else ''}, {horizon} quarter"
            f"{'s' if horizon != 1 else ''} ahead, {len(metrics)} metric"
            f"{'s' if len(metrics) != 1 else ''}.{confidence_text}"
        )
        forecast_question = (
            "Summarize the forecast outlook, distinguishing historical actuals "
            "from projected values and highlighting the largest projected movement."
        )

    return {
        "trend": trend_text,
        "variance_or_quality": variance_text,
        "forecast": forecast_text,
        "trend_question": trend_question,
        "forecast_question": forecast_question,
    }


def _recent_user_question(chat_history):
    """Return the most recent user question without making an AI call."""
    for message in reversed(chat_history or []):
        if message.get("role") == "user" and str(message.get("content", "")).strip():
            return str(message.get("content")).strip()
    return ""


def _question_topic(text):
    """Coarse deterministic topic routing for next-best-question suggestions."""
    lower = str(text or "").lower()
    topic_words = {
        "forecast": ("forecast", "scenario", "outlook", "projection", "confidence", "wape", "backtest", "model"),
        "drivers": ("driver", "contributor", "contribution", "why", "cause", "variance bridge", "bridge"),
        "quality": ("quality", "missing", "duplicate", "schema", "trust", "reliable", "confidence in data"),
        "watchlist": ("watchlist", "bookmark", "priority metric", "watched"),
        "comparison": ("compare", "quarter over quarter", "quarter-on-quarter", "qoq", "versus", " vs "),
        "trend": ("trend", "increase", "decrease", "movement", "change", "peak", "low", "anomaly"),
        "priority": ("priority", "material", "important", "management", "attention"),
    }
    scores = {
        topic: sum(1 for word in words if word in lower)
        for topic, words in topic_words.items()
    }
    best = max(scores, key=scores.get) if scores else None
    return best if best and scores[best] > 0 else None


def build_contextual_ai_questions(
    summary,
    forecast_result=None,
    previous_quarter=None,
    selected_quarter=None,
    comparison_result=None,
    chat_history=None,
):
    """Return only the three most useful next questions for the current state.

    Suggestions are ranked deterministically from already-calculated evidence.
    Creating the suggestions never calls Copilot; an AI call happens only after
    the user chooses a suggestion or types a question.
    """
    pulse = build_management_pulse(summary, forecast_result)
    ranking = build_materiality_ranking(
        summary,
        forecast_result=forecast_result,
        comparison=comparison_result,
        limit=6,
    )
    context = summary.get("analytical_context", {}) or {}
    quality = context.get("data_quality", {}) or {}
    quality_gate = summary.get("analysis_confidence", {}) or {}
    watchlist = list(summary.get("management_watchlist", []) or [])
    last_question = _recent_user_question(chat_history)
    last_topic = _question_topic(last_question)

    candidates = []

    def add(label, prompt, topic, score, reason):
        # Avoid immediately repeating the same conversational lane unless it is
        # clearly the strongest unresolved decision path.
        repeat_penalty = 16 if last_topic and topic == last_topic else 0
        candidates.append({
            "label": str(label),
            "prompt": str(prompt),
            "topic": topic,
            "score": float(score) - repeat_penalty,
            "reason": str(reason),
        })

    # 1) Highest-ranked calculated signal is generally the strongest next step.
    if ranking:
        top = ranking[0]
        metric = top.get("metric") or top.get("title")
        add(
            "Explain top priority",
            f"Explain why {top.get('title')} is currently ranked as the top decision signal. "
            f"Use the supplied score components and evidence only, and tell me what management should inspect next.",
            "priority",
            100 + float(top.get("score", 0)) / 10,
            f"Current #1 deterministic signal: {top.get('title')} (score {top.get('score')}).",
        )

    # 2) A reconciled bridge is more actionable than a generic comparison.
    bridge = (comparison_result or {}).get("variance_bridge") if comparison_result else None
    if bridge:
        add(
            "Explain variance bridge",
            f"Walk me through the {bridge.get('first_quarter')} to {bridge.get('second_quarter')} variance bridge for "
            f"{bridge.get('metric')} by {bridge.get('dimension')}. Identify the largest positive and negative contributions, "
            "show how they reconcile to the ending total, and do not claim causation.",
            "drivers",
            96,
            "A reconciled cross-quarter bridge is available.",
        )

    # 3) Forecast questions adapt to what has actually been built.
    if forecast_result and forecast_result.get("forecasts"):
        backtests = forecast_result.get("backtests", {}) or {}
        confidence_pairs = [
            (metric, details.get("confidence"), details.get("wape_pct"))
            for metric, details in backtests.items()
            if details
        ]
        limited = [item for item in confidence_pairs if item[1] == "Limited"]
        if limited:
            metric, confidence, wape = limited[0]
            add(
                "Test forecast reliability",
                f"Assess how much reliance management should place on the {metric} forecast. Use its walk-forward backtest, "
                "selected model and supplied confidence label. Explain the limitations without upgrading the confidence.",
                "forecast",
                94,
                f"{metric} currently has Limited forecast confidence.",
            )
        else:
            add(
                "Compare forecast scenarios",
                "Compare the current baseline, upside and downside forecast scenarios. Identify the largest end-of-horizon "
                "spread and explain the management implications without assigning scenario probabilities.",
                "forecast",
                90,
                "A forecast with scenario paths is currently available.",
            )

    # 4) Deterministic driver analysis where available.
    driver_metrics = ((context.get("driver_analysis", {}) or {}).get("metrics", {}) or {})
    if driver_metrics:
        first_metric = next(iter(driver_metrics))
        dimensions = driver_metrics.get(first_metric, {}) or {}
        if dimensions:
            first_dimension = next(iter(dimensions))
            add(
                "Explain key drivers",
                f"Explain the calculated contributors to the change in {first_metric} by {first_dimension}. "
                "Separate arithmetic contribution from causal interpretation and identify the strongest offsetting movements.",
                "drivers",
                88,
                f"Driver decomposition is available for {first_metric} by {first_dimension}.",
            )

    # 5) Data-quality caveats should jump ahead when analytical readiness is weak.
    quality_status = str(quality_gate.get("status", quality.get("status", "")) or "")
    if quality_status in {"Watch", "Limited"}:
        add(
            "Check analysis reliability",
            "Explain which current data-quality issues most limit the reliability of the trends, comparisons or forecasts, "
            "and which issue should be resolved first. Use only the deterministic quality checks.",
            "quality",
            99 if quality_status == "Limited" else 92,
            f"Analysis confidence is currently {quality_status}.",
        )
    elif float(quality.get("missing_rate_pct", 0) or 0) > 0:
        add(
            "Review data quality",
            "Summarize the current data-quality findings and tell me whether any of them materially change how I should interpret the analysis.",
            "quality",
            72,
            "The dataset contains some missing values.",
        )

    # 6) Watchlist is useful, but user preference should not dominate evidence.
    if watchlist:
        focus = ", ".join(map(str, watchlist[:3]))
        add(
            "Review watchlist",
            f"Review my management watchlist ({focus}). Which watched metric currently deserves the most attention, "
            "what evidence supports that, and is any non-watchlisted signal more important?",
            "watchlist",
            80,
            "Management watchlist preferences are available.",
        )

    # 7) Offer a comparison if one has not yet been created.
    if not comparison_result and previous_quarter and selected_quarter:
        add(
            "Compare prior quarter",
            f"Compare {selected_quarter} and {previous_quarter}. Focus on the most material KPI and category movements and distinguish arithmetic drivers from causes.",
            "comparison",
            76,
            "A prior quarter is available for comparison.",
        )

    # 8) General deterministic ranking remains a useful fallback.
    add(
        "Prioritize watch items",
        "Explain the deterministic materiality ranking in priority order. Use the supplied scores and evidence only; "
        "do not invent a different ranking or accounting materiality threshold.",
        "priority",
        68,
        "A deterministic decision-priority ranking is available.",
    )

    # Deduplicate near-identical prompts/topics, retaining the strongest version.
    deduped = {}
    for item in candidates:
        key = (item["label"], item["topic"])
        if key not in deduped or item["score"] > deduped[key]["score"]:
            deduped[key] = item

    ranked = sorted(deduped.values(), key=lambda item: item["score"], reverse=True)

    # Keep the surface intentionally small. Prefer topic diversity when possible.
    chosen = []
    seen_topics = set()
    for item in ranked:
        if item["topic"] in seen_topics and len(chosen) < 2:
            continue
        chosen.append(item)
        seen_topics.add(item["topic"])
        if len(chosen) == 3:
            break
    if len(chosen) < 3:
        for item in ranked:
            if item not in chosen:
                chosen.append(item)
            if len(chosen) == 3:
                break

    return chosen


@st.cache_data(ttl=1800, show_spinner=False, max_entries=32)
def build_dashboard_panels(summary):
    """Build fast, deterministic layout panels from pre-calculated facts."""
    context = summary.get("analytical_context", {}) or {}
    financial = context.get("financial_variance", {}) or {}
    trends = context.get("trends", {}) or {}
    breakdowns = context.get("top_breakdowns", {}) or {}
    quality = context.get("data_quality", {}) or {}

    variance_items = []
    risk_items = []
    questions = []

    actual_budget = financial.get("actual_vs_budget")
    if actual_budget:
        variance = actual_budget.get("variance")
        variance_pct = actual_budget.get("variance_pct")
        actual_name = actual_budget.get("actual_metric", "Actual")
        budget_name = actual_budget.get("budget_metric", "Budget")
        direction = "above" if (variance or 0) >= 0 else "below"
        pct_text = (
            f" ({abs(variance_pct):.1f}%)"
            if variance_pct is not None
            else ""
        )
        variance_items.append(
            f"**{actual_name} vs {budget_name}:** "
            f"{_compact_value(abs(variance))} {direction} plan{pct_text}."
        )
        if variance is not None and variance < 0:
            risk_items.append(
                f"{actual_name} is below {budget_name} by "
                f"{_compact_value(abs(variance))}{pct_text}."
            )
        questions.append(
            f"What are the main drivers of the {actual_name} versus {budget_name} variance?"
        )

    revenue_costs = financial.get("revenue_vs_costs")
    if revenue_costs:
        spread = revenue_costs.get("spread")
        spread_pct = revenue_costs.get("spread_pct_of_revenue")
        variance_items.append(
            "**Revenue less costs:** "
            f"{_compact_value(spread)}"
            + (
                f" ({spread_pct:.1f}% of revenue)."
                if spread_pct is not None
                else "."
            )
        )
        if spread is not None and spread < 0:
            risk_items.append(
                "Costs exceed revenue by "
                f"{_compact_value(abs(spread))}."
            )

    strongest_negative = None
    for metric, details in trends.items():
        pct = details.get("change_pct")
        if pct is None or pct >= 0:
            continue
        if strongest_negative is None or pct < strongest_negative[1]:
            strongest_negative = (metric, pct, details)

    if strongest_negative:
        metric, pct, details = strongest_negative
        risk_items.append(
            f"{metric} declined {abs(pct):.1f}% from the first to the latest observed period."
        )
        questions.append(
            f"What explains the decline in {metric}, and is it concentrated in a specific category?"
        )

    for breakdown_name, details in breakdowns.items():
        top = details.get("top", [])
        if not top:
            continue
        total = sum(abs(float(item.get("value", 0) or 0)) for item in top)
        lead_value = abs(float(top[0].get("value", 0) or 0))
        share = (lead_value / total * 100) if total else 0
        if share >= 50:
            risk_items.append(
                f"Concentration watch: {top[0].get('category')} represents about "
                f"{share:.0f}% of the top-five total for {breakdown_name}."
            )
        questions.append(
            f"Why is {top[0].get('category')} the leading contributor to {breakdown_name}?"
        )
        break

    missing_rate = quality.get("missing_rate_pct")
    if missing_rate is not None and missing_rate > 5:
        risk_items.append(
            f"Data quality watch: {missing_rate:.1f}% of cells are missing."
        )

    if not variance_items:
        variance_items.append(
            "No supported budget/variance relationship was detected in the current dataset."
        )

    if not risk_items:
        risk_items.append(
            "No material risk signal was detected by the deterministic checks currently available."
        )

    fallback_questions = [
        "Which categories are contributing most to the current result?",
        "Which movement deserves management attention first?",
        "What should management investigate before the next reporting cycle?",
    ]
    for item in fallback_questions:
        if len(questions) >= 3:
            break
        if item not in questions:
            questions.append(item)

    return {
        "variance": variance_items[:3],
        "risks": risk_items[:3],
        "questions": questions[:3],
    }


# ============================================================
# DEFAULT MANUAL DASHBOARDS
# ============================================================

def category_is_small(
    profile,
    category,
):

    unique = (
        profile[
            "details"
        ]
        .get(
            category,
            {},
        )
        .get(
            "unique",
            999,
        )
    )


    return (
        2
        <= unique
        <= 12
    )


def manual_plan(
    mode,
    profile,
    business,
):

    numeric = (
        profile[
            "numeric"
        ]
    )

    categories = (
        profile[
            "categories"
        ]
    )

    dates = (
        profile[
            "dates"
        ]
    )


    metric = (

        business[
            "revenue"
        ]

        or business[
            "budget"
        ]

        or business[
            "costs"
        ]

        or (
            numeric[0]
            if numeric
            else None
        )
    )


    category = (

        business[
            "dimension"
        ]

        or (
            categories[0]
            if categories
            else None
        )
    )


    plan = {

        "title":
            f"{mode} Dashboard",

        "reason":
            (
                "Using the selected "
                "default dashboard."
            ),

        "kpis":
            [

                {
                    "column":
                        column,

                    "aggregation":
                        "sum",

                    "label":
                        column,
                }

                for column in numeric[
                    :4
                ]
            ],

        "charts":
            [],
    }


    if (
        mode == "Executive"
        and category
        and metric
    ):

        plan[
            "charts"
        ].append(
            {

                "type":
                    "bar",

                "title":
                    (
                        f"{metric} "
                        f"by {category}"
                    ),

                "x":
                    category,

                "y":
                    metric,

                "aggregation":
                    "sum",
            }
        )


        if category_is_small(
            profile,
            category,
        ):

            plan[
                "charts"
            ].append(
                {

                    "type":
                        "donut",

                    "title":
                        (
                            f"{metric} mix"
                        ),

                    "x":
                        category,

                    "y":
                        metric,

                    "aggregation":
                        "sum",
                }
            )


    elif (
        mode == "Trends"
        and dates
        and metric
    ):

        plan[
            "charts"
        ].append(
            {

                "type":
                    "line",

                "title":
                    (
                        f"{metric} trend"
                    ),

                "x":
                    dates[0],

                "y":
                    metric,

                "aggregation":
                    "sum",
            }
        )


    elif (
        mode == "Breakdown"
        and category
        and metric
    ):

        plan[
            "charts"
        ].append(
            {

                "type":
                    "bar",

                "title":
                    (
                        f"{metric} "
                        f"by {category}"
                    ),

                "x":
                    category,

                "y":
                    metric,

                "aggregation":
                    "sum",
            }
        )


        if category_is_small(
            profile,
            category,
        ):

            plan[
                "charts"
            ].append(
                {

                    "type":
                        "donut",

                    "title":
                        (
                            f"{metric} "
                            "distribution"
                        ),

                    "x":
                        category,

                    "y":
                        metric,

                    "aggregation":
                        "sum",
                }
            )


    return plan


# ============================================================
# WORKIVA CACHE
#
# Spreadsheet/sheet discovery is cached for 30 minutes.
# Actual sheet values are cached for 10 minutes.
# This improves speed while Refresh Workiva remains available.
# ============================================================

@st.cache_data(
    ttl=1800,
    show_spinner=False,
)
def cached_quarters():

    return (
        discover_available_quarters()
    )


@st.cache_data(
    ttl=1800,
    show_spinner=False,
)
def cached_source(quarter):

    return (
        find_quarter_source_smart(
            quarter
        )
    )


@st.cache_data(
    ttl=1800,
    show_spinner=False,
)
def cached_data(
    spreadsheet_id,
    sheet_id,
):

    response = (
        get_sheet_data(
            spreadsheet_id,
            sheet_id,
        )
    )


    return (
        clean_dataframe(
            workiva_sheet_to_dataframe(
                response
            )
        )
    )


@st.cache_data(ttl=1800, show_spinner=False)
def quarter_bundle(quarter):

    year_match = re.search(r"\b(20\d{2})\b", str(quarter))
    reporting_year = int(year_match.group(1)) if year_match else None

    source = (
        cached_source(
            quarter
        )
    )


    if not source:

        raise ValueError(
            f"No Workiva source "
            f"could be found for {quarter}."
        )


    raw = (
        cached_data(
            source[
                "spreadsheet_id"
            ],
            source[
                "sheet_id"
            ],
        )
    )


    if raw.empty:

        raise ValueError(
            f"The Workiva sheet for "
            f"{quarter} contains no usable data."
        )


    profile = (
        profile_dataframe(
            raw,
            default_year=reporting_year,
        )
    )

    data = (
        prepare_dataframe(
            raw,
            profile,
            default_year=reporting_year,
        )
    )

    business = (
        detect_business_columns(
            data
        )
    )

    # Add deterministic finance measures once per cached quarter bundle.
    # This is vectorized pandas work and does not call Copilot.
    data, derived_measures = add_financial_intelligence(
        data,
        business,
    )

    if derived_measures:
        # Derived measures are guaranteed numeric by add_financial_intelligence().
        # Extend the existing profile instead of re-running date/type inference over
        # every cell a second time. Keep the original business mapping so derived
        # labels cannot be mistaken for source concepts such as Actual.
        profile = extend_profile_with_derived_numeric(
            data,
            profile,
            derived_measures,
        )

    business["derived_measures"] = derived_measures

    summary = (
        build_data_summary(
            data,
            profile,
            quarter,
        )
    )
    summary["financial_intelligence"] = {
        "detected_concepts": {
            key: value
            for key, value in business.items()
            if key != "derived_measures" and value is not None
        },
        "derived_measures": derived_measures,
    }

    summary["analytical_context"] = build_analytical_context(
        data,
        profile,
        business,
    )
    summary["data_quality_intelligence"] = assess_data_quality_intelligence(
        data,
        profile,
        business,
        quarter,
    )
    # Expose the deterministic quality gate to every downstream AI surface.
    summary["analytical_context"]["quality_intelligence"] = compact_quality_evidence(
        summary["data_quality_intelligence"]
    )


    return {

        "quarter":
            quarter,

        "source":
            source,

        "data":
            data,

        "profile":
            profile,

        "business":
            business,

        "summary":
            summary,
    }


# ============================================================
# MULTI-QUARTER FORECASTING
# ============================================================

def quarter_period(value):
    """Return a pandas quarterly Period for labels such as Q1 2026."""
    text = str(value).upper().strip()
    match = re.search(r"\bQ([1-4])[\s\-_\/]*(20\d{2})\b", text)
    if not match:
        match = re.search(r"\b(20\d{2})[\s\-_\/]*Q([1-4])\b", text)
        if not match:
            return None
        year = int(match.group(1))
        quarter = int(match.group(2))
    else:
        quarter = int(match.group(1))
        year = int(match.group(2))
    return pd.Period(year=year, quarter=quarter, freq="Q")


def _forecast_aggregation(metric):
    """Use averages for rates/ratios and sums for additive measures."""
    lower = str(metric).strip().lower()
    average_tokens = ("%", "ratio", "margin", "rate", "average", "avg")
    return "average" if any(token in lower for token in average_tokens) else "sum"


def build_multi_quarter_history(bundles):
    """Create one quarterly observation per common numeric metric."""
    ordered = sorted(
        bundles,
        key=lambda item: quarter_period(item["quarter"]) or pd.Period("1900Q1", freq="Q"),
    )
    if not ordered:
        return pd.DataFrame(), []

    common_numeric = set(ordered[0]["profile"]["numeric"])
    for item in ordered[1:]:
        common_numeric &= set(item["profile"]["numeric"])

    # Keep the source-column order from the newest/first bundle rather than
    # alphabetizing finance measures into an unfamiliar order.
    common_numeric = [
        column
        for column in ordered[-1]["profile"]["numeric"]
        if column in common_numeric
    ]

    rows = []
    for item in ordered:
        row = {
            "Quarter": item["quarter"],
            "Quarter Period": quarter_period(item["quarter"]),
        }
        for metric in common_numeric:
            values = pd.to_numeric(item["data"][metric], errors="coerce").dropna()
            if values.empty:
                row[metric] = pd.NA
            elif _forecast_aggregation(metric) == "average":
                row[metric] = float(values.mean())
            else:
                row[metric] = float(values.sum())
        rows.append(row)

    history = pd.DataFrame(rows)
    if not history.empty:
        history = history.sort_values("Quarter Period").reset_index(drop=True)
    return history, common_numeric


def _linear_fit(values):
    """Fit a transparent least-squares line to an ordered numeric sequence."""
    values = [float(value) for value in values]
    n = len(values)
    if not values:
        return 0.0, 0.0
    if n == 1:
        return 0.0, values[0]

    x = list(range(n))
    x_mean = sum(x) / n
    y_mean = sum(values) / n
    denominator = sum((value - x_mean) ** 2 for value in x)
    slope = (
        sum((x[i] - x_mean) * (values[i] - y_mean) for i in range(n))
        / denominator
        if denominator
        else 0.0
    )
    intercept = y_mean - slope * x_mean
    return float(slope), float(intercept)


def _forecast_model_catalog(history_points):
    """Return simple, interpretable models supported by the available history."""
    models = [
        {"id": "flat", "name": "Flat baseline", "complexity": 0},
    ]
    if history_points >= 2:
        models.append({"id": "linear", "name": "Linear trend", "complexity": 1})
    if history_points >= 3:
        models.extend([
            {"id": "moving_average_2", "name": "2-quarter moving average", "complexity": 1},
            {"id": "exp_smoothing", "name": "Exponential smoothing", "complexity": 1},
        ])
    return models


def _forecast_next_value(train_values, model_id):
    """Predict one step ahead from only values available at the prediction time."""
    values = [float(value) for value in train_values]
    if not values:
        return 0.0

    if model_id == "flat" or len(values) == 1:
        return float(values[-1])

    if model_id == "linear":
        slope, intercept = _linear_fit(values)
        return float(intercept + slope * len(values))

    if model_id == "moving_average_2":
        window = values[-2:] if len(values) >= 2 else values
        return float(sum(window) / len(window))

    if model_id == "exp_smoothing":
        # Fixed alpha keeps the method transparent and avoids fitting another
        # parameter to a very short quarterly history.
        alpha = 0.5
        level = float(values[0])
        for value in values[1:]:
            level = alpha * float(value) + (1.0 - alpha) * level
        return float(level)

    return float(values[-1])


def backtest_quarter_forecast(history, metric, model_id="linear"):
    """Walk-forward validation for one explicit forecast model.

    Each known quarter is predicted using only earlier observations. This is the
    same information boundary used by the model-selection routine and prevents
    look-ahead bias.
    """
    working = history[["Quarter", "Quarter Period", metric]].dropna().copy()
    if working.empty or len(working) < 2:
        return {
            "model_id": model_id,
            "tests": pd.DataFrame(),
            "history_points": int(len(working)),
            "backtest_points": 0,
            "mae": None,
            "rmse": None,
            "mape_pct": None,
            "wape_pct": None,
            "smape_pct": None,
            "confidence": "Limited",
            "confidence_reason": "At least two historical quarters are needed for a backtest.",
        }

    values = [float(value) for value in working[metric].tolist()]
    rows = []

    for test_index in range(1, len(values)):
        train_values = values[:test_index]
        prediction = _forecast_next_value(train_values, model_id)
        actual = values[test_index]
        error = prediction - actual
        abs_error = abs(error)
        ape = (abs_error / abs(actual) * 100) if actual else None
        smape_denominator = abs(actual) + abs(prediction)
        smape = 200 * abs_error / smape_denominator if smape_denominator else 0.0
        rows.append({
            "Quarter": str(working.iloc[test_index]["Quarter"]),
            "Actual": float(actual),
            "Backtest Forecast": float(prediction),
            "Error": float(error),
            "Absolute Error": float(abs_error),
            "Absolute % Error": float(ape) if ape is not None else None,
            "sMAPE %": float(smape),
            "Training Quarters": int(test_index),
        })

    tests = pd.DataFrame(rows)
    abs_errors = tests["Absolute Error"]
    mae = float(abs_errors.mean())
    rmse = float((tests["Error"].pow(2).mean()) ** 0.5)
    nonzero_ape = tests["Absolute % Error"].dropna()
    mape = float(nonzero_ape.mean()) if not nonzero_ape.empty else None
    denominator = float(tests["Actual"].abs().sum())
    wape = float(abs_errors.sum() / denominator * 100) if denominator else None
    smape = float(tests["sMAPE %"].mean())

    backtest_points = len(tests)
    if backtest_points < 2:
        confidence = "Limited"
        confidence_reason = (
            "Only one historical holdout can be tested, so forecast reliability cannot be established."
        )
    elif wape is None:
        confidence = "Limited"
        confidence_reason = "Relative forecast error cannot be calculated reliably for this series."
    elif backtest_points >= 3 and wape <= 20:
        confidence = "Moderate"
        confidence_reason = (
            f"Walk-forward WAPE is {wape:.1f}% across {backtest_points} historical holdouts. "
            "Confidence is capped at Moderate because the selected history contains at most four quarters."
        )
    elif backtest_points >= 2 and wape <= 10:
        confidence = "Moderate"
        confidence_reason = (
            f"Walk-forward WAPE is {wape:.1f}%, but only {backtest_points} historical holdouts are available. "
            "Confidence is capped at Moderate because the sample is short."
        )
    else:
        confidence = "Limited"
        confidence_reason = (
            f"Walk-forward WAPE is {wape:.1f}% across {backtest_points} historical holdout"
            f"{'s' if backtest_points != 1 else ''}; treat the projection as directional."
        )

    return {
        "model_id": model_id,
        "tests": tests,
        "history_points": int(len(working)),
        "backtest_points": int(backtest_points),
        "mae": mae,
        "rmse": rmse,
        "mape_pct": mape,
        "wape_pct": wape,
        "smape_pct": smape,
        "confidence": confidence,
        "confidence_reason": confidence_reason,
    }


def select_forecast_model(history, metric):
    """Choose an interpretable model using walk-forward error plus simplicity.

    Because the UI currently allows at most four historical quarters, this is a
    conservative selector rather than an AutoML routine. Models are ranked on
    WAPE when available (otherwise MAE), with a small complexity penalty and a
    simplicity tie-break so tiny apparent improvements do not trigger a more
    complex model.
    """
    working = history[["Quarter", "Quarter Period", metric]].dropna().copy()
    history_points = int(len(working))
    catalog = _forecast_model_catalog(history_points)

    if history_points <= 1:
        return {
            "model_id": "flat",
            "model_name": "Flat baseline",
            "reason": "Only one usable historical quarter is available, so no trend can be estimated.",
            "candidate_scores": [],
        }

    candidates = []
    for model in catalog:
        result = backtest_quarter_forecast(history, metric, model["id"])
        wape = result.get("wape_pct")
        mae = result.get("mae")
        # Relative error is easier to compare across measures. A tiny penalty
        # favours simpler models when backtest performance is essentially tied.
        if wape is not None:
            raw_score = float(wape)
            score_basis = "WAPE"
        elif mae is not None:
            raw_score = float(mae)
            score_basis = "MAE"
        else:
            raw_score = float("inf")
            score_basis = "insufficient"
        selection_score = raw_score + float(model["complexity"]) * 0.25
        candidates.append({
            "model_id": model["id"],
            "model_name": model["name"],
            "complexity": model["complexity"],
            "selection_score": selection_score,
            "score_basis": score_basis,
            "wape_pct": wape,
            "mae": mae,
            "backtest_points": result.get("backtest_points", 0),
        })

    viable = [item for item in candidates if item["selection_score"] != float("inf")]
    if not viable:
        chosen_id = "flat"
    else:
        viable.sort(key=lambda item: (item["selection_score"], item["complexity"], item["model_name"]))
        best = viable[0]
        # If another simpler model is within 1 percentage point WAPE of the best,
        # prefer the simpler model to reduce overfitting to only a few holdouts.
        chosen = best
        if best.get("wape_pct") is not None:
            near = [
                item for item in viable
                if item.get("wape_pct") is not None
                and float(item["wape_pct"]) <= float(best["wape_pct"]) + 1.0
            ]
            if near:
                chosen = sorted(near, key=lambda item: (item["complexity"], item["selection_score"]))[0]
        chosen_id = chosen["model_id"]

    chosen_model = next(item for item in catalog if item["id"] == chosen_id)
    chosen_candidate = next((item for item in candidates if item["model_id"] == chosen_id), None)
    if chosen_candidate and chosen_candidate.get("wape_pct") is not None:
        reason = (
            f"{chosen_model['name']} was selected from {len(catalog)} supported simple models using "
            f"walk-forward validation (WAPE {chosen_candidate['wape_pct']:.1f}%). "
            "When models perform similarly, the simpler method is preferred."
        )
    else:
        reason = (
            f"{chosen_model['name']} was selected because the available history is too short "
            "for a stronger comparative validation result."
        )

    return {
        "model_id": chosen_id,
        "model_name": chosen_model["name"],
        "reason": reason,
        "candidate_scores": candidates,
    }


def model_quarter_forecast(history, metric, horizon, model_id):
    """Generate actual + future values with the selected transparent model."""
    working = history[["Quarter", "Quarter Period", metric]].dropna().copy()
    if working.empty:
        return pd.DataFrame()

    values = [float(value) for value in working[metric].tolist()]
    rows = [
        {
            "Quarter": str(row["Quarter"]),
            "Quarter Period": row["Quarter Period"],
            "Value": float(row[metric]),
            "Series": "Actual",
        }
        for _, row in working.iterrows()
    ]

    last_period = working["Quarter Period"].iloc[-1]
    forecast_values = list(values)
    for step in range(1, int(horizon) + 1):
        period = last_period + step
        estimate = _forecast_next_value(forecast_values, model_id)
        rows.append({
            "Quarter": f"Q{period.quarter} {period.year}",
            "Quarter Period": period,
            "Value": float(estimate),
            "Series": "Forecast",
        })
        # Recursive forecasting is needed for moving-average/smoothing models.
        # Linear forecasts also remain coherent under this one-step recursion.
        forecast_values.append(float(estimate))

    return pd.DataFrame(rows)

def forecast_trend_commentary(history, metric, forecast_frame):
    """Deterministic trend commentary so the board is useful without an AI call."""
    actual = forecast_frame[forecast_frame["Series"] == "Actual"]
    projected = forecast_frame[forecast_frame["Series"] == "Forecast"]
    if actual.empty or projected.empty:
        return "There is not enough usable data to describe this forecast."

    latest = float(actual["Value"].iloc[-1])
    future = float(projected["Value"].iloc[-1])
    change = future - latest
    change_pct = (change / abs(latest) * 100) if latest else None

    if len(actual) == 1:
        return (
            "Only one historical quarter is selected, so the projection uses the "
            "latest value as a flat baseline. Add another quarter to estimate a trend."
        )

    first = float(actual["Value"].iloc[0])
    historical_change = latest - first
    direction = "upward" if historical_change > 0 else "downward" if historical_change < 0 else "flat"

    actual_values = actual["Value"].tolist()
    deltas = [actual_values[i] - actual_values[i - 1] for i in range(1, len(actual_values))]
    reversals = sum(
        1 for i in range(1, len(deltas))
        if deltas[i] != 0 and deltas[i - 1] != 0 and (deltas[i] > 0) != (deltas[i - 1] > 0)
    )
    stability = "volatile" if reversals else "directionally consistent"

    projected_text = f"{change:+,.2f}"
    if change_pct is not None:
        projected_text += f" ({change_pct:+.1f}%)"

    return (
        f"The selected history is {direction} and {stability}. "
        f"The selected-model projection moves {metric} by {projected_text} from the latest "
        "actual quarter to the end of the forecast horizon. This is a data-driven extrapolation "
        "of the selected quarterly history, not a causal forecast."
    )


def build_forecast_scenarios(forecast_frame, sensitivity_pct):
    """Build transparent baseline/upside/downside cases from the baseline forecast.

    Scenario values are deterministic. For forecast step n, upside/downside apply a
    cumulative +/- sensitivity adjustment to the baseline value. Historical actuals
    are unchanged. "Upside" means a numerically higher value, which is not
    necessarily economically favourable for cost/expense measures.
    """
    if forecast_frame is None or forecast_frame.empty:
        return pd.DataFrame()

    sensitivity = max(float(sensitivity_pct), 0.0) / 100.0
    actual = forecast_frame[forecast_frame["Series"] == "Actual"].copy()
    projected = forecast_frame[forecast_frame["Series"] == "Forecast"].copy()

    rows = []
    for _, row in actual.iterrows():
        rows.append({
            "Quarter": row["Quarter"],
            "Quarter Period": row["Quarter Period"],
            "Value": float(row["Value"]),
            "Scenario": "Actual",
        })

    for step, (_, row) in enumerate(projected.iterrows(), start=1):
        baseline = float(row["Value"])
        rows.extend([
            {
                "Quarter": row["Quarter"],
                "Quarter Period": row["Quarter Period"],
                "Value": baseline,
                "Scenario": "Baseline",
            },
            {
                "Quarter": row["Quarter"],
                "Quarter Period": row["Quarter Period"],
                "Value": baseline * ((1.0 + sensitivity) ** step),
                "Scenario": "Upside",
            },
            {
                "Quarter": row["Quarter"],
                "Quarter Period": row["Quarter Period"],
                "Value": baseline * ((1.0 - sensitivity) ** step),
                "Scenario": "Downside",
            },
        ])

    return pd.DataFrame(rows)


@st.cache_data(ttl=1800, show_spinner=False, max_entries=32)
def build_forecast_result(bundles, selected_metrics, horizon, scenario_sensitivity_pct=5.0):
    history, common_metrics = build_multi_quarter_history(bundles)
    metrics = [metric for metric in selected_metrics if metric in common_metrics]
    forecasts = {}
    commentary = {}
    backtests = {}
    scenarios = {}
    model_selection = {}
    for metric in metrics:
        selection = select_forecast_model(history, metric)
        model_selection[metric] = selection
        frame = model_quarter_forecast(
            history,
            metric,
            horizon,
            selection.get("model_id", "flat"),
        )
        if frame.empty:
            continue
        forecasts[metric] = frame
        commentary[metric] = forecast_trend_commentary(history, metric, frame)
        backtests[metric] = backtest_quarter_forecast(
            history,
            metric,
            selection.get("model_id", "flat"),
        )
        scenarios[metric] = build_forecast_scenarios(frame, scenario_sensitivity_pct)
    return {
        "quarters": [item["quarter"] for item in bundles],
        "history": history,
        "common_metrics": common_metrics,
        "metrics": metrics,
        "horizon": int(horizon),
        "forecasts": forecasts,
        "commentary": commentary,
        "backtests": backtests,
        "scenarios": scenarios,
        "model_selection": model_selection,
        "scenario_sensitivity_pct": float(scenario_sensitivity_pct),
    }

def generate_forecast_ai_commentary(forecast_result):
    payload = {
        "historical_quarters": forecast_result.get("quarters", []),
        "forecast_horizon_quarters": forecast_result.get("horizon", 0),
        "scenario_sensitivity_pct": forecast_result.get("scenario_sensitivity_pct", 0),
        "scenario_definition": (
            "Baseline is the selected-model forecast. Upside is a numerically higher case and Downside "
            "a numerically lower case using the explicit per-quarter sensitivity assumption; "
            "for costs/expenses, a numerically higher Upside is not necessarily favourable."
        ),
        "metrics": {},
    }
    for metric, frame in forecast_result.get("forecasts", {}).items():
        backtest = (forecast_result.get("backtests", {}) or {}).get(metric, {}) or {}
        selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
        tests = backtest.get("tests")
        scenario_frame = (forecast_result.get("scenarios", {}) or {}).get(metric)
        payload["metrics"][metric] = {
            "selected_model": {
                "model_id": selection.get("model_id"),
                "model_name": selection.get("model_name"),
                "reason": selection.get("reason"),
                "candidate_scores": selection.get("candidate_scores", []),
            },
            "series": frame[["Quarter", "Value", "Series"]].to_dict(orient="records"),
            "scenarios": (
                scenario_frame[["Quarter", "Value", "Scenario"]].to_dict(orient="records")
                if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty
                else []
            ),
            "backtest": {
                "history_points": backtest.get("history_points"),
                "backtest_points": backtest.get("backtest_points"),
                "mae": backtest.get("mae"),
                "rmse": backtest.get("rmse"),
                "mape_pct": backtest.get("mape_pct"),
                "wape_pct": backtest.get("wape_pct"),
                "smape_pct": backtest.get("smape_pct"),
                "confidence": backtest.get("confidence"),
                "confidence_reason": backtest.get("confidence_reason"),
                "tests": (
                    tests.to_dict(orient="records")
                    if isinstance(tests, pd.DataFrame) and not tests.empty
                    else []
                ),
            },
        }

    prompt = f"""
You are an executive FP&A analyst reviewing a short quarterly forecasting dashboard.

Forecast dataset:
{payload}

Rules:
- Clearly separate historical actuals from projected values.
- The projection uses the supplied selected simple forecasting model, chosen deterministically from methods supported by the available history using walk-forward validation.
- Do not invent causes, seasonality, external drivers, management actions already taken, or confidence intervals.
- Point out direction, acceleration/deceleration only if the supplied sequence supports it, and material risks/opportunities.
- Use the supplied model-selection and walk-forward backtest results when discussing forecast reliability. Never claim the selected model is globally optimal; it is only the best supported simple candidate under the available short history.
- Never describe confidence as higher than the supplied confidence label.
- If history is too short, backtest error is high, or the supplied label is Limited, explicitly say reliability is limited.
- Backtest metrics are historical diagnostics, not guarantees of future accuracy.
- Scenario values are deterministic sensitivity cases, not separate predictive models or probabilities.
- Never describe Upside/Downside as likely outcomes or assign probabilities unless supplied.
- For costs and expenses, a numerically higher Upside case is not automatically favourable.
- Keep the response concise and management-oriented.

Write sections:
1. Forecast takeaway
2. Trend signals
3. Risks / watch items
4. Management questions
"""
    return grounded_ai_response(
        prompt,
        payload,
        purpose="forecast commentary",
    )



# ============================================================
# UNIFIED AI MANAGEMENT BRIEF
# ============================================================

def _forecast_management_brief_payload(forecast_result):
    """Compact forecast evidence for the unified executive brief."""
    if not forecast_result:
        return {"status": "not_built"}

    payload = {
        "status": "available",
        "historical_quarters": forecast_result.get("quarters", []),
        "forecast_horizon_quarters": forecast_result.get("horizon"),
        "scenario_sensitivity_pct": forecast_result.get("scenario_sensitivity_pct"),
        "metrics": {},
    }

    forecasts = forecast_result.get("forecasts", {}) or {}
    backtests = forecast_result.get("backtests", {}) or {}
    scenarios = forecast_result.get("scenarios", {}) or {}

    for metric in (forecast_result.get("metrics", []) or [])[:4]:
        frame = forecasts.get(metric)
        if not isinstance(frame, pd.DataFrame) or frame.empty:
            continue

        actual = frame[frame["Series"] == "Actual"]
        projected = frame[frame["Series"] == "Forecast"]
        details = {
            "latest_actual": (
                float(actual.iloc[-1]["Value"]) if not actual.empty else None
            ),
            "forecast_end": (
                float(projected.iloc[-1]["Value"]) if not projected.empty else None
            ),
            "forecast_end_quarter": (
                str(projected.iloc[-1]["Quarter"]) if not projected.empty else None
            ),
        }

        backtest = backtests.get(metric, {}) or {}
        selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
        details["selected_model"] = {
            "model_name": selection.get("model_name"),
            "reason": selection.get("reason"),
        }
        details["validation"] = {
            "confidence": backtest.get("confidence"),
            "confidence_reason": backtest.get("confidence_reason"),
            "backtest_points": backtest.get("backtest_points"),
            "mae": backtest.get("mae"),
            "wape_pct": backtest.get("wape_pct"),
            "smape_pct": backtest.get("smape_pct"),
        }

        scenario_frame = scenarios.get(metric)
        if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty:
            future = scenario_frame[scenario_frame["Scenario"] != "Actual"]
            if not future.empty:
                final_quarter = future["Quarter Period"].max()
                final_rows = future[future["Quarter Period"] == final_quarter]
                details["scenario_end_values"] = {
                    str(row["Scenario"]): float(row["Value"])
                    for _, row in final_rows.iterrows()
                }

        payload["metrics"][metric] = details

    return payload


def _comparison_management_brief_payload(comparison):
    """Compact already-calculated comparison evidence, when the user built one."""
    if not comparison:
        return {"status": "not_built"}

    metrics = comparison.get("metrics")
    metric_rows = []
    if isinstance(metrics, pd.DataFrame) and not metrics.empty:
        ranked = metrics.copy()
        if "Change %" in ranked.columns:
            ranked["_rank"] = pd.to_numeric(ranked["Change %"], errors="coerce").abs()
            ranked = ranked.sort_values("_rank", ascending=False).drop(columns=["_rank"])
        metric_rows = ranked.head(6).to_dict(orient="records")

    breakdown = comparison.get("breakdown")
    breakdown_rows = (
        breakdown.head(6).to_dict(orient="records")
        if isinstance(breakdown, pd.DataFrame) and not breakdown.empty
        else []
    )

    return {
        "status": "available",
        "first_quarter": comparison.get("first_quarter"),
        "second_quarter": comparison.get("second_quarter"),
        "focus_metric": comparison.get("focus_metric"),
        "largest_metric_changes": metric_rows,
        "focus_breakdown": breakdown_rows,
        "variance_bridge": comparison.get("variance_bridge"),
    }


def generate_unified_management_brief(summary, plan, forecast_result=None, comparison=None):
    """Synthesize all available deterministic evidence into one executive brief."""
    evidence = {
        "current_quarter": summary.get("quarter"),
        "dashboard": {
            "title": (plan or {}).get("title"),
            "reason": (plan or {}).get("reason"),
        },
        "current_quarter_evidence": compact_ai_context(summary),
        "current_quarter_statistics": compact_prompt_summary(summary),
        "forecasting": _forecast_management_brief_payload(forecast_result),
        "quarter_comparison": _comparison_management_brief_payload(comparison),
        "data_quality_intelligence": summary.get("data_quality_intelligence", {}),
        "cross_quarter_quality": st.session_state.get("cross_quarter_quality"),
        "management_watchlist": summary.get("management_watchlist", []),
        "materiality_ranking": build_materiality_ranking(
            summary,
            forecast_result=forecast_result,
            comparison=comparison,
            limit=8,
        ),
    }

    prompt = f"""
You are preparing a concise executive management brief for a Workiva decision-support application.

Available evidence:
{evidence}

Synthesize the strongest information across the current-quarter dashboard, trends, deterministic driver analysis,
financial variances, anomalies, data quality, forecasting, forecast scenarios/backtesting, and quarter comparison.
Only discuss a section when its evidence is available.

Required structure:
### Executive takeaway
Give the 2 to 4 most important evidence-supported conclusions.

### Performance and drivers
Explain the most material movements and where the calculated contribution occurred. Do not convert contribution into causal claims.

### Outlook and scenarios
If a forecast exists, distinguish historical actuals from projections, summarize the baseline and material scenario range, and state the supplied forecast confidence. If no forecast exists, say it has not been built yet.

### Confidence and controls
Highlight forecast-validation limitations, data-quality issues, or other reasons management should be cautious. Never increase the supplied forecast confidence label.

### Management actions
Give 2 to 4 concise questions or actions management should investigate next. These must be framed as decisions/questions, not invented facts.

Rules:
- Follow the supplied deterministic materiality ranking when deciding what to emphasize; do not invent a different priority order.
- Treat the management watchlist as a user preference. Give relevant bookmarked metrics attention, but never let preference override stronger contrary evidence.
- The materiality score is a decision-priority heuristic, not an accounting/statutory materiality threshold.
- Prioritize materiality; do not repeat every metric.
- Never invent causes, probabilities, targets, benchmarks, external events, or missing historical data.
- Scenario cases are deterministic sensitivities, not probabilities.
- A numerically higher cost/expense scenario is not automatically favourable.
- Keep the brief concise enough to read in roughly one minute.
"""

    return grounded_ai_response(
        prompt,
        evidence,
        purpose="unified management brief",
    )


# ============================================================
# QUARTER COMPARISON
# ============================================================

def extract_quarters(text):

    matches = (
        re.findall(
            (
                r"\bQ([1-4])"
                r"[\s\-_\/]*"
                r"(20\d{2})\b"
            ),
            str(text),
            flags=re.IGNORECASE,
        )
    )


    quarters = []


    for q, year in matches:

        quarter = (
            f"Q{q} {year}"
        )


        if quarter not in quarters:

            quarters.append(
                quarter
            )


    return quarters


def normalized_column_map(
    columns,
):

    return {

        str(column)
        .strip()
        .lower():
            column

        for column in columns
    }


def common_columns(
    first,
    second,
):

    first_map = (
        normalized_column_map(
            first
        )
    )

    second_map = (
        normalized_column_map(
            second
        )
    )


    result = []


    for key, first_name in (
        first_map.items()
    ):

        if key in second_map:

            result.append(
                (
                    first_name,
                    second_map[
                        key
                    ],
                )
            )


    return result


@st.cache_data(ttl=1800, show_spinner=False, max_entries=32)
def build_cross_quarter_variance_bridge(
    first_bundle,
    second_bundle,
    first_metric,
    second_metric,
    first_dimension,
    second_dimension,
    max_drivers=6,
):
    """Build an auditable bridge from one quarter total to another.

    Contributions are deterministic arithmetic: each category's second-quarter
    value minus its first-quarter value. The displayed top contributors plus an
    optional Other residual always reconcile exactly to the net movement.
    """
    if not all([first_metric, second_metric, first_dimension, second_dimension]):
        return None

    first_data = first_bundle.get("data")
    second_data = second_bundle.get("data")
    if not isinstance(first_data, pd.DataFrame) or not isinstance(second_data, pd.DataFrame):
        return None
    if first_metric not in first_data.columns or first_dimension not in first_data.columns:
        return None
    if second_metric not in second_data.columns or second_dimension not in second_data.columns:
        return None
    if not _is_additive_driver_metric(first_metric) or not _is_additive_driver_metric(second_metric):
        return None

    def grouped_values(frame, dimension, metric):
        working = pd.DataFrame({
            "Category": frame[dimension].astype("string"),
            "Value": pd.to_numeric(frame[metric], errors="coerce"),
        }).dropna(subset=["Category", "Value"])
        if working.empty:
            return pd.Series(dtype=float)
        return working.groupby("Category", dropna=False)["Value"].sum()

    first_values = grouped_values(first_data, first_dimension, first_metric)
    second_values = grouped_values(second_data, second_dimension, second_metric)
    categories = first_values.index.union(second_values.index)
    if len(categories) < 1:
        return None

    first_values = first_values.reindex(categories, fill_value=0.0)
    second_values = second_values.reindex(categories, fill_value=0.0)
    changes = second_values - first_values

    first_total = float(first_values.sum())
    second_total = float(second_values.sum())
    total_change = float(second_total - first_total)
    ranked = changes.reindex(changes.abs().sort_values(ascending=False).index)

    shown = ranked.head(max_drivers)
    rows = []
    for category, change in shown.items():
        change = float(change)
        rows.append({
            "Category": str(category),
            "First": float(first_values.get(category, 0.0)),
            "Second": float(second_values.get(category, 0.0)),
            "Contribution": change,
            "Share of net change %": (float(change / total_change * 100) if total_change else None),
        })

    residual = float(total_change - sum(row["Contribution"] for row in rows))
    if abs(residual) > max(1e-9, abs(total_change) * 1e-10):
        rows.append({
            "Category": "Other",
            "First": None,
            "Second": None,
            "Contribution": residual,
            "Share of net change %": (float(residual / total_change * 100) if total_change else None),
        })

    reconciliation = first_total + sum(row["Contribution"] for row in rows)
    return {
        "metric": str(first_metric),
        "dimension": str(first_dimension),
        "first_quarter": first_bundle.get("quarter"),
        "second_quarter": second_bundle.get("quarter"),
        "first_total": first_total,
        "second_total": second_total,
        "total_change": total_change,
        "total_change_pct": (float(total_change / first_total * 100) if first_total else None),
        "contributions": rows,
        "reconciliation_difference": float(second_total - reconciliation),
    }


def detect_comparison_focus(
    question,
    first_bundle,
    second_bundle,
    common_numeric,
):

    lower = (
        question.lower()
    )


    groups = [

        (
            [
                "budget",
                "plan",
                "forecast",
                "target",
            ],
            "budget",
        ),

        (
            [
                "cost",
                "costs",
                "expense",
                "expenses",
            ],
            "costs",
        ),

        (
            [
                "revenue",
                "sales",
                "income",
                "actual",
            ],
            "revenue",
        ),
    ]


    for words, key in groups:

        if any(
            word in lower
            for word in words
        ):

            first_metric = (
                first_bundle[
                    "business"
                ].get(
                    key
                )
            )

            second_metric = (
                second_bundle[
                    "business"
                ].get(
                    key
                )
            )


            if (
                first_metric
                and second_metric
            ):

                return (
                    first_metric,
                    second_metric,
                )


    for (
        first_name,
        second_name,
    ) in common_numeric:

        if (
            str(
                first_name
            ).lower()
            in lower
        ):

            return (
                first_name,
                second_name,
            )


    if common_numeric:

        return (
            common_numeric[0]
        )


    return (
        None,
        None,
    )


def compare_quarter_bundles(
    first_bundle,
    second_bundle,
    question="",
):

    common_numeric = (
        common_columns(
            first_bundle[
                "profile"
            ][
                "numeric"
            ],
            second_bundle[
                "profile"
            ][
                "numeric"
            ],
        )
    )


    common_categories = (
        common_columns(
            first_bundle[
                "profile"
            ][
                "categories"
            ],
            second_bundle[
                "profile"
            ][
                "categories"
            ],
        )
    )


    rows = []


    for (
        first_metric,
        second_metric,
    ) in common_numeric:

        first_total = (
            pd.to_numeric(
                first_bundle[
                    "data"
                ][
                    first_metric
                ],
                errors="coerce",
            )
            .sum()
        )


        second_total = (
            pd.to_numeric(
                second_bundle[
                    "data"
                ][
                    second_metric
                ],
                errors="coerce",
            )
            .sum()
        )


        change = (
            second_total
            - first_total
        )


        change_pct = (

            change
            / first_total
            * 100

            if first_total != 0

            else None
        )


        rows.append(
            {

                "Metric":
                    str(
                        first_metric
                    ),

                first_bundle[
                    "quarter"
                ]:
                    float(
                        first_total
                    ),

                second_bundle[
                    "quarter"
                ]:
                    float(
                        second_total
                    ),

                "Change":
                    float(
                        change
                    ),

                "Change %":
                    (
                        float(
                            change_pct
                        )

                        if change_pct
                        is not None

                        else None
                    ),
            }
        )


    metrics_df = (
        pd.DataFrame(
            rows
        )
    )


    (
        focus_first,
        focus_second,
    ) = (
        detect_comparison_focus(
            question,
            first_bundle,
            second_bundle,
            common_numeric,
        )
    )


    category_first = None

    category_second = None


    first_dimension = (
        first_bundle[
            "business"
        ].get(
            "dimension"
        )
    )

    second_dimension = (
        second_bundle[
            "business"
        ].get(
            "dimension"
        )
    )


    if (
        first_dimension
        and second_dimension
    ):

        if (
            str(
                first_dimension
            )
            .strip()
            .lower()

            ==

            str(
                second_dimension
            )
            .strip()
            .lower()
        ):

            category_first = (
                first_dimension
            )

            category_second = (
                second_dimension
            )


    if (
        category_first is None
        and common_categories
    ):

        (
            category_first,
            category_second,
        ) = (
            common_categories[0]
        )


    breakdown = None


    if (
        focus_first
        and focus_second
        and category_first
        and category_second
    ):

        first_group = (

            first_bundle[
                "data"
            ]

            .groupby(
                category_first,
                dropna=False,
            )[
                focus_first
            ]

            .sum()

            .reset_index()

            .rename(
                columns={
                    category_first:
                        "Category",

                    focus_first:
                        first_bundle[
                            "quarter"
                        ],
                }
            )
        )


        second_group = (

            second_bundle[
                "data"
            ]

            .groupby(
                category_second,
                dropna=False,
            )[
                focus_second
            ]

            .sum()

            .reset_index()

            .rename(
                columns={
                    category_second:
                        "Category",

                    focus_second:
                        second_bundle[
                            "quarter"
                        ],
                }
            )
        )


        breakdown = (
            first_group
            .merge(
                second_group,
                on="Category",
                how="outer",
            )
            .fillna(0)
        )


    variance_bridge = build_cross_quarter_variance_bridge(
        first_bundle,
        second_bundle,
        focus_first,
        focus_second,
        category_first,
        category_second,
    )


    return {

        "first_quarter":
            first_bundle[
                "quarter"
            ],

        "second_quarter":
            second_bundle[
                "quarter"
            ],

        "first_source":
            first_bundle[
                "source"
            ],

        "second_source":
            second_bundle[
                "source"
            ],

        "metrics":
            metrics_df,

        "focus_metric":
            (
                str(
                    focus_first
                )

                if focus_first

                else None
            ),

        "breakdown":
            breakdown,

        "variance_bridge":
            variance_bridge,
    }


# ============================================================
# COMPARISON COMMENTARY
# ============================================================

@st.cache_data(
    ttl=600,
    show_spinner=False,
)
def comparison_commentary(
    comparison,
):

    metrics = (
        comparison[
            "metrics"
        ].to_dict(
            orient="records"
        )
    )


    breakdown = (

        comparison[
            "breakdown"
        ]
        .head(20)
        .to_dict(
            orient="records"
        )

        if comparison[
            "breakdown"
        ] is not None

        else []
    )


    prompt = f"""
You are an executive analyst comparing two Workiva quarters.

Python has already calculated all values.

Never invent values.

Quarter 1:
{comparison["first_quarter"]}

Quarter 2:
{comparison["second_quarter"]}

Metric comparison:
{metrics}

Focus metric:
{comparison["focus_metric"]}

Category comparison:
{breakdown}

Reconciled variance bridge:
{comparison.get("variance_bridge")}

Explain the most important movements. When bridge evidence is available, explain which categories
mathematically contributed most to the movement and confirm that contribution is not proof of causation.

Highlight material increases or decreases.

Give 2 or 3 management questions.

Be concise.
"""


    evidence = {
        "first_quarter": comparison.get("first_quarter"),
        "second_quarter": comparison.get("second_quarter"),
        "metric_comparison": metrics,
        "focus_metric": comparison.get("focus_metric"),
        "category_comparison": breakdown,
        "variance_bridge": comparison.get("variance_bridge"),
    }

    return grounded_ai_response(
        prompt,
        evidence,
        purpose="quarter comparison commentary",
    )


def render_comparison(
    comparison,
    key_prefix="comparison",
):

    first_key = re.sub(r"[^A-Za-z0-9_-]+", "_", str(comparison["first_quarter"]))
    second_key = re.sub(r"[^A-Za-z0-9_-]+", "_", str(comparison["second_quarter"]))
    comparison_key = f"{key_prefix}_{first_key}_{second_key}"

    st.subheader(
        (
            f"{comparison['first_quarter']} "
            f"vs "
            f"{comparison['second_quarter']}"
        )
    )


    st.caption(
        (
            f"Source 1: "
            f"{comparison['first_source']['spreadsheet_name']} "
            f"→ "
            f"{comparison['first_source']['sheet_name']} "
            f"| "
            f"Source 2: "
            f"{comparison['second_source']['spreadsheet_name']} "
            f"→ "
            f"{comparison['second_source']['sheet_name']}"
        )
    )


    metrics = (
        comparison[
            "metrics"
        ]
    )


    if metrics.empty:

        st.warning(
            (
                "The two quarters do not "
                "contain matching numeric "
                "columns that can be compared safely."
            )
        )

        return


    cards_data = (
        metrics.head(
            4
        )
    )


    cards = (
        st.columns(
            len(
                cards_data
            )
        )
    )


    for index, (
        _,
        row,
    ) in enumerate(
        cards_data.iterrows()
    ):

        delta = (
            row[
                "Change %"
            ]
        )


        delta_text = (

            f"{delta:+.1f}%"

            if pd.notna(
                delta
            )

            else "n/a"
        )


        cards[
            index
        ].metric(
            row[
                "Metric"
            ],
            (
                f"{row[comparison['second_quarter']]:,.2f}"
            ),
            delta_text,
        )


    long_metrics = (

        metrics[
            [
                "Metric",
                comparison[
                    "first_quarter"
                ],
                comparison[
                    "second_quarter"
                ],
            ]
        ]

        .melt(
            id_vars="Metric",
            var_name="Quarter",
            value_name="Value",
        )
    )


    figure = (
        px.bar(
            long_metrics,
            x="Metric",
            y="Value",
            color="Quarter",
            barmode="group",
            title=(
                "Quarter comparison "
                "by metric"
            ),
        )
    )


    st.plotly_chart(
        figure,
        use_container_width=True,
        key=f"{comparison_key}_metrics",
    )
    render_visual_ai_action(
        {
            "title": "Quarter comparison by metric",
            "first_quarter": comparison.get("first_quarter"),
            "second_quarter": comparison.get("second_quarter"),
            "metrics": metrics.to_dict(orient="records"),
        },
        key=f"{comparison_key}_metrics_visual",
    )


    pct_data = (
        metrics
        .dropna(
            subset=[
                "Change %"
            ]
        )
        .sort_values(
            "Change %",
            ascending=False,
        )
    )


    if not pct_data.empty:

        figure = (
            px.bar(
                pct_data,
                x="Metric",
                y="Change %",
                title=(
                    "Percentage change"
                ),
            )
        )


        st.plotly_chart(
            figure,
            use_container_width=True,
            key=f"{comparison_key}_pct_change",
        )


    if (
        comparison[
            "breakdown"
        ]
        is not None
    ):

        long_breakdown = (

            comparison[
                "breakdown"
            ]

            .melt(
                id_vars="Category",
                var_name="Quarter",
                value_name="Value",
            )
        )


        figure = (
            px.bar(
                long_breakdown,
                x="Category",
                y="Value",
                color="Quarter",
                barmode="group",
                title=(
                    f"{comparison['focus_metric']} "
                    "by category"
                ),
            )
        )


        st.plotly_chart(
            figure,
            use_container_width=True,
            key=f"{comparison_key}_breakdown",
        )


    bridge = comparison.get("variance_bridge")
    if bridge and bridge.get("contributions"):
        st.markdown("#### Cross-quarter variance bridge")
        st.caption(
            "A deterministic reconciliation from the first-quarter total to the second-quarter total. "
            "Each step is a category contribution (second quarter minus first quarter); contribution is not causation."
        )

        b1, b2, b3 = st.columns(3)
        b1.metric(str(bridge.get("first_quarter")), f"{float(bridge.get('first_total', 0)):,.2f}")
        change_pct = bridge.get("total_change_pct")
        b2.metric(
            "Net change",
            f"{float(bridge.get('total_change', 0)):+,.2f}",
            f"{float(change_pct):+.1f}%" if change_pct is not None else "n/a",
        )
        b3.metric(str(bridge.get("second_quarter")), f"{float(bridge.get('second_total', 0)):,.2f}")

        contribution_rows = bridge.get("contributions", [])
        labels = [str(bridge.get("first_quarter"))] + [row["Category"] for row in contribution_rows] + [str(bridge.get("second_quarter"))]
        measures = ["absolute"] + ["relative"] * len(contribution_rows) + ["total"]
        values = [float(bridge.get("first_total", 0))] + [float(row.get("Contribution", 0)) for row in contribution_rows] + [0.0]
        texts = [f"{float(bridge.get('first_total', 0)):,.2f}"] + [f"{float(row.get('Contribution', 0)):+,.2f}" for row in contribution_rows] + [f"{float(bridge.get('second_total', 0)):,.2f}"]

        bridge_figure = go.Figure(go.Waterfall(
            x=labels,
            measure=measures,
            y=values,
            text=texts,
            textposition="outside",
            connector={"line": {"width": 1}},
        ))
        bridge_figure.update_layout(
            title=f"{bridge.get('metric')} variance bridge by {bridge.get('dimension')}",
            yaxis_title=str(bridge.get("metric")),
            showlegend=False,
        )
        st.plotly_chart(
            bridge_figure,
            use_container_width=True,
            key=f"{comparison_key}_variance_bridge",
        )
        render_visual_ai_action(
            {
                "title": f"{bridge.get('metric')} variance bridge by {bridge.get('dimension')}",
                "first_quarter": bridge.get("first_quarter"),
                "second_quarter": bridge.get("second_quarter"),
                "first_total": bridge.get("first_total"),
                "second_total": bridge.get("second_total"),
                "total_change": bridge.get("total_change"),
                "total_change_pct": bridge.get("total_change_pct"),
                "contributions": bridge.get("contributions", []),
                "reconciliation_difference": bridge.get("reconciliation_difference"),
                "interpretation_constraint": "Contribution identifies where the movement occurred, not root cause.",
            },
            key=f"{comparison_key}_variance_bridge_visual",
        )

        with st.expander("Bridge contribution details", expanded=False):
            bridge_table = pd.DataFrame(contribution_rows).rename(columns={
                "First": str(bridge.get("first_quarter")),
                "Second": str(bridge.get("second_quarter")),
            })
            st.dataframe(bridge_table, use_container_width=True, hide_index=True)
            reconciliation_difference = float(bridge.get("reconciliation_difference", 0) or 0)
            st.caption(
                f"Reconciliation check: displayed bridge differs from the ending total by {reconciliation_difference:,.6f}. "
                "A value near zero confirms the bridge reconciles."
            )


    with st.spinner(
        "Generating comparison commentary..."
    ):

        commentary = (
            comparison_commentary(
                comparison
            )
        )


    st.session_state["comparison_ai_commentary"] = commentary

    st.markdown(
        commentary
    )


# ============================================================
# EXPORT HELPERS
# ============================================================

def _figure_png_bytes(figure, width=1200, height=650):
    if figure is None:
        return None
    try:
        return figure.to_image(format="png", width=width, height=height, scale=1.4)
    except Exception:
        return None


def _dashboard_kpi_rows(data, plan):
    rows = []
    for item in plan.get("kpis", [])[:6]:
        column = item.get("column")
        if column not in data.columns:
            continue
        value = aggregate_value(data[column], item.get("aggregation", "sum"))
        rows.append([item.get("label", column), f"{value:,.2f}"])
    return rows


def _export_text(value):
    """Safe, compact text for PDF/PPTX export without interpreting markdown."""
    if value is None:
        return ""
    text = str(value).replace("**", "").replace("### ", "").replace("#### ", "")
    return text.strip()


def _pdf_paragraph(value, style):
    text = escape(_export_text(value)).replace("\n", "<br/>")
    return Paragraph(text or "Not available in this session.", style)


def _pdf_table(rows, widths=None, repeat_rows=1, font_size=8):
    safe_rows = []
    for row in rows:
        safe_rows.append([Paragraph(escape(_export_text(cell)), getSampleStyleSheet()["BodyText"]) for cell in row])
    table = Table(safe_rows, colWidths=widths, repeatRows=repeat_rows)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("FONTSIZE", (0, 0), (-1, -1), font_size),
        ("PADDING", (0, 0), (-1, -1), 4),
    ]))
    return table


def _report_sections(summary):
    panels = build_dashboard_panels(summary)
    return {
        "risks": panels.get("risks", [])[:8],
        "questions": panels.get("questions", [])[:8],
        "variance": panels.get("variance", [])[:8],
    }


def _financial_export_payload(data, business):
    metrics = []
    for key in ["revenue", "budget", "costs", "profit", "ebitda", "opex", "capex", "headcount"]:
        column = business.get(key)
        if column and column in data.columns and column not in metrics:
            values = pd.to_numeric(data[column], errors="coerce")
            metrics.append({"concept": key, "column": column, "total": float(values.sum())})

    dimension_chart = None
    dimension = business.get("dimension")
    metric = business.get("revenue") or business.get("budget") or business.get("costs")
    if dimension and metric and dimension in data.columns and metric in data.columns:
        grouped = grouped_data(data, dimension, metric, "sum")
        if not grouped.empty:
            dimension_chart = px.bar(grouped, x=dimension, y=metric, title=f"{metric} by {dimension}")

    return {
        "metrics": metrics,
        "dimension": dimension,
        "dimension_metric": metric,
        "dimension_chart": dimension_chart,
        "derived_measures": business.get("derived_measures", {}) or {},
    }


def _trend_export_figures(data, profile):
    figures = []
    if not profile.get("dates") or not profile.get("numeric"):
        return figures
    date_column = profile["dates"][0]
    for metric in profile["numeric"][:3]:
        if date_column not in data.columns or metric not in data.columns:
            continue
        trend_data = grouped_data(data, date_column, metric, "sum").sort_values(date_column)
        if trend_data.empty:
            continue
        figures.append((metric, px.line(
            trend_data,
            x=date_column,
            y=metric,
            markers=True,
            title=f"{metric} trend",
        )))
    return figures


def _comparison_export_figures(comparison):
    figures = []
    if not comparison:
        return figures
    metrics = comparison.get("metrics")
    if isinstance(metrics, pd.DataFrame) and not metrics.empty:
        first_q = comparison.get("first_quarter")
        second_q = comparison.get("second_quarter")
        long_metrics = metrics[["Metric", first_q, second_q]].melt(
            id_vars="Metric", var_name="Quarter", value_name="Value"
        )
        figures.append(("Quarter comparison by metric", px.bar(
            long_metrics, x="Metric", y="Value", color="Quarter", barmode="group",
            title="Quarter comparison by metric"
        )))
        pct_data = metrics.dropna(subset=["Change %"]).sort_values("Change %", ascending=False)
        if not pct_data.empty:
            figures.append(("Percentage change", px.bar(
                pct_data, x="Metric", y="Change %", title="Percentage change"
            )))

    breakdown = comparison.get("breakdown")
    if isinstance(breakdown, pd.DataFrame) and not breakdown.empty:
        long_breakdown = breakdown.melt(id_vars="Category", var_name="Quarter", value_name="Value")
        figures.append((f"{comparison.get('focus_metric')} by category", px.bar(
            long_breakdown, x="Category", y="Value", color="Quarter", barmode="group",
            title=f"{comparison.get('focus_metric')} by category"
        )))

    bridge = comparison.get("variance_bridge") or {}
    if bridge.get("contributions"):
        contributions = bridge.get("contributions", [])
        labels = [str(bridge.get("first_quarter"))] + [str(row.get("Category")) for row in contributions] + [str(bridge.get("second_quarter"))]
        measures = ["absolute"] + ["relative"] * len(contributions) + ["total"]
        values = [float(bridge.get("first_total", 0))] + [float(row.get("Contribution", 0)) for row in contributions] + [0.0]
        fig = go.Figure(go.Waterfall(x=labels, measure=measures, y=values, connector={"line": {"width": 1}}))
        fig.update_layout(
            title=f"{bridge.get('metric')} variance bridge by {bridge.get('dimension')}",
            yaxis_title=str(bridge.get("metric")), showlegend=False,
        )
        figures.append(("Cross-quarter variance bridge", fig))
    return figures


def _forecast_export_figures(forecast_result):
    figures = []
    if not forecast_result or not forecast_result.get("forecasts"):
        return figures
    for metric in forecast_result.get("metrics", []):
        frame = forecast_result.get("forecasts", {}).get(metric)
        if isinstance(frame, pd.DataFrame) and not frame.empty:
            chart_frame = frame.copy()
            chart_frame["Quarter order"] = chart_frame["Quarter Period"].astype(str)
            figures.append((f"{metric}: actuals and forecast", px.line(
                chart_frame, x="Quarter order", y="Value", color="Series", markers=True,
                title=f"{metric}: actuals and selected-model forecast",
                labels={"Quarter order": "Quarter", "Value": metric},
            )))
        scenario_frame = forecast_result.get("scenarios", {}).get(metric)
        if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty:
            scenario_plot = scenario_frame.copy()
            scenario_plot["Quarter order"] = scenario_plot["Quarter Period"].astype(str)
            figures.append((f"{metric}: scenarios", px.line(
                scenario_plot, x="Quarter order", y="Value", color="Scenario", markers=True,
                title=f"{metric}: baseline and sensitivity scenarios",
                labels={"Quarter order": "Quarter", "Value": metric},
            )))
    return figures


def _full_export_state(summary, profile, business, management_pulse, materiality):
    """Capture only already-built session evidence; never trigger AI or Workiva calls."""
    return {
        "management_pulse": management_pulse,
        "materiality": materiality,
        "watchlist": list(summary.get("management_watchlist", [])),
        "unified_management_brief": st.session_state.get("unified_management_brief"),
        "trends_ai_summary": st.session_state.get("trends_ai_summary"),
        "forecast_result": st.session_state.get("forecast_result"),
        "forecast_ai_commentary": st.session_state.get("forecast_ai_commentary"),
        "comparison_result": st.session_state.get("comparison_result"),
        "comparison_ai_commentary": st.session_state.get("comparison_ai_commentary"),
        "cross_quarter_quality": st.session_state.get("cross_quarter_quality"),
        "visual_ai_explanations": dict(st.session_state.get("visual_ai_explanations", {})),
        "chat_messages": list(st.session_state.get("chat_messages", []))[-12:],
        "audit_log": list(st.session_state.get("audit_log", [])),
        "profile": profile,
        "business": business,
    }


# ============================================================
# PDF EXPORT - FULL DASHBOARD PACK
# ============================================================

def build_pdf_export(quarter, source, summary, management_summary, data, plan, export_state=None):
    """Build a comprehensive board-pack from all dashboard evidence already available."""
    export_state = export_state or {}
    profile = export_state.get("profile", {}) or {}
    business = export_state.get("business", {}) or {}
    sections = _report_sections(summary)
    financial = _financial_export_payload(data, business)
    trends = _trend_export_figures(data, profile)
    forecast_result = export_state.get("forecast_result")
    comparison = export_state.get("comparison_result")

    buffer = BytesIO()
    document = SimpleDocTemplate(
        buffer, pagesize=landscape(A4), rightMargin=11 * mm, leftMargin=11 * mm,
        topMargin=10 * mm, bottomMargin=10 * mm,
    )
    styles = getSampleStyleSheet()
    story = []

    # Cover / contents
    story += [
        Paragraph(f"{escape(_export_text(plan.get('title', 'Workiva AI Dashboard')))} - {escape(_export_text(quarter))}", styles["Title"]),
        Spacer(1, 8),
        _pdf_paragraph("Comprehensive dashboard board-pack. It contains the information already available in the current app session; sections not yet built are identified explicitly.", styles["BodyText"]),
        Spacer(1, 8),
        _pdf_paragraph(
            f"Spreadsheet: {source.get('spreadsheet_name')}\nSheet: {source.get('sheet_name')}\nSource match confidence: {source.get('confidence', 'Unknown')}\nRows analysed: {summary.get('rows', 0)}\nColumns analysed: {len(summary.get('columns', []))}",
            styles["BodyText"],
        ),
        Spacer(1, 10),
        _pdf_paragraph(
            "Contents: Executive view; Overview dashboard; Financials; Trends and drivers; Forecasting/scenarios/backtesting; Quarter comparison and variance bridge; Data quality; Decision signals/watchlist; AI interaction appendix; Session audit log; Source provenance.",
            styles["BodyText"],
        ),
        PageBreak(),
    ]

    # Executive view
    story += [Paragraph("1. Executive view", styles["Heading1"])]
    pulse = export_state.get("management_pulse") or {}
    pulse_rows = [["Management pulse", "Current signal"]] + [
        ["Top ranked signal", pulse.get("trend", "Not available")],
        ["Control signal", pulse.get("variance_or_quality", "Not available")],
        ["Forecast state", pulse.get("forecast", "Not available")],
    ]
    story += [_pdf_table(pulse_rows, widths=[55 * mm, 190 * mm]), Spacer(1, 8)]

    unified = export_state.get("unified_management_brief")
    if unified:
        story += [Paragraph("Unified AI Management Brief", styles["Heading2"]), _pdf_paragraph(unified, styles["BodyText"]), Spacer(1, 8)]
    if management_summary:
        story += [Paragraph("Overview AI Management Summary", styles["Heading2"]), _pdf_paragraph(management_summary, styles["BodyText"])]
    else:
        story += [Paragraph("Overview AI Management Summary", styles["Heading2"]), _pdf_paragraph("Not generated in this session.", styles["BodyText"])]

    # Overview
    story += [PageBreak(), Paragraph("2. Overview dashboard", styles["Heading1"])]
    story += [_pdf_paragraph(f"Dashboard rationale: {plan.get('reason', '')}", styles["BodyText"]), Spacer(1, 6)]
    kpi_rows = [["KPI", "Value"]] + _dashboard_kpi_rows(data, plan)
    if len(kpi_rows) == 1:
        kpi_rows.append(["No KPI selected", ""])
    story += [_pdf_table(kpi_rows, widths=[120 * mm, 70 * mm]), Spacer(1, 8)]
    if sections["variance"]:
        story += [Paragraph("Variance / performance notes", styles["Heading2"])]
        for item in sections["variance"]:
            story.append(_pdf_paragraph(f"- {item}", styles["BodyText"]))
    story += [Spacer(1, 6), Paragraph("Risks / watch items", styles["Heading2"])]
    for item in (sections["risks"] or ["No deterministic risk signal was detected."]):
        story.append(_pdf_paragraph(f"- {item}", styles["BodyText"]))
    story += [Spacer(1, 6), Paragraph("Management questions / actions", styles["Heading2"])]
    for item in (sections["questions"] or ["No management question was generated from deterministic panels."]):
        story.append(_pdf_paragraph(f"- {item}", styles["BodyText"]))

    overview_images = []
    for chart in plan.get("charts", [])[:6]:
        png = _figure_png_bytes(build_chart_figure(data, chart))
        if png:
            overview_images.append((chart.get("title", "Dashboard chart"), png))
    for title, png in overview_images:
        story += [PageBreak(), Paragraph(escape(_export_text(title)), styles["Heading2"]), Image(BytesIO(png), width=250 * mm, height=135 * mm)]

    # Financials
    story += [PageBreak(), Paragraph("3. Financials", styles["Heading1"])]
    fin_rows = [["Concept", "Detected column", "Current total"]]
    for item in financial["metrics"]:
        fin_rows.append([item["concept"].title(), item["column"], f"{item['total']:,.2f}"])
    if len(fin_rows) == 1:
        fin_rows.append(["No primary finance concepts detected", "", ""])
    story += [_pdf_table(fin_rows, widths=[55 * mm, 100 * mm, 70 * mm]), Spacer(1, 8)]
    detected_concepts = summary.get("financial_intelligence", {}).get("detected_concepts", {}) or {}
    story += [Paragraph("Detected business concepts", styles["Heading2"])]
    if detected_concepts:
        rows = [["Concept", "Detected column"]] + [[key, value] for key, value in detected_concepts.items()]
        story += [_pdf_table(rows, widths=[70*mm, 150*mm], font_size=7), Spacer(1, 7)]
    else:
        story.append(_pdf_paragraph("No additional business concepts were detected.", styles["BodyText"]))
    derived = financial.get("derived_measures") or {}
    story += [Paragraph("Derived financial intelligence", styles["Heading2"])]
    if derived:
        for name, formula in derived.items():
            story.append(_pdf_paragraph(f"- {name} = {formula}", styles["BodyText"]))
    else:
        story.append(_pdf_paragraph("No derived financial measures were created.", styles["BodyText"]))
    if financial.get("dimension_chart") is not None:
        png = _figure_png_bytes(financial["dimension_chart"])
        if png:
            story += [Spacer(1, 8), Image(BytesIO(png), width=245 * mm, height=125 * mm)]

    # Trends and drivers
    story += [PageBreak(), Paragraph("4. Trends and drivers", styles["Heading1"])]
    trend_summary = export_state.get("trends_ai_summary")
    story += [Paragraph("AI Trend Summary", styles["Heading2"]), _pdf_paragraph(trend_summary or "Not generated in this session.", styles["BodyText"])]
    for metric, fig in trends:
        png = _figure_png_bytes(fig)
        if png:
            story += [PageBreak(), Paragraph(f"{escape(_export_text(metric))} trend", styles["Heading2"]), Image(BytesIO(png), width=250 * mm, height=135 * mm)]

    driver_analysis = summary.get("analytical_context", {}).get("driver_analysis", {}) or {}
    driver_metrics = driver_analysis.get("metrics", {}) or {}
    if driver_metrics:
        story += [PageBreak(), Paragraph("Automatic driver analysis", styles["Heading2"])]
        for metric, dimensions in list(driver_metrics.items())[:4]:
            for dimension, details in list((dimensions or {}).items())[:2]:
                story.append(Paragraph(f"{escape(_export_text(metric))} by {escape(_export_text(dimension))}", styles["Heading3"]))
                rows = [[dimension, "Earlier", "Latest", "Change", "Contribution to net change %"]]
                for item in (details.get("drivers", []) or [])[:10]:
                    rows.append([
                        item.get("category"), f"{float(item.get('first_value', 0)):,.2f}", f"{float(item.get('last_value', 0)):,.2f}",
                        f"{float(item.get('change', 0)):+,.2f}",
                        "" if item.get("contribution_pct_of_net_change") is None else f"{float(item.get('contribution_pct_of_net_change')):+.1f}%",
                    ])
                story += [_pdf_table(rows, widths=[65 * mm, 43 * mm, 43 * mm, 43 * mm, 50 * mm], font_size=7), Spacer(1, 7)]

    # Forecasting
    story += [PageBreak(), Paragraph("5. Forecasting, scenarios and validation", styles["Heading1"])]
    if not forecast_result or not forecast_result.get("forecasts"):
        story.append(_pdf_paragraph("Forecasting board has not been built in this session.", styles["BodyText"]))
    else:
        story.append(_pdf_paragraph(
            f"Historical quarters: {', '.join(map(str, forecast_result.get('quarters', [])))}\nForecast horizon: {forecast_result.get('horizon')} quarter(s)\nScenario sensitivity: +/-{forecast_result.get('scenario_sensitivity_pct', 0):.0f}% per quarter",
            styles["BodyText"],
        ))
        forecast_commentary = export_state.get("forecast_ai_commentary")
        if forecast_commentary:
            story += [Paragraph("Executive forecast commentary", styles["Heading2"]), _pdf_paragraph(forecast_commentary, styles["BodyText"])]
        for metric in forecast_result.get("metrics", []):
            frame = forecast_result.get("forecasts", {}).get(metric)
            if not isinstance(frame, pd.DataFrame) or frame.empty:
                continue
            actual = frame[frame["Series"] == "Actual"]
            projected = frame[frame["Series"] == "Forecast"]
            backtest = (forecast_result.get("backtests", {}) or {}).get(metric, {}) or {}
            selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
            rows = [["Measure", "Selected model", "Latest actual", "End forecast", "Confidence", "WAPE"]]
            rows.append([
                metric,
                selection.get("model_name", "Flat baseline"),
                f"{float(actual['Value'].iloc[-1]):,.2f}" if not actual.empty else "",
                f"{float(projected['Value'].iloc[-1]):,.2f}" if not projected.empty else "",
                backtest.get("confidence", "Limited"),
                "" if backtest.get("wape_pct") is None else f"{float(backtest.get('wape_pct')):.1f}%",
            ])
            story += [PageBreak(), Paragraph(escape(_export_text(metric)), styles["Heading2"]), _pdf_table(rows, widths=[45*mm,55*mm,42*mm,42*mm,38*mm,35*mm], font_size=7)]
            commentary = forecast_result.get("commentary", {}).get(metric)
            if commentary:
                story += [Spacer(1, 6), _pdf_paragraph(commentary, styles["BodyText"])]
            candidate_scores = selection.get("candidate_scores", []) or []
            if candidate_scores:
                candidate_rows = [["Candidate model", "WAPE %", "MAE", "Holdouts"]]
                for candidate in candidate_scores:
                    candidate_rows.append([candidate.get("model_name"), candidate.get("wape_pct"), candidate.get("mae"), candidate.get("backtest_points")])
                story += [Spacer(1, 6), Paragraph("Candidate model comparison", styles["Heading3"]), _pdf_table(candidate_rows, font_size=7)]
            scenario_frame = (forecast_result.get("scenarios", {}) or {}).get(metric)
            if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty:
                future = scenario_frame[scenario_frame["Scenario"] != "Actual"]
                if not future.empty:
                    final_period = future["Quarter Period"].max()
                    final_rows = future[future["Quarter Period"] == final_period]
                    scenario_rows = [["End quarter", "Scenario", "Value"]] + [[str(row["Quarter Period"]), row["Scenario"], f"{float(row['Value']):,.2f}"] for _, row in final_rows.iterrows()]
                    story += [Spacer(1, 6), Paragraph("End-of-horizon scenarios", styles["Heading3"]), _pdf_table(scenario_rows, font_size=7)]
            for title, fig in _forecast_export_figures({
                **forecast_result,
                "metrics": [metric],
            }):
                png = _figure_png_bytes(fig)
                if png:
                    story += [Spacer(1, 8), Paragraph(escape(_export_text(title)), styles["Heading3"]), Image(BytesIO(png), width=235 * mm, height=118 * mm)]
            tests = backtest.get("tests")
            if isinstance(tests, pd.DataFrame) and not tests.empty:
                display_cols = [c for c in ["Quarter", "Actual", "Backtest Forecast", "Error", "Absolute % Error", "Training Quarters"] if c in tests.columns]
                table_rows = [display_cols] + tests[display_cols].head(12).astype(str).values.tolist()
                story += [Spacer(1, 7), Paragraph("Backtest details", styles["Heading3"]), _pdf_table(table_rows, font_size=6)]

    # Comparison
    story += [PageBreak(), Paragraph("6. Quarter comparison and variance bridge", styles["Heading1"])]
    if not comparison:
        story.append(_pdf_paragraph("Quarter comparison has not been built in this session.", styles["BodyText"]))
    else:
        story.append(_pdf_paragraph(f"{comparison.get('first_quarter')} vs {comparison.get('second_quarter')} - focus metric: {comparison.get('focus_metric') or 'automatic'}", styles["BodyText"]))
        metrics = comparison.get("metrics")
        if isinstance(metrics, pd.DataFrame) and not metrics.empty:
            cols = list(metrics.columns)
            rows = [cols] + metrics.head(30).astype(str).values.tolist()
            story += [_pdf_table(rows, font_size=6), Spacer(1, 6)]
        for title, fig in _comparison_export_figures(comparison):
            png = _figure_png_bytes(fig)
            if png:
                story += [PageBreak(), Paragraph(escape(_export_text(title)), styles["Heading2"]), Image(BytesIO(png), width=245 * mm, height=128 * mm)]
        breakdown = comparison.get("breakdown")
        if isinstance(breakdown, pd.DataFrame) and not breakdown.empty:
            rows = [list(breakdown.columns)] + breakdown.head(50).astype(str).values.tolist()
            story += [PageBreak(), Paragraph("Category comparison details", styles["Heading2"]), _pdf_table(rows, font_size=6)]
        bridge = comparison.get("variance_bridge") or {}
        if bridge.get("contributions"):
            bridge_rows = [["Category", "First", "Second", "Contribution"]] + [[row.get("Category"), row.get("First"), row.get("Second"), row.get("Contribution")] for row in bridge.get("contributions", [])]
            story += [Spacer(1, 7), Paragraph("Variance bridge contribution details", styles["Heading2"]), _pdf_table(bridge_rows, font_size=7), _pdf_paragraph(f"Reconciliation difference: {float(bridge.get('reconciliation_difference', 0) or 0):,.6f}", styles["BodyText"])]
        commentary = export_state.get("comparison_ai_commentary")
        story += [Spacer(1, 8), Paragraph("Comparison commentary", styles["Heading2"]), _pdf_paragraph(commentary or "Not generated/displayed in this session.", styles["BodyText"])]

    # Data quality
    story += [PageBreak(), Paragraph("7. Data quality and analysis readiness", styles["Heading1"])]
    quality = summary.get("data_quality_intelligence", {}) or {}
    story += [_pdf_paragraph(f"Readiness: {quality.get('status', 'n/a')} - {quality.get('score', 0):.0f}/100", styles["BodyText"])]
    checks = quality.get("checks", []) or []
    if checks:
        rows = [["Check", "Status", "Finding"]] + [[i.get("check"), i.get("status"), i.get("detail")] for i in checks]
        story += [_pdf_table(rows, widths=[55*mm,30*mm,160*mm], font_size=7)]
    cross_quality = export_state.get("cross_quarter_quality")
    story += [Spacer(1, 8), Paragraph("Cross-quarter consistency", styles["Heading2"])]
    if cross_quality:
        story.append(_pdf_paragraph(f"Status: {cross_quality.get('status')} - {cross_quality.get('score', 0):.0f}/100", styles["BodyText"]))
        rows = [["Check", "Status", "Finding"]] + [[i.get("check"), i.get("status"), i.get("detail")] for i in (cross_quality.get("checks", []) or [])]
        story += [_pdf_table(rows, widths=[55*mm,30*mm,160*mm], font_size=7)]
    else:
        story.append(_pdf_paragraph("Cross-quarter quality scan has not been run in this session.", styles["BodyText"]))

    profile_rows = [["Column", "Detected type", "Missing", "Unique"]]
    for column in data.columns:
        detail = profile.get("details", {}).get(column, {})
        profile_rows.append([column, detail.get("type", "unknown"), int(data[column].isna().sum()), int(data[column].nunique(dropna=True))])
    story += [PageBreak(), Paragraph("Column profile", styles["Heading2"]), _pdf_table(profile_rows, widths=[100*mm,45*mm,35*mm,35*mm], font_size=7)]

    # Decision support
    story += [PageBreak(), Paragraph("8. Decision signals and management preferences", styles["Heading1"])]
    watchlist = export_state.get("watchlist") or []
    story += [_pdf_paragraph("Management watchlist: " + (", ".join(map(str, watchlist)) if watchlist else "None selected"), styles["BodyText"]), Spacer(1, 6)]
    ranking = export_state.get("materiality") or []
    if ranking:
        rows = [["Rank", "Signal", "Type", "Priority score", "Evidence"]]
        for item in ranking[:10]:
            rows.append([item.get("rank"), ("* " if item.get("watchlisted") else "") + str(item.get("title")), item.get("type"), item.get("score"), item.get("detail")])
        story += [_pdf_table(rows, widths=[18*mm,65*mm,36*mm,30*mm,100*mm], font_size=6)]

    # Generated visual explanations
    explanations = export_state.get("visual_ai_explanations") or {}
    story += [PageBreak(), Paragraph("9. AI visual explanations", styles["Heading1"])]
    if explanations:
        for key, explanation in explanations.items():
            story += [Paragraph(escape(_export_text(key)), styles["Heading3"]), _pdf_paragraph(explanation, styles["BodyText"]), Spacer(1, 5)]
    else:
        story.append(_pdf_paragraph("No chart-level AI explanations have been generated in this session.", styles["BodyText"]))

    # Ask AI appendix - only already-visible conversation
    messages = export_state.get("chat_messages") or []
    story += [PageBreak(), Paragraph("10. Ask AI interaction appendix", styles["Heading1"])]
    if not messages:
        story.append(_pdf_paragraph("No Ask AI conversation is available in this session.", styles["BodyText"]))
    else:
        for msg in messages:
            role = str(msg.get("role", "message")).title()
            story += [Paragraph(escape(role), styles["Heading3"]), _pdf_paragraph(msg.get("content", ""), styles["BodyText"]), Spacer(1, 5)]

    # Session audit log
    audit_rows = export_state.get("audit_log") or []
    story += [PageBreak(), Paragraph("11. Session audit log", styles["Heading1"])]
    if audit_rows:
        rows = [["Time", "Action", "Detail"]]
        for item in audit_rows:
            rows.append([item.get("time", ""), item.get("action", ""), item.get("detail", "")])
        story += [_pdf_table(rows, widths=[48*mm,58*mm,145*mm], font_size=7)]
    else:
        story.append(_pdf_paragraph("No session activity has been recorded.", styles["BodyText"]))

    # Provenance
    story += [PageBreak(), Paragraph("12. Source and provenance", styles["Heading1"])]
    story += [_pdf_paragraph(
        f"Quarter: {quarter}\nSpreadsheet: {source.get('spreadsheet_name')}\nSheet: {source.get('sheet_name')}\nSource match confidence: {source.get('confidence', 'Unknown')}\nRows analysed: {summary.get('rows', 0)}\nColumns analysed: {len(summary.get('columns', []))}\nGenerated from the evidence currently available in the Workiva AI application. AI narrative is grounded against deterministic Python-calculated evidence; optional AI sections are included only when they exist in the session.",
        styles["BodyText"],
    )]

    document.build(story)
    buffer.seek(0)
    return buffer.getvalue()


# ============================================================
# POWERPOINT EXPORT - FULL DASHBOARD PACK
# ============================================================

def _ppt_add_bullets(slide, title, items, left=0.7, top=1.15, width=12.0, height=5.8, font_size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = _export_text(title)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    for item in items:
        p = tf.add_paragraph()
        p.text = _export_text(item)
        p.level = 0
        p.font.size = Pt(font_size)
    return box


def _ppt_add_table(slide, rows, left=0.55, top=1.35, width=12.2, height=5.5, font_size=10):
    if not rows:
        return None
    max_rows = min(len(rows), 18)
    rows = rows[:max_rows]
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = _export_text(value)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                if r == 0:
                    p.font.bold = True
    return table


def _ppt_add_paginated_table(presentation, title, rows, chunk_size=14, font_size=8):
    if not rows:
        return
    header = rows[0]
    body = rows[1:]
    if not body:
        body = [[""] * len(header)]
    for index in range(0, len(body), chunk_size):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        suffix = f" ({index // chunk_size + 1})" if len(body) > chunk_size else ""
        slide.shapes.title.text = _export_text(title) + suffix
        _ppt_add_table(slide, [header] + body[index:index + chunk_size], top=1.25, height=5.7, font_size=font_size)


def _ppt_chart_slide(presentation, title, figure):
    png = _figure_png_bytes(figure)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = _export_text(title)
    if png:
        slide.shapes.add_picture(BytesIO(png), Inches(0.7), Inches(1.25), width=Inches(11.95), height=Inches(5.65))
    else:
        _ppt_add_bullets(slide, "Chart export", ["Chart image unavailable. Install/verify kaleido for Plotly image export."], top=1.4)
    return slide


def build_powerpoint_export(quarter, source, summary, management_summary, data, plan, export_state=None):
    """Build a comprehensive PowerPoint containing all dashboard evidence available in-session."""
    export_state = export_state or {}
    profile = export_state.get("profile", {}) or {}
    business = export_state.get("business", {}) or {}
    financial = _financial_export_payload(data, business)
    forecast_result = export_state.get("forecast_result")
    comparison = export_state.get("comparison_result")
    sections = _report_sections(summary)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # Cover
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = _export_text(plan.get("title", "Workiva AI Dashboard"))
    slide.placeholders[1].text = (
        f"{quarter}\n{source.get('spreadsheet_name')} -> {source.get('sheet_name')}\n"
        "Comprehensive dashboard board-pack"
    )

    # Executive pulse + brief
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Executive management view"
    pulse = export_state.get("management_pulse") or {}
    items = [
        f"Top ranked signal: {pulse.get('trend', 'Not available')}",
        f"Control signal: {pulse.get('variance_or_quality', 'Not available')}",
        f"Forecast state: {pulse.get('forecast', 'Not available')}",
        f"Watchlist: {', '.join(map(str, export_state.get('watchlist') or [])) or 'None selected'}",
    ]
    _ppt_add_bullets(slide, "Management pulse", items, top=1.25, height=2.4)
    brief = export_state.get("unified_management_brief") or management_summary or "Unified/overview AI summary was not generated in this session."
    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.65), Inches(12.0), Inches(3.1))
    box.text_frame.word_wrap = True
    box.text_frame.text = _export_text(brief)
    for p in box.text_frame.paragraphs:
        p.font.size = Pt(12)

    # Overview KPI
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Overview - KPI and management focus"
    kpi_rows = [["KPI", "Value"]] + _dashboard_kpi_rows(data, plan)
    _ppt_add_table(slide, kpi_rows, top=1.3, height=2.4, font_size=11)
    _ppt_add_bullets(slide, "Variance / performance notes", sections["variance"] or ["No deterministic variance note available."], top=4.0, height=1.1, font_size=11)
    _ppt_add_bullets(slide, "Risks / watch items", sections["risks"] or ["No deterministic risk signal detected."], top=5.0, height=0.9, font_size=10)
    _ppt_add_bullets(slide, "Management questions / actions", sections["questions"] or ["No deterministic question available."], top=5.85, height=0.9, font_size=10)

    for chart in plan.get("charts", [])[:6]:
        _ppt_chart_slide(presentation, chart.get("title", "Overview chart"), build_chart_figure(data, chart))

    # Financials
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Financials"
    rows = [["Concept", "Detected column", "Current total"]] + [[i["concept"].title(), i["column"], f"{i['total']:,.2f}"] for i in financial["metrics"]]
    _ppt_add_table(slide, rows if len(rows) > 1 else [["Financials", "Status"], ["Primary concepts", "Not detected"]], top=1.25, height=3.1, font_size=10)
    derived_items = [f"{name} = {formula}" for name, formula in financial.get("derived_measures", {}).items()] or ["No derived measures created."]
    _ppt_add_bullets(slide, "Derived financial intelligence", derived_items[:8], top=4.6, height=2.2, font_size=11)
    detected_concepts = summary.get("financial_intelligence", {}).get("detected_concepts", {}) or {}
    if detected_concepts:
        concept_rows = [["Concept", "Detected column"]] + [[key, value] for key, value in detected_concepts.items()]
        _ppt_add_paginated_table(presentation, "Detected business concepts", concept_rows, chunk_size=14, font_size=9)
    if financial.get("dimension_chart") is not None:
        _ppt_chart_slide(presentation, f"{financial.get('dimension_metric')} by {financial.get('dimension')}", financial["dimension_chart"])

    # Trends
    trend_summary = export_state.get("trends_ai_summary")
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Trends - AI summary and calculated evidence"
    _ppt_add_bullets(slide, "AI Trend Summary", [trend_summary or "Not generated in this session."], top=1.25, height=5.8, font_size=12)
    for metric, fig in _trend_export_figures(data, profile):
        _ppt_chart_slide(presentation, f"{metric} trend", fig)

    drivers = summary.get("analytical_context", {}).get("driver_analysis", {}).get("metrics", {}) or {}
    for metric, dimensions in list(drivers.items())[:3]:
        for dimension, details in list((dimensions or {}).items())[:1]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = f"Driver analysis - {metric} by {dimension}"
            rows = [[dimension, "Earlier", "Latest", "Change", "Contribution %"]]
            for item in (details.get("drivers", []) or [])[:10]:
                rows.append([item.get("category"), item.get("first_value"), item.get("last_value"), item.get("change"), item.get("contribution_pct_of_net_change")])
            _ppt_add_table(slide, rows, top=1.3, height=5.4, font_size=9)

    # Forecasting
    if forecast_result and forecast_result.get("forecasts"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Forecasting - board setup"
        _ppt_add_bullets(slide, "Forecast configuration", [
            f"Historical quarters: {', '.join(map(str, forecast_result.get('quarters', [])))}",
            f"Forecast horizon: {forecast_result.get('horizon')} quarters",
            f"Scenario sensitivity: +/-{forecast_result.get('scenario_sensitivity_pct', 0):.0f}% per quarter",
            "Scenarios are deterministic sensitivities, not probabilities.",
        ], top=1.3)
        if export_state.get("forecast_ai_commentary"):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Forecasting - executive AI commentary"
            _ppt_add_bullets(slide, "Grounded commentary", [export_state.get("forecast_ai_commentary")], top=1.2, font_size=11)

        for metric in forecast_result.get("metrics", []):
            frame = forecast_result.get("forecasts", {}).get(metric)
            if not isinstance(frame, pd.DataFrame) or frame.empty:
                continue
            backtest = (forecast_result.get("backtests", {}) or {}).get(metric, {}) or {}
            selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = f"Forecast summary - {metric}"
            actual = frame[frame["Series"] == "Actual"]
            projected = frame[frame["Series"] == "Forecast"]
            rows = [["Selected model", "Latest actual", "End forecast", "Confidence", "WAPE", "MAE"]]
            rows.append([
                selection.get("model_name", "Flat baseline"),
                f"{float(actual['Value'].iloc[-1]):,.2f}" if not actual.empty else "",
                f"{float(projected['Value'].iloc[-1]):,.2f}" if not projected.empty else "",
                backtest.get("confidence", "Limited"),
                "" if backtest.get("wape_pct") is None else f"{float(backtest.get('wape_pct')):.1f}%",
                "" if backtest.get("mae") is None else f"{float(backtest.get('mae')):,.2f}",
            ])
            _ppt_add_table(slide, rows, top=1.25, height=1.5, font_size=9)
            _ppt_add_bullets(slide, "Calculated commentary", [forecast_result.get("commentary", {}).get(metric, "No metric commentary available.")], top=3.0, height=3.6, font_size=11)
            candidate_scores = selection.get("candidate_scores", []) or []
            if candidate_scores:
                candidate_rows = [["Candidate model", "WAPE %", "MAE", "Holdouts"]] + [[c.get("model_name"), c.get("wape_pct"), c.get("mae"), c.get("backtest_points")] for c in candidate_scores]
                _ppt_add_paginated_table(presentation, f"Candidate models - {metric}", candidate_rows, chunk_size=12, font_size=9)
            scenario_frame = (forecast_result.get("scenarios", {}) or {}).get(metric)
            if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty:
                future = scenario_frame[scenario_frame["Scenario"] != "Actual"]
                if not future.empty:
                    final_period = future["Quarter Period"].max()
                    final_rows = future[future["Quarter Period"] == final_period]
                    scenario_rows = [["End quarter", "Scenario", "Value"]] + [[str(row["Quarter Period"]), row["Scenario"], f"{float(row['Value']):,.2f}"] for _, row in final_rows.iterrows()]
                    _ppt_add_paginated_table(presentation, f"End scenarios - {metric}", scenario_rows, chunk_size=10, font_size=9)
            for title, fig in _forecast_export_figures({**forecast_result, "metrics": [metric]}):
                _ppt_chart_slide(presentation, title, fig)
            tests = backtest.get("tests")
            if isinstance(tests, pd.DataFrame) and not tests.empty:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                slide.shapes.title.text = f"Backtest validation - {metric}"
                cols = [c for c in ["Quarter", "Actual", "Backtest Forecast", "Error", "Absolute % Error", "Training Quarters"] if c in tests.columns]
                rows = [cols] + tests[cols].head(12).astype(str).values.tolist()
                _ppt_add_table(slide, rows, top=1.25, height=5.6, font_size=8)
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Forecasting"
        _ppt_add_bullets(slide, "Status", ["Forecasting board has not been built in this session."], top=1.3)

    # Comparison
    if comparison:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Quarter comparison - {comparison.get('first_quarter')} vs {comparison.get('second_quarter')}"
        metrics = comparison.get("metrics")
        if isinstance(metrics, pd.DataFrame) and not metrics.empty:
            rows = [list(metrics.columns)] + metrics.head(14).astype(str).values.tolist()
            _ppt_add_table(slide, rows, top=1.25, height=5.5, font_size=8)
        for title, fig in _comparison_export_figures(comparison):
            _ppt_chart_slide(presentation, title, fig)
        breakdown = comparison.get("breakdown")
        if isinstance(breakdown, pd.DataFrame) and not breakdown.empty:
            rows = [list(breakdown.columns)] + breakdown.astype(str).values.tolist()
            _ppt_add_paginated_table(presentation, "Category comparison details", rows, chunk_size=14, font_size=8)
        bridge = comparison.get("variance_bridge") or {}
        if bridge.get("contributions"):
            rows = [["Category", "First", "Second", "Contribution"]] + [[r.get("Category"), r.get("First"), r.get("Second"), r.get("Contribution")] for r in bridge.get("contributions", [])]
            _ppt_add_paginated_table(presentation, "Variance bridge contribution details", rows, chunk_size=14, font_size=8)
        commentary = export_state.get("comparison_ai_commentary")
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Quarter comparison - AI commentary"
        _ppt_add_bullets(slide, "Grounded comparison commentary", [commentary or "Not generated/displayed in this session."], top=1.25, font_size=11)
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Quarter comparison"
        _ppt_add_bullets(slide, "Status", ["Quarter comparison has not been built in this session."], top=1.3)

    # Data quality
    quality = summary.get("data_quality_intelligence", {}) or {}
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Data quality and analysis readiness"
    checks = quality.get("checks", []) or []
    rows = [["Check", "Status", "Finding"]] + [[i.get("check"), i.get("status"), i.get("detail")] for i in checks]
    _ppt_add_bullets(slide, "Readiness", [f"{quality.get('status', 'n/a')} - {quality.get('score', 0):.0f}/100"], top=1.0, height=0.8)
    _ppt_add_table(slide, rows if len(rows)>1 else [["Check", "Status", "Finding"], ["Current-quarter quality", "n/a", "No checks available"]], top=2.0, height=4.7, font_size=8)

    cross_quality = export_state.get("cross_quarter_quality")
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Cross-quarter consistency"
    if cross_quality:
        rows = [["Check", "Status", "Finding"]] + [[i.get("check"), i.get("status"), i.get("detail")] for i in (cross_quality.get("checks", []) or [])]
        _ppt_add_bullets(slide, "Status", [f"{cross_quality.get('status')} - {cross_quality.get('score', 0):.0f}/100"], top=1.0, height=0.8)
        _ppt_add_table(slide, rows, top=2.0, height=4.7, font_size=8)
    else:
        _ppt_add_bullets(slide, "Status", ["Cross-quarter quality scan has not been run in this session."], top=1.3)

    rows = [["Column", "Detected type", "Missing", "Unique"]]
    for column in data.columns:
        detail = profile.get("details", {}).get(column, {})
        rows.append([column, detail.get("type", "unknown"), int(data[column].isna().sum()), int(data[column].nunique(dropna=True))])
    _ppt_add_paginated_table(presentation, "Column profile", rows, chunk_size=14, font_size=8)

    # Decision signals
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Top decision signals"
    ranking = export_state.get("materiality") or []
    if ranking:
        rows = [["Rank", "Signal", "Type", "Score", "Evidence"]]
        for item in ranking[:10]:
            rows.append([item.get("rank"), ("* " if item.get("watchlisted") else "") + str(item.get("title")), item.get("type"), item.get("score"), item.get("detail")])
        _ppt_add_table(slide, rows, top=1.25, height=5.7, font_size=8)
    else:
        _ppt_add_bullets(slide, "Status", ["No ranked decision signals are available."], top=1.3)

    # Visual explanation appendix
    explanations = export_state.get("visual_ai_explanations") or {}
    if explanations:
        items = [f"{key}: {_export_text(value)}" for key, value in explanations.items()]
        for start in range(0, len(items), 4):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "AI visual explanations"
            _ppt_add_bullets(slide, "Generated chart explanations", items[start:start+4], top=1.2, font_size=10)

    # Ask AI appendix
    messages = export_state.get("chat_messages") or []
    if messages:
        for start in range(0, len(messages), 4):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Ask AI interaction appendix"
            items = [f"{str(m.get('role', 'message')).title()}: {_export_text(m.get('content', ''))}" for m in messages[start:start+4]]
            _ppt_add_bullets(slide, "Recent conversation", items, top=1.2, font_size=11)

    # Session audit log
    audit_rows = export_state.get("audit_log") or []
    if audit_rows:
        rows = [["Time", "Action", "Detail"]] + [
            [item.get("time", ""), item.get("action", ""), item.get("detail", "")]
            for item in audit_rows
        ]
        _ppt_add_paginated_table(presentation, "Session audit log", rows, chunk_size=14, font_size=8)
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Session audit log"
        _ppt_add_bullets(slide, "Status", ["No session activity has been recorded."], top=1.3)

    # Provenance
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Source and provenance"
    _ppt_add_bullets(slide, "Report basis", [
        f"Quarter: {quarter}",
        f"Rows analysed: {summary.get('rows', 0)}",
        f"Columns analysed: {len(summary.get('columns', []))}",
        f"Source match confidence: {source.get('confidence', 'Unknown')}",
        f"Spreadsheet: {source.get('spreadsheet_name')}",
        f"Sheet: {source.get('sheet_name')}",
        "The deck includes evidence already available in this application session; unavailable optional analyses are labeled rather than inferred.",
        "AI narrative remains grounded against deterministic Python-calculated evidence.",
    ], top=1.2, font_size=12)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output.getvalue()


@st.cache_data(ttl=1800, show_spinner=False, max_entries=16)
def dataframe_csv_bytes(data):
    """Serialize a dataframe once per cached dataset instead of on every rerun."""
    return data.to_csv(index=False).encode("utf-8")


# ============================================================
# LANDING PAGE
# ============================================================

if not st.session_state[
    "app_open"
]:

    st.title(
        APP_NAME
    )


    st.subheader(
        (
            "Management reporting, "
            "analytics and AI commentary"
        )
    )


    st.write(
        (
            "Connect to Workiva, discover quarter data, "
            "build adaptive dashboards, compare reporting "
            "periods, and ask management questions "
            "in natural language."
        )
    )


    try:

        discovered = (
            cached_quarters()
        )

    except Exception:

        discovered = []


    latest = (

        discovered[0]

        if discovered

        else "Not yet discovered"
    )


    c1, c2, c3 = (
        st.columns(
            3
        )
    )


    c1.metric(
        "Latest discovered quarter",
        latest,
    )


    c2.metric(
        "Workiva access",
        "Read-only",
    )


    c3.metric(
        "AI mode",
        "Dashboard + analysis",
    )


    st.info(
        (
            "Workiva write-back is disabled. "
            "The app only reads data until explicit "
            "write permissions and confirmation "
            "controls are designed later."
        )
    )


    if st.button(
        "Open dashboard",
        type="primary",
        use_container_width=True,
    ):

        st.session_state[
            "app_open"
        ] = True


        log_event(
            "Open dashboard",
            "User entered the dashboard.",
        )


        st.rerun()


    st.stop()


# ============================================================
# MAIN PAGE
# ============================================================

st.title(
    APP_NAME
)


st.caption(
    (
        "Smart Workiva discovery, adaptive dashboards, "
        "quarter comparison and AI management analysis."
    )
)

# Subtle application-wide accessibility and responsive refinements.
# These preserve the existing visual language instead of introducing a
# separate, visually heavier accessibility skin.
st.markdown(
    """
    <style>
    /* Respect the operating system's reduced-motion preference in every mode. */
    @media (prefers-reduced-motion: reduce) {
        *, *::before, *::after {
            scroll-behavior: auto !important;
            animation-duration: 0.001ms !important;
            animation-iteration-count: 1 !important;
            transition-duration: 0.001ms !important;
        }
    }

    /* Clear keyboard focus without changing the sophisticated neutral design. */
    button:focus-visible,
    [role="button"]:focus-visible,
    [role="tab"]:focus-visible,
    input:focus-visible,
    textarea:focus-visible,
    a:focus-visible {
        outline: 2px solid currentColor !important;
        outline-offset: 3px !important;
    }

    /* Tabs remain usable on narrower screens rather than compressing labels. */
    [data-testid="stTabs"] [role="tablist"] {
        overflow-x: auto;
        scrollbar-width: thin;
        gap: 1.15rem !important;
        align-items: center;
    }
    [data-testid="stTabs"] [role="tab"] {
        white-space: nowrap;
        min-height: 42px;
        flex: 0 0 auto !important;
        padding-left: 0.25rem !important;
        padding-right: 0.25rem !important;
    }

    /* Preserve the spacing on compact screens; horizontal scrolling is
       preferable to squeezing navigation labels together. */
    @media (max-width: 900px) {
        [data-testid="stTabs"] [role="tablist"] {
            gap: 0.9rem !important;
        }
    }


    /* Horizontal workspace navigator: spaced, keyboard friendly, and scrollable. */
    [data-testid="stRadio"] div[role="radiogroup"] {
        display: flex !important;
        flex-wrap: nowrap !important;
        gap: 1.15rem !important;
        overflow-x: auto !important;
        padding: 0.2rem 0 0.45rem 0 !important;
        scrollbar-width: thin;
    }
    [data-testid="stRadio"] div[role="radiogroup"] > label {
        flex: 0 0 auto !important;
        min-height: 42px !important;
        padding: 0.35rem 0.2rem !important;
        white-space: nowrap !important;
    }
    [data-testid="stRadio"] div[role="radiogroup"] > label:has(input:checked) {
        font-weight: 750 !important;
        text-decoration: underline !important;
        text-decoration-thickness: 2px !important;
        text-underline-offset: 0.38rem !important;
    }

    /* Keep secondary captions compact in Focus mode while retaining readability. */
    .compact-meta {
        line-height: 1.35;
        margin: 0.1rem 0;
    }

    /* Comfortable touch targets while keeping controls visually compact. */
    button, [role="button"], [data-baseweb="select"] > div {
        min-height: 40px;
    }

    /* Prevent overly wide prose while letting analytical visuals stay expansive. */
    [data-testid="stAlert"] p,
    [data-testid="stCaptionContainer"] p {
        line-height: 1.5;
    }
    </style>
    """,
    unsafe_allow_html=True,
)


# ============================================================
# SIDEBAR
# ============================================================

with st.sidebar:

    st.header(
        "Agent controls"
    )


    try:

        quarter_options = (
            cached_quarters()
        )

    except Exception:

        quarter_options = []


    if not quarter_options:

        quarter_options = [

            "Q1 2026",
            "Q4 2025",
            "Q3 2025",
            "Q2 2025",
            "Q1 2025",
        ]


    selected_quarter = (
        st.selectbox(
            "Quarter",
            quarter_options,
        )
    )


    board_mode = "AI designed"


    st.divider()

    st.subheader(
        "Display & accessibility"
    )

    display_mode = st.selectbox(
        "Display mode",
        [
            "Standard",
            "Focus",
            "Accessible reading",
        ],
        index=1,
        help=(
            "Focus is the recommended default: it shows the most important information first. "
            "Standard keeps the full dashboard. Accessible reading uses larger text, clearer controls, "
            "single-column supporting visuals, and text alternatives for charts."
        ),
        key="display_mode",
    )

    show_advanced_tools = st.checkbox(
        "Show advanced management tools",
        value=False,
        help=(
            "Shows the watchlist, full decision-signal table, unified management brief, "
            "and AI dashboard builder. Leave off for the fastest, least crowded experience."
        ),
        key="show_advanced_management_tools",
    )


focus_mode = (
    display_mode == "Focus"
)

accessible_mode = (
    display_mode == "Accessible reading"
)


if accessible_mode:

    st.markdown(
        """
        <style>
        /* Accessible reading: larger type, generous spacing and strong focus. */
        html {
            font-size: 18px !important;
            scroll-behavior: auto !important;
        }

        body, [data-testid="stAppViewContainer"], [data-testid="stSidebar"] {
            font-size: 1rem !important;
        }

        p, li, label, [data-testid="stCaptionContainer"] {
            line-height: 1.7 !important;
            letter-spacing: 0.01em !important;
        }

        h1 { font-size: 2.25rem !important; line-height: 1.2 !important; }
        h2 { font-size: 1.75rem !important; line-height: 1.25 !important; }
        h3 { font-size: 1.45rem !important; line-height: 1.3 !important; }

        /* Larger pointer/touch targets. */
        button,
        [role="button"],
        [role="tab"],
        [data-baseweb="select"] > div,
        input,
        textarea {
            min-height: 48px !important;
            font-size: 1rem !important;
        }

        textarea {
            line-height: 1.6 !important;
        }

        /* Make keyboard focus impossible to miss. */
        button:focus-visible,
        [role="button"]:focus-visible,
        [role="tab"]:focus-visible,
        input:focus-visible,
        textarea:focus-visible,
        [data-baseweb="select"] *:focus-visible,
        a:focus-visible {
            outline: 3px solid currentColor !important;
            outline-offset: 3px !important;
            box-shadow: 0 0 0 2px Canvas !important;
        }

        /* Tabs need a clear selected state that does not rely on colour alone. */
        [data-testid="stTabs"] [role="tablist"] {
            gap: 1rem !important;
        }

        [role="tab"] {
            padding: 0.75rem 0.35rem !important;
            font-weight: 600 !important;
            flex: 0 0 auto !important;
        }

        [role="tab"][aria-selected="true"] {
            font-weight: 800 !important;
            text-decoration: underline !important;
            text-decoration-thickness: 3px !important;
            text-underline-offset: 0.35rem !important;
        }

        /* Metrics, alerts and expanders become easier to scan. */
        [data-testid="stMetric"] {
            padding: 1rem !important;
            border: 1px solid currentColor !important;
            border-radius: 0.5rem !important;
        }

        [data-testid="stMetricLabel"] {
            font-size: 1rem !important;
            font-weight: 700 !important;
        }

        [data-testid="stMetricValue"] {
            font-size: 2.15rem !important;
            line-height: 1.2 !important;
        }

        [data-testid="stAlert"] {
            border-width: 2px !important;
        }

        [data-testid="stExpander"] summary {
            min-height: 48px !important;
            font-size: 1.05rem !important;
            font-weight: 700 !important;
        }

        /* Dataframes/tables: larger cells and visible row separation. */
        [data-testid="stDataFrame"],
        [data-testid="stTable"] {
            font-size: 1rem !important;
        }

        [data-testid="stDataFrame"] * {
            line-height: 1.45 !important;
        }

        /* Avoid motion when accessible reading is selected. */
        *, *::before, *::after {
            animation-duration: 0.001ms !important;
            animation-iteration-count: 1 !important;
            transition-duration: 0.001ms !important;
        }

        /* Keep content comfortably readable on very wide screens. */
        [data-testid="stMainBlockContainer"] {
            max-width: 1200px !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.info(
        "Accessible reading is on: larger text and controls, stronger keyboard "
        "focus, reduced motion, clearer tabs, single-column supporting charts, "
        "and enlarged chart labels are enabled."
    )


# ============================================================
# LOAD QUARTER
# ============================================================

try:

    with st.spinner(
        (
            f"Finding {selected_quarter} "
            "in Workiva..."
        )
    ):

        bundle = (
            quarter_bundle(
                selected_quarter
            )
        )


except Exception as error:

    st.error(
        friendly_error(
            error
        )
    )


    st.stop()


source = (
    bundle[
        "source"
    ]
)

data = (
    bundle[
        "data"
    ]
)

profile = (
    bundle[
        "profile"
    ]
)

business = (
    bundle[
        "business"
    ]
)

summary = (
    bundle[
        "summary"
    ]
)


source_key = (
    f"{source['spreadsheet_id']}"
    f":"
    f"{source['sheet_id']}"
)


if (
    st.session_state[
        "current_source_key"
    ]
    != source_key
):

    st.session_state[
        "current_source_key"
    ] = source_key


    st.session_state[
        "agent_plan"
    ] = None


    st.session_state[
        "chat_messages"
    ] = []


    st.session_state[
        "comparison_result"
    ] = None


    st.session_state[
        "pdf_bytes"
    ] = None


    st.session_state[
        "pptx_bytes"
    ] = None

    st.session_state[
        "management_summary"
    ] = None

    st.session_state[
        "management_summary_key"
    ] = None


    # Preserve bookmarked metrics across quarters when possible. If none survive,
    # seed a concise finance-oriented watchlist from the current source.
    existing_watchlist = [
        item for item in st.session_state.get("management_watchlist", [])
        if item in profile.get("numeric", [])
    ]
    if not existing_watchlist:
        for candidate in [
            business.get("revenue"),
            business.get("profit"),
            business.get("ebitda"),
            business.get("costs"),
            business.get("headcount"),
        ]:
            if candidate and candidate in profile.get("numeric", []) and candidate not in existing_watchlist:
                existing_watchlist.append(candidate)
            if len(existing_watchlist) >= 4:
                break
    st.session_state["management_watchlist"] = existing_watchlist[:5]
    st.session_state.pop("management_watchlist_selector", None)


    st.session_state[
        "forecast_ai_commentary"
    ] = None


    st.session_state[
        "forecast_ai_commentary_key"
    ] = None


    st.session_state[
        "last_refreshed"
    ] = (
        now_text()
    )


    log_event(
        "Load source",
        (
            f"{selected_quarter}: "
            f"{source['spreadsheet_name']} "
            f"→ "
            f"{source['sheet_name']}"
        ),
    )


# ============================================================
# SOURCE INFORMATION
# ============================================================

confidence = source.get("confidence", "Unknown")
current_quality = summary.get("data_quality_intelligence", {}) or {}
quality_status = current_quality.get("status", "Unavailable")
quality_score = current_quality.get("score")
quality_guidance = current_quality.get("analysis_guidance", "")
quality_label = f"Analysis confidence: {quality_status}"
if quality_score is not None:
    quality_label += f" · {quality_score:.0f}/100"

# One compact status surface replaces several vertically stacked alerts.
with st.container(border=True):
    if accessible_mode:
        st.markdown(f"**Source: {source['spreadsheet_name']} → {source['sheet_name']}**")
        st.caption(
            f"{selected_quarter} · {profile['rows']:,} rows · {profile['columns']:,} columns · "
            f"{len(profile['numeric'])} numeric · {len(profile['categories'])} categorical · "
            f"{len(profile['dates'])} date"
        )
        if quality_status in {"Strong", "Good"}:
            st.success(quality_label)
        elif quality_status == "Watch":
            st.warning(quality_label)
        else:
            st.error(quality_label)
        if quality_guidance:
            st.caption(quality_guidance)
    else:
        source_col, quality_col = st.columns([1.35, 1])
        with source_col:
            st.markdown(f"**{source['spreadsheet_name']} → {source['sheet_name']}**")
            st.caption(
                f"{selected_quarter} · {profile['rows']:,} rows · {profile['columns']:,} columns · "
                f"{len(profile['numeric'])} numeric · {len(profile['categories'])} categorical · "
                f"{len(profile['dates'])} date"
            )
        with quality_col:
            if quality_status in {"Strong", "Good"}:
                st.success(quality_label)
            elif quality_status == "Watch":
                st.warning(quality_label)
            else:
                st.error(quality_label)
            if quality_guidance and not focus_mode:
                st.caption(quality_guidance)

    if not focus_mode:
        with st.expander("Source & refresh details", expanded=False):
            st.caption(f"Last refreshed: {st.session_state['last_refreshed']}")
            if confidence == "Exact match":
                st.caption("Match confidence: Exact Workiva match")
            else:
                st.warning(f"Match confidence: {confidence}. Fallback matching was used.")
            if quality_guidance:
                st.caption(quality_guidance)

# ============================================================
# MANAGEMENT WATCHLIST
# ============================================================

available_watch_metrics = list(profile.get("numeric", []))
current_watchlist = [
    item for item in st.session_state.get("management_watchlist", [])
    if item in available_watch_metrics
]
st.session_state["management_watchlist"] = current_watchlist

if show_advanced_tools:
    with st.expander("Management watchlist", expanded=False):
        st.caption(
            "Bookmark up to five measures that deserve recurring management attention. "
            "Bookmarks influence prioritization and defaults, but never replace calculated materiality."
        )
        selected_watchlist = st.multiselect(
            "Priority measures",
            available_watch_metrics,
            default=current_watchlist,
            max_selections=5,
            key="management_watchlist_selector",
            help="These metrics receive a modest attention boost across ranking, forecasting defaults, summaries and Ask AI.",
        )
        st.session_state["management_watchlist"] = selected_watchlist
        if selected_watchlist:
            st.caption("Watching: " + " · ".join(map(str, selected_watchlist)))
        else:
            st.caption("No metrics are bookmarked. The app will rely entirely on calculated signal ranking.")

# Use a local summary copy so cached AI outputs are keyed to watchlist state rather
# than silently reusing commentary generated under a previous preference set.
summary = deepcopy(summary)
summary["management_watchlist"] = list(st.session_state.get("management_watchlist", []))

# ============================================================
# MATERIALITY + MANAGEMENT PULSE
# Compute the ranking once per rerun and reuse it in both surfaces.
# ============================================================
current_materiality = build_materiality_ranking(
    summary,
    forecast_result=st.session_state.get("forecast_result"),
    comparison=st.session_state.get("comparison_result"),
    limit=8,
)

management_pulse = build_management_pulse(
    summary,
    st.session_state.get("forecast_result"),
    materiality_ranking=current_materiality,
)
with st.container(border=True):
    if focus_mode:
        st.caption("Management pulse · the three highest-value signals are shown first; deeper evidence is available in the workspaces.")
    if accessible_mode:
        st.markdown("**Top ranked signal**")
        st.write(management_pulse["trend"])
        st.markdown("**Control signal**")
        st.write(management_pulse["variance_or_quality"])
        st.markdown("**Forecast state**")
        st.write(management_pulse["forecast"])
    else:
        pulse_a, pulse_b, pulse_c = st.columns(3)
        with pulse_a:
            st.caption("Top ranked signal")
            st.markdown(management_pulse["trend"])
        with pulse_b:
            st.caption("Control signal")
            st.markdown(management_pulse["variance_or_quality"])
        with pulse_c:
            st.caption("Forecast state")
            st.markdown(management_pulse["forecast"])



# ============================================================
# MATERIALITY-BASED SIGNAL RANKING
# ============================================================

if show_advanced_tools:

    with st.expander("Top decision signals", expanded=False):
        st.caption(
            "Deterministic priority ranking: 50% movement magnitude, 30% business relevance, "
            "and 20% consistency/reliability. This is a management-priority heuristic, not an "
            "accounting materiality threshold."
        )
        if not current_materiality:
            st.info("No material signal could be ranked from the currently available evidence.")
        else:
            priority_rows = []
            for item in current_materiality[:5]:
                priority_rows.append({
                    "Rank": item.get("rank"),
                    "Signal": ("★ " if item.get("watchlisted") else "") + str(item.get("title")),
                    "Type": item.get("type"),
                    "Priority score": item.get("score"),
                    "Evidence": item.get("detail"),
                })
            st.dataframe(
                pd.DataFrame(priority_rows),
                use_container_width=True,
                hide_index=True,
            )
            st.caption(
                "Score components: Magnitude 50% · Business relevance 30% · "
                "Consistency/reliability 20%. Forecast magnitude is reliability-adjusted so a "
                "large low-confidence projection cannot dominate the ranking on size alone."
            )



# ============================================================
# UNIFIED MANAGEMENT BRIEF
# ============================================================

if show_advanced_tools:

    # This is deliberately one compact global surface rather than another dashboard tab.
    # It uses only evidence already calculated in the application and runs AI on demand.
    brief_forecast = st.session_state.get("forecast_result")
    brief_comparison = st.session_state.get("comparison_result")
    brief_state = {
        "source_key": source_key,
        "forecast": {
            "quarters": (brief_forecast or {}).get("quarters", []),
            "metrics": (brief_forecast or {}).get("metrics", []),
            "horizon": (brief_forecast or {}).get("horizon"),
            "scenario_sensitivity_pct": (brief_forecast or {}).get("scenario_sensitivity_pct"),
        },
        "comparison": {
            "first_quarter": (brief_comparison or {}).get("first_quarter"),
            "second_quarter": (brief_comparison or {}).get("second_quarter"),
            "focus_metric": (brief_comparison or {}).get("focus_metric"),
        },
        "management_watchlist": summary.get("management_watchlist", []),
    }
    brief_cache_key = json.dumps(brief_state, sort_keys=True, default=str)

    if st.session_state.get("unified_management_brief_key") != brief_cache_key:
        st.session_state["unified_management_brief"] = None
        st.session_state["unified_management_brief_key"] = brief_cache_key

    with st.expander("Executive AI Management Brief", expanded=False):
        st.caption(
            "One grounded brief across performance, trends, drivers, forecast scenarios, validation, "
            "comparison and data quality. Generated only when requested."
        )

        current_brief = st.session_state.get("unified_management_brief")
        if current_brief:
            st.markdown(current_brief)
            render_ai_evidence_indicator(
                summary,
                purpose="unified management brief",
                forecast_result=brief_forecast,
                comparison_result=brief_comparison,
                key="unified_brief_evidence",
            )
            if st.button("Refresh management brief", key="refresh_unified_management_brief"):
                st.session_state["unified_management_brief"] = None
                st.rerun()
        elif st.button("Generate unified management brief", key="generate_unified_management_brief"):
            try:
                with st.spinner("Synthesizing management evidence..."):
                    current_brief = generate_unified_management_brief(
                        summary,
                        {"title": "Unified management view", "reason": "Cross-board evidence synthesis"},
                        brief_forecast,
                        brief_comparison,
                    )
                st.session_state["unified_management_brief"] = current_brief
                log_event("Generate unified management brief", selected_quarter)
                st.markdown(current_brief)
                render_ai_evidence_indicator(
                    summary,
                    purpose="unified management brief",
                    forecast_result=brief_forecast,
                    comparison_result=brief_comparison,
                    key="unified_brief_evidence_new",
                )
            except Exception as error:
                st.warning("AI management brief unavailable: " + friendly_error(error))


# ============================================================
# AI DASHBOARD BUILDER
# ============================================================

if show_advanced_tools:

    st.subheader(
        "AI Dashboard Builder"
    )


    st.caption(
        "Choose a suggested dashboard or describe exactly what you want."
    )


    SUGGESTED_DASHBOARD_PROMPTS = {

        "Executive overview": (
            "Build an executive dashboard highlighting the most important "
            "KPIs, trends, risks, and business drivers in this dataset."
        ),

        "Financial performance": (
            "Build a financial performance dashboard focused on revenue, "
            "costs, budget variance, and the most important areas requiring "
            "management attention."
        ),

        "Trends & anomalies": (
            "Build a dashboard focused on trends, unusual movements, "
            "outliers, and areas of improving or declining performance."
        ),

        "Management deep dive": (
            "Build a management dashboard showing the strongest and weakest "
            "categories, key contributors to performance, and actionable insights."
        ),
    }


    prompt_columns = st.columns(4)


    for prompt_index, (prompt_label, prompt_text) in enumerate(
        SUGGESTED_DASHBOARD_PROMPTS.items()
    ):

        if prompt_columns[prompt_index].button(
            prompt_label,
            use_container_width=True,
            key=f"suggested_dashboard_prompt_{prompt_index}",
        ):

            st.session_state[
                "agent_request_input"
            ] = prompt_text

            st.session_state[
                "agent_request"
            ] = prompt_text


    st.caption(
        "AI grounding: Workiva/Python calculations remain the source of truth. "
        "AI selects, explains and prioritizes evidence; unsupported numeric claims are rejected."
    )

    agent_request = st.text_area(
        "Tell the agent what you want",
        key="agent_request_input",
        placeholder=(
            "Example: Build a CFO dashboard focused on revenue versus budget, "
            "largest negative variances, regional performance and key risks."
        ),
        height=100,
    )


    if st.button(
        "Build AI dashboard",
        type="primary",
    ):

        if not agent_request.strip():

            st.warning(
                "Type a dashboard request first."
            )


        else:

            try:

                with st.spinner(
                    "Designing dashboard..."
                ):

                    raw_plan = create_dashboard_plan(
                        agent_request,
                        summary,
                    )

                    plan = validate_plan(
                        raw_plan,
                        profile,
                    )

                    st.session_state["agent_plan"] = plan


                st.session_state[
                    "agent_request"
                ] = agent_request


                st.session_state[
                    "pdf_bytes"
                ] = None


                st.session_state[
                    "pptx_bytes"
                ] = None


                log_event(
                    "Build AI dashboard",
                    agent_request,
                )


                st.rerun()


            except Exception as error:

                st.error(
                    friendly_error(
                        error
                    )
                )


# ============================================================
# ACTIVE PLAN
# ============================================================

if board_mode == "AI designed":

    # FAST START: render a useful deterministic executive dashboard immediately.
    # Copilot is only called when the user explicitly clicks "Build AI dashboard".
    # This removes a full network/AI round trip from initial page load.
    if st.session_state["agent_plan"] is None:
        st.session_state["agent_plan"] = manual_plan(
            "Executive",
            profile,
            business,
        )

    active_plan = st.session_state["agent_plan"]


elif board_mode == "Data quality":

    active_plan = None


else:

    active_plan = (
        manual_plan(
            board_mode,
            profile,
            business,
        )
    )


# The visible plan may be locally customized in the Overview tab.
# It starts as the AI/manual plan and remains available to exports and chat.
display_plan = active_plan


summary_plan = (

    {
        "title":
            "Data Quality Dashboard",

        "reason":
            (
                "Reviewing dataset "
                "completeness."
            ),
    }

    if board_mode
    == "Data quality"

    else active_plan
)


# ============================================================
# WORKSPACE NAVIGATION
# ============================================================

workspace_options = [
    "Overview",
    "Financials",
    "Trends",
    "Forecasting",
    "Comparison",
    "Data quality",
    "Ask AI",
]

active_workspace = st.radio(
    "Workspace",
    workspace_options,
    horizontal=True,
    label_visibility="collapsed",
    key="active_workspace",
    help="Only the selected workspace is executed, improving responsiveness on large datasets.",
)

# ============================================================
# OVERVIEW TAB
# ============================================================

if active_workspace == "Overview":

    if board_mode == "Data quality":

        st.info(
            (
                "Data quality is selected. "
                "Open the Data quality tab."
            )
        )


    else:

        st.subheader(
            active_plan[
                "title"
            ]
        )


        if active_plan[
            "reason"
        ]:

            st.caption(
                active_plan[
                    "reason"
                ]
            )


        # ----------------------------------------------------
        # LOCAL DASHBOARD CUSTOMIZER
        # ----------------------------------------------------
        display_plan = customize_dashboard_plan(
            active_plan,
            data,
            profile,
            key_prefix=f"{selected_quarter}_{board_mode}",
        )

        # ----------------------------------------------------
        # KPI ROW
        # ----------------------------------------------------
        # Standard and Accessible reading keep all planned KPIs.
        # Focus mode intentionally limits the first view to four KPIs.
        visible_kpis = (
            display_plan["kpis"][:4]
            if focus_mode
            else display_plan["kpis"]
        )

        render_kpis(
            data,
            visible_kpis,
        )


        charts = [
            chart
            for chart in display_plan.get("charts", [])
            if chart.get("visible", True)
        ]


        # ----------------------------------------------------
        # PRIMARY CHART
        # The AI's first chart gets visual priority.
        # ----------------------------------------------------
        if charts:

            st.markdown(
                "#### Primary view"
            )

            render_chart(
                data,
                charts[0],
                key=f"dashboard_primary_{selected_quarter}",
            )


        # ----------------------------------------------------
        # SUPPORTING BREAKDOWNS
        # Standard = existing two-column layout.
        # Focus = secondary charts collapsed.
        # Accessible reading = single-column charts.
        # ----------------------------------------------------
        supporting_charts = charts[1:]

        if supporting_charts:

            st.markdown(
                "#### Supporting breakdowns"
            )

            if focus_mode:

                with st.expander(
                    "Show supporting charts",
                    expanded=False,
                ):

                    for chart_index, chart in enumerate(
                        supporting_charts
                    ):

                        render_chart(
                            data,
                            chart,
                            key=(
                                f"dashboard_supporting_"
                                f"{selected_quarter}_"
                                f"{chart_index}"
                            ),
                        )

            elif accessible_mode:

                for chart_index, chart in enumerate(
                    supporting_charts
                ):

                    render_chart(
                        data,
                        chart,
                        key=(
                            f"dashboard_supporting_"
                            f"{selected_quarter}_"
                            f"{chart_index}"
                        ),
                    )

            else:

                for index in range(
                    0,
                    len(supporting_charts),
                    2,
                ):

                    row = st.columns(2)

                    for offset in range(2):

                        chart_index = index + offset

                        if chart_index >= len(supporting_charts):
                            break

                        with row[offset]:
                            render_chart(
                                data,
                                supporting_charts[chart_index],
                                key=(
                                    f"dashboard_supporting_"
                                    f"{selected_quarter}_"
                                    f"{chart_index}"
                                ),
                            )


        # ----------------------------------------------------
        # MANAGEMENT PANELS
        # These use already-calculated Python facts only, so
        # they add no extra Copilot or Workiva calls.
        # ----------------------------------------------------
        layout_panels = build_dashboard_panels(
            summary
        )

        st.markdown(
            "#### Management focus"
        )

        if focus_mode:

            with st.expander(
                "Show management details",
                expanded=False,
            ):

                st.markdown("**Variance panel**")
                for item in layout_panels["variance"]:
                    st.markdown(f"- {item}")

                st.markdown("**Risk / watch panel**")
                for item in layout_panels["risks"]:
                    st.markdown(f"- {item}")

                st.markdown("**Management questions**")
                for item in layout_panels["questions"]:
                    st.markdown(f"- {item}")

        elif accessible_mode:

            st.markdown("**Variance panel**")
            for item in layout_panels["variance"]:
                st.markdown(f"- {item}")

            st.markdown("**Risk / watch panel**")
            for item in layout_panels["risks"]:
                st.markdown(f"- {item}")

            with st.expander(
                "Management questions",
                expanded=False,
            ):
                for item in layout_panels["questions"]:
                    st.markdown(f"- {item}")

        else:

            variance_col, risk_col = st.columns(2)

            with variance_col:
                st.markdown("**Variance panel**")
                for item in layout_panels["variance"]:
                    st.markdown(f"- {item}")

            with risk_col:
                st.markdown("**Risk / watch panel**")
                for item in layout_panels["risks"]:
                    st.markdown(f"- {item}")

            with st.expander(
                "Management questions",
                expanded=False,
            ):
                for item in layout_panels["questions"]:
                    st.markdown(f"- {item}")


    st.divider()


    st.subheader(
        "AI Management Summary"
    )

    # FAST START: do not block every dashboard load on a second Copilot call.
    # Generate the narrative only when requested, then keep it for this source/plan.
    summary_cache_key = (
        f"{source_key}:"
        f"{abs(hash(json.dumps({
            'plan': summary_plan,
            'management_watchlist': summary.get('management_watchlist', []),
        }, sort_keys=True, default=str)))}"
    )

    if st.session_state.get("management_summary_key") != summary_cache_key:
        st.session_state["management_summary"] = None
        st.session_state["management_summary_key"] = summary_cache_key

    if st.session_state.get("management_summary"):
        management_summary = st.session_state["management_summary"]
        st.markdown(management_summary)
        render_ai_evidence_indicator(summary, purpose="management summary", key="management_summary_evidence")
    else:
        st.caption(
            "The dashboard is ready now. Generate the narrative only when you need it."
        )
        if st.button(
            "Generate AI management summary",
            key="generate_management_summary_button",
        ):
            try:
                with st.spinner("Generating management summary..."):
                    management_summary = generate_management_summary(
                        summary,
                        summary_plan,
                    )
                st.session_state["management_summary"] = management_summary
                st.markdown(management_summary)
                render_ai_evidence_indicator(summary, purpose="management summary", key="management_summary_evidence_new")
            except Exception as error:
                management_summary = (
                    "AI summary unavailable: " + friendly_error(error)
                )
                st.warning(management_summary)
        else:
            management_summary = ""


    # ========================================================
    # EXPORT
    # ========================================================

    st.divider()


    st.subheader(
        "Export"
    )


    csv_data = dataframe_csv_bytes(data)


    if st.download_button(
        "Download current quarter as CSV",
        data=csv_data,
        file_name=(
            selected_quarter
            .replace(
                " ",
                "_",
            )
            + "_workiva.csv"
        ),
        mime="text/csv",
    ):

        log_event(
            "Export CSV",
            selected_quarter,
        )


    full_export_state = _full_export_state(
        summary,
        profile,
        business,
        management_pulse,
        current_materiality,
    )

    export_fingerprint = json.dumps({
        "source": source_key,
        "management_summary": bool(management_summary),
        "unified_brief": bool(full_export_state.get("unified_management_brief")),
        "trends_summary": bool(full_export_state.get("trends_ai_summary")),
        "forecast": {
            "quarters": (full_export_state.get("forecast_result") or {}).get("quarters", []),
            "metrics": (full_export_state.get("forecast_result") or {}).get("metrics", []),
            "horizon": (full_export_state.get("forecast_result") or {}).get("horizon"),
            "scenario_sensitivity_pct": (full_export_state.get("forecast_result") or {}).get("scenario_sensitivity_pct"),
        },
        "comparison": {
            "first": (full_export_state.get("comparison_result") or {}).get("first_quarter"),
            "second": (full_export_state.get("comparison_result") or {}).get("second_quarter"),
            "focus": (full_export_state.get("comparison_result") or {}).get("focus_metric"),
        },
        "cross_quality_key": st.session_state.get("cross_quarter_quality_key"),
        "watchlist": full_export_state.get("watchlist", []),
        "materiality": [(item.get("rank"), item.get("title"), item.get("score")) for item in full_export_state.get("materiality", [])],
        "visual_explanations": sorted((full_export_state.get("visual_ai_explanations") or {}).keys()),
        "chat_count": len(full_export_state.get("chat_messages") or []),
        "audit_count": len(full_export_state.get("audit_log") or []),
        "audit_last": (full_export_state.get("audit_log") or [{}])[-1].get("time") if full_export_state.get("audit_log") else None,
    }, sort_keys=True, default=str)
    if st.session_state.get("full_export_fingerprint") != export_fingerprint:
        st.session_state["pdf_bytes"] = None
        st.session_state["pptx_bytes"] = None
        st.session_state["full_export_fingerprint"] = export_fingerprint

    st.caption(
        "PDF and PowerPoint are comprehensive board packs: they include all dashboard evidence "
        "already available in this session. Optional analyses that have not been built are labeled as unavailable."
    )
    with st.expander("Export coverage", expanded=False):
        st.markdown(
            "**Always included:** Overview, Financials, Trends/driver evidence, current-quarter Data Quality, "
            "decision signals, watchlist, source provenance and column profile."
        )
        st.markdown(
            f"**Forecasting:** {'included' if full_export_state.get('forecast_result') else 'not built'}  \n"
            f"**Quarter comparison / variance bridge:** {'included' if full_export_state.get('comparison_result') else 'not built'}  \n"
            f"**Cross-quarter quality scan:** {'included' if full_export_state.get('cross_quarter_quality') else 'not run'}  \n"
            f"**Unified AI brief:** {'included' if full_export_state.get('unified_management_brief') else 'not generated'}  \n"
            f"**Chart-level AI explanations:** {len(full_export_state.get('visual_ai_explanations') or {})} available  \n"
            f"**Recent Ask AI messages:** {len(full_export_state.get('chat_messages') or [])} included  \n"
            f"**Session audit log:** {len(full_export_state.get('audit_log') or [])} entries included"
        )

    # --------------------------------------------------------
    # LAZY EXPORTS
    #
    # The expensive files are only generated when requested.
    # This is intentionally faster than preparing them on
    # every Streamlit rerun.
    # --------------------------------------------------------

    prepare1, prepare2 = (
        st.columns(
            2
        )
    )


    with prepare1:

        if st.button(
            "Prepare full PDF board pack",
            use_container_width=True,
        ):

            try:

                with st.spinner(
                    "Preparing PDF..."
                ):

                    st.session_state[
                        "pdf_bytes"
                    ] = (
                        build_pdf_export(
                            selected_quarter,
                            source,
                            summary,
                            management_summary,
                            data,
                            display_plan,
                            full_export_state,
                        )
                    )


                log_event(
                    "Prepare PDF",
                    selected_quarter,
                )


            except Exception as error:

                st.error(
                    friendly_error(
                        error
                    )
                )


    with prepare2:

        if st.button(
            "Prepare full PowerPoint board pack",
            use_container_width=True,
        ):

            try:

                with st.spinner(
                    "Preparing PowerPoint..."
                ):

                    st.session_state[
                        "pptx_bytes"
                    ] = (
                        build_powerpoint_export(
                            selected_quarter,
                            source,
                            summary,
                            management_summary,
                            data,
                            display_plan,
                            full_export_state,
                        )
                    )


                log_event(
                    "Prepare PowerPoint",
                    selected_quarter,
                )


            except Exception as error:

                st.error(
                    friendly_error(
                        error
                    )
                )


    download1, download2 = (
        st.columns(
            2
        )
    )


    with download1:

        if (
            st.session_state[
                "pdf_bytes"
            ]
            is not None
        ):

            st.download_button(
                "Download full PDF board pack",
                data=(
                    st.session_state[
                        "pdf_bytes"
                    ]
                ),
                file_name=(
                    selected_quarter
                    .replace(
                        " ",
                        "_",
                    )
                    + "_Workiva_Full_Board_Pack.pdf"
                ),
                mime="application/pdf",
                use_container_width=True,
            )


    with download2:

        if (
            st.session_state[
                "pptx_bytes"
            ]
            is not None
        ):

            st.download_button(
                "Download full PowerPoint board pack",
                data=(
                    st.session_state[
                        "pptx_bytes"
                    ]
                ),
                file_name=(
                    selected_quarter
                    .replace(
                        " ",
                        "_",
                    )
                    + "_Workiva_Full_Board_Pack.pptx"
                ),
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                use_container_width=True,
            )


# ============================================================
# FINANCIALS TAB
# ============================================================

if active_workspace == "Financials":

    st.subheader(
        (
            f"{selected_quarter} "
            "Financial View"
        )
    )


    financial_metrics = []


    for column in [

        business[
            "revenue"
        ],

        business[
            "budget"
        ],

        business[
            "costs"
        ],
    ]:

        if (
            column
            and column
            not in financial_metrics
        ):

            financial_metrics.append(
                column
            )


    if not financial_metrics:

        st.info(
            (
                "No obvious revenue, budget "
                "or cost columns were detected."
            )
        )


    else:

        cards = (
            st.columns(
                len(
                    financial_metrics
                )
            )
        )


        for index, metric in enumerate(
            financial_metrics
        ):

            value = (
                pd.to_numeric(
                    data[
                        metric
                    ],
                    errors="coerce",
                )
                .sum()
            )


            cards[
                index
            ].metric(
                metric,
                f"{value:,.2f}",
            )


        if (
            business[
                "revenue"
            ]
            and business[
                "budget"
            ]
        ):

            revenue_total = (
                pd.to_numeric(
                    data[
                        business[
                            "revenue"
                        ]
                    ],
                    errors="coerce",
                )
                .sum()
            )


            budget_total = (
                pd.to_numeric(
                    data[
                        business[
                            "budget"
                        ]
                    ],
                    errors="coerce",
                )
                .sum()
            )


            variance = (
                revenue_total
                - budget_total
            )


            variance_pct = (

                variance
                / budget_total
                * 100

                if budget_total != 0

                else None
            )


            st.metric(
                "Revenue vs Budget variance",
                f"{variance:,.2f}",
                (
                    f"{variance_pct:+.1f}%"

                    if variance_pct
                    is not None

                    else "n/a"
                ),
            )


        if (
            business[
                "dimension"
            ]
            and financial_metrics
        ):

            metric = (

                business[
                    "revenue"
                ]

                or business[
                    "budget"
                ]

                or business[
                    "costs"
                ]
            )


            grouped = (
                grouped_data(
                    data,
                    business[
                        "dimension"
                    ],
                    metric,
                    "sum",
                )
            )


            figure = (
                px.bar(
                    grouped,
                    x=(
                        business[
                            "dimension"
                        ]
                    ),
                    y=metric,
                    title=(
                        f"{metric} "
                        f"by "
                        f"{business['dimension']}"
                    ),
                )
            )


            st.plotly_chart(
                figure,
                use_container_width=True,
                key=f"finance_dimension_{selected_quarter}_{business['dimension']}_{metric}",
            )
            render_visual_ai_action(
                {
                    "title": f"{metric} by {business['dimension']}",
                    "metric": metric,
                    "dimension": business["dimension"],
                    "aggregation": "sum",
                    "values": grouped.sort_values(metric, ascending=False).head(12).to_dict(orient="records"),
                },
                key=f"finance_dimension_{selected_quarter}_{business['dimension']}_{metric}_visual",
            )


# ============================================================
# TRENDS TAB
# ============================================================

if active_workspace == "Trends":

    st.subheader(
        (
            f"{selected_quarter} "
            "Trends"
        )
    )


    if not profile[
        "dates"
    ]:

        st.info(
            (
                "No reliable date column "
                "was detected in this sheet."
            )
        )


    elif not profile[
        "numeric"
    ]:

        st.info(
            (
                "No numeric measure was "
                "detected for trend analysis."
            )
        )


    else:

        date_column = (
            profile[
                "dates"
            ][0]
        )


        for trend_index, metric in enumerate(
            profile[
                "numeric"
            ][:3]
        ):

            trend_data = (
                grouped_data(
                    data,
                    date_column,
                    metric,
                    "sum",
                )
                .sort_values(
                    date_column
                )
            )


            figure = (
                px.line(
                    trend_data,
                    x=date_column,
                    y=metric,
                    markers=True,
                    title=(
                        f"{metric} trend"
                    ),
                )
            )


            st.plotly_chart(
                figure,
                use_container_width=True,
                key=f"trend_{selected_quarter}_{trend_index}",
            )
            trend_values = trend_data[[date_column, metric]].copy()
            trend_first = float(trend_values[metric].iloc[0]) if not trend_values.empty else None
            trend_last = float(trend_values[metric].iloc[-1]) if not trend_values.empty else None
            render_visual_ai_action(
                {
                    "title": f"{metric} trend",
                    "metric": metric,
                    "date_column": str(date_column),
                    "aggregation": "sum",
                    "values": [
                        {"period": str(row[date_column]), "value": float(row[metric])}
                        for _, row in trend_values.head(16).iterrows()
                    ],
                    "first_value": trend_first,
                    "last_value": trend_last,
                    "change": (trend_last - trend_first) if trend_first is not None and trend_last is not None else None,
                    "change_pct": ((trend_last - trend_first) / abs(trend_first) * 100) if trend_first not in (None, 0) and trend_last is not None else None,
                },
                key=f"trend_{selected_quarter}_{trend_index}_visual",
            )


        # ----------------------------------------------------
        # AUTOMATIC DRIVER ANALYSIS
        # Deterministic arithmetic first; AI only interprets it.
        # ----------------------------------------------------
        driver_analysis = (
            summary.get("analytical_context", {}).get("driver_analysis", {}) or {}
        )
        driver_metrics = driver_analysis.get("metrics", {}) or {}

        if driver_metrics:
            st.markdown("#### What drove the change?")
            st.caption(
                "Calculated decomposition of the earliest-to-latest movement. "
                "Contribution shows where the change occurred; it does not assert root cause."
            )

            driver_metric_options = list(driver_metrics.keys())
            selected_driver_metric = st.selectbox(
                "Measure",
                driver_metric_options,
                key=f"driver_metric_{selected_quarter}",
            )

            driver_dimensions = driver_metrics.get(selected_driver_metric, {}) or {}
            if driver_dimensions:
                selected_driver_dimension = st.selectbox(
                    "Driver dimension",
                    list(driver_dimensions.keys()),
                    key=f"driver_dimension_{selected_quarter}_{selected_driver_metric}",
                )
                driver_details = driver_dimensions[selected_driver_dimension]

                total_change = driver_details.get("total_change")
                total_change_pct = driver_details.get("total_change_pct")
                metric_columns = st.columns(3)
                metric_columns[0].metric(
                    "Earlier period",
                    _compact_value(driver_details.get("total_first")),
                )
                metric_columns[1].metric(
                    "Latest period",
                    _compact_value(driver_details.get("total_last")),
                )
                metric_columns[2].metric(
                    "Net change",
                    _compact_value(total_change),
                    (f"{total_change_pct:+.1f}%" if total_change_pct is not None else None),
                )

                driver_rows = []
                for item in (driver_details.get("drivers", []) or [])[:8]:
                    driver_rows.append({
                        selected_driver_dimension: item.get("category"),
                        "Earlier": item.get("first_value"),
                        "Latest": item.get("last_value"),
                        "Change": item.get("change"),
                        "Contribution to net change %": item.get("contribution_pct_of_net_change"),
                    })

                if driver_rows:
                    driver_frame = pd.DataFrame(driver_rows)
                    st.dataframe(
                        driver_frame,
                        use_container_width=True,
                        hide_index=True,
                    )

                    top_driver = driver_rows[0]
                    direction = "increased" if (top_driver.get("Change") or 0) > 0 else "decreased"
                    st.caption(
                        f"Largest absolute contributor: {top_driver[selected_driver_dimension]} "
                        f"{direction} {selected_driver_metric} by "
                        f"{_compact_value(abs(top_driver.get('Change') or 0))}."
                    )

        st.divider()
        st.markdown("#### AI trend summary")
        st.caption(
            "Ask the AI analyst to synthesize the calculated trend, anomaly, and correlation "
            "signals. All numerical claims are checked against the Python-calculated evidence."
        )

        trends_summary_key = abs(hash(json.dumps({
            "quarter": selected_quarter,
            "trends": compact_ai_context(summary).get("trends", {}),
            "anomalies": compact_ai_context(summary).get("anomalies", {}),
            "management_watchlist": summary.get("management_watchlist", []),
        }, sort_keys=True, default=str)))

        if st.session_state.get("trends_ai_summary_key") != trends_summary_key:
            st.session_state["trends_ai_summary"] = None
            st.session_state["trends_ai_summary_key"] = trends_summary_key

        if st.session_state.get("trends_ai_summary"):
            st.markdown(st.session_state["trends_ai_summary"])
            render_ai_evidence_indicator(summary, purpose="trends summary", key="trends_summary_evidence")
        elif st.button("Generate AI trend summary", key="generate_trends_ai_summary"):
            try:
                with st.spinner("Analyzing trend signals..."):
                    trends_summary = generate_trends_ai_summary(summary)
                st.session_state["trends_ai_summary"] = trends_summary
                st.markdown(trends_summary)
                render_ai_evidence_indicator(summary, purpose="trends summary", key="trends_summary_evidence_new")
            except Exception as error:
                st.warning("AI trend summary unavailable: " + friendly_error(error))


# ============================================================
# FORECASTING TAB
# ============================================================

if active_workspace == "Forecasting":

    st.subheader("Multi-quarter forecasting")
    st.caption(
        "Combine one to four historical Workiva quarters and project the selected "
        "measures forward. Forecast values use a simple quarterly trend extrapolation, "
        "with walk-forward backtesting when enough historical observations exist."
    )

    default_forecast_quarters = quarter_options[: min(4, len(quarter_options))]

    forecast_quarters = st.multiselect(
        "Historical quarters",
        quarter_options,
        default=default_forecast_quarters,
        max_selections=4,
        help="Select between 1 and 4 quarters. More history generally makes the direction more informative.",
        key="forecast_quarters",
    )

    horizon = st.slider(
        "Forecast horizon (quarters)",
        min_value=1,
        max_value=4,
        value=4,
        step=1,
        key="forecast_horizon",
    )

    scenario_sensitivity = st.slider(
        "Scenario sensitivity per quarter (%)",
        min_value=0,
        max_value=25,
        value=5,
        step=1,
        help=(
            "Applied deterministically around the baseline forecast. Upside is the numerically "
            "higher case and Downside the numerically lower case; this is not a probability."
        ),
        key="forecast_scenario_sensitivity",
    )

    # Discover common measures from the already-loaded quarter first. The full
    # cross-quarter intersection is validated when the user builds the board.
    preferred_metrics = []
    for candidate in summary.get("management_watchlist", []):
        if candidate in profile["numeric"] and candidate not in preferred_metrics:
            preferred_metrics.append(candidate)
    for candidate in [
        business.get("revenue"),
        business.get("budget"),
        business.get("costs"),
        business.get("profit"),
        business.get("ebitda"),
    ]:
        if candidate and candidate in profile["numeric"] and candidate not in preferred_metrics:
            preferred_metrics.append(candidate)
    if not preferred_metrics:
        preferred_metrics = profile["numeric"][:3]

    forecast_metrics = st.multiselect(
        "Measures to forecast",
        profile["numeric"],
        default=preferred_metrics[:3],
        max_selections=4,
        help="Choose up to four numeric measures for the forecasting board. Watchlisted measures are preferred as defaults when available.",
        key="forecast_metrics",
    )

    if st.button(
        "Build forecasting board",
        type="primary",
        key="build_forecasting_board",
    ):
        if not forecast_quarters:
            st.warning("Select at least one historical quarter.")
        elif not forecast_metrics:
            st.warning("Select at least one measure to forecast.")
        else:
            try:
                with st.spinner("Loading selected quarters and building forecast..."):
                    forecast_bundles = [quarter_bundle(item) for item in forecast_quarters]
                    result = build_forecast_result(
                        forecast_bundles,
                        forecast_metrics,
                        horizon,
                        scenario_sensitivity_pct=scenario_sensitivity,
                    )

                missing_metrics = [
                    metric for metric in forecast_metrics
                    if metric not in result["common_metrics"]
                ]
                if missing_metrics:
                    st.warning(
                        "These measures are not available as numeric fields in every selected quarter "
                        "and were excluded: " + ", ".join(map(str, missing_metrics))
                    )

                if not result["forecasts"]:
                    st.warning(
                        "No selected measure had enough compatible numeric data across the chosen quarters."
                    )
                else:
                    st.session_state["forecast_result"] = result
                    st.session_state["forecast_ai_commentary"] = None
                    st.session_state["forecast_ai_commentary_key"] = None
                    log_event(
                        "Build forecast board",
                        f"History: {', '.join(forecast_quarters)}; horizon: {horizon} quarters",
                    )
            except Exception as error:
                st.error(friendly_error(error))

    forecast_result = st.session_state.get("forecast_result")

    if forecast_result and forecast_result.get("forecasts"):
        st.divider()
        st.markdown("#### Forecast board")

        history_count = len(forecast_result.get("history", []))
        metric_count = len(forecast_result.get("forecasts", {}))
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Historical quarters", history_count)
        c2.metric("Forecast horizon", f"{forecast_result['horizon']} quarters")
        c3.metric("Measures", metric_count)
        c4.metric("Scenario sensitivity", f"±{forecast_result.get('scenario_sensitivity_pct', 0):.0f}% / qtr")

        if history_count < 2:
            st.warning(
                "Only one historical quarter is included. The board therefore uses a flat baseline, "
                "not an estimated trend."
            )
        elif history_count < 4:
            st.info(
                "The projection uses a short history. The app compares only simple models supported by "
                "the available observations and treats the result as directional planning support."
            )
        else:
            st.info(
                "Four historical quarters are included. The app selects among simple interpretable models "
                "using walk-forward validation; confidence remains capped because the history is short."
            )

        for metric_index, metric in enumerate(forecast_result["metrics"]):
            frame = forecast_result["forecasts"].get(metric)
            if frame is None or frame.empty:
                continue

            actual = frame[frame["Series"] == "Actual"]
            projected = frame[frame["Series"] == "Forecast"]
            latest_actual = float(actual["Value"].iloc[-1])
            end_forecast = float(projected["Value"].iloc[-1])
            delta = end_forecast - latest_actual
            delta_pct = (delta / abs(latest_actual) * 100) if latest_actual else None

            st.markdown(f"##### {metric}")
            backtest = (forecast_result.get("backtests", {}) or {}).get(metric, {}) or {}
            selection = (forecast_result.get("model_selection", {}) or {}).get(metric, {}) or {}
            confidence = backtest.get("confidence", "Limited")
            wape = backtest.get("wape_pct")

            model_name = selection.get("model_name", "Flat baseline")
            st.caption(f"**Selected forecast model:** {model_name}. {selection.get('reason', '')}")

            card1, card2, card3, card4 = st.columns(4)
            card1.metric("Latest actual", f"{latest_actual:,.2f}")
            card2.metric("End-of-horizon forecast", f"{end_forecast:,.2f}")
            card3.metric(
                "Projected change",
                f"{delta:+,.2f}",
                f"{delta_pct:+.1f}%" if delta_pct is not None else None,
            )
            card4.metric(
                "Forecast confidence",
                confidence,
                f"WAPE {wape:.1f}%" if wape is not None else "Backtest limited",
            )

            scenario_frame = (forecast_result.get("scenarios", {}) or {}).get(metric)
            if isinstance(scenario_frame, pd.DataFrame) and not scenario_frame.empty:
                future_scenarios = scenario_frame[scenario_frame["Scenario"] != "Actual"]
                final_period = future_scenarios["Quarter Period"].max()
                final_cases = future_scenarios[future_scenarios["Quarter Period"] == final_period]
                case_values = {
                    row["Scenario"]: float(row["Value"])
                    for _, row in final_cases.iterrows()
                }
                s1, s2, s3 = st.columns(3)
                s1.metric("Downside case", f"{case_values.get('Downside', float('nan')):,.2f}")
                s2.metric("Baseline case", f"{case_values.get('Baseline', float('nan')):,.2f}")
                s3.metric("Upside case", f"{case_values.get('Upside', float('nan')):,.2f}")

                scenario_plot = scenario_frame.copy()
                scenario_plot["Quarter order"] = scenario_plot["Quarter Period"].astype(str)
                scenario_figure = px.line(
                    scenario_plot,
                    x="Quarter order",
                    y="Value",
                    color="Scenario",
                    markers=True,
                    title=f"{metric}: baseline and sensitivity scenarios",
                    labels={"Quarter order": "Quarter", "Value": metric},
                )
                st.plotly_chart(
                    scenario_figure,
                    use_container_width=True,
                    key=f"forecast_scenario_chart_{metric_index}_{metric}",
                )
                render_visual_ai_action(
                    {
                        "title": f"{metric}: baseline and sensitivity scenarios",
                        "metric": metric,
                        "scenario_sensitivity_pct": forecast_result.get("scenario_sensitivity_pct"),
                        "forecast_confidence": confidence,
                        "selected_model": model_name,
                        "scenario_values": scenario_frame[["Quarter Period", "Scenario", "Value"]].to_dict(orient="records"),
                        "interpretation_constraint": "Scenarios are deterministic sensitivity cases, not probabilities; higher numeric values are not automatically favourable.",
                    },
                    key=f"forecast_scenario_chart_{metric_index}_{metric}_visual",
                )
                st.caption(
                    f"Scenario assumption: ±{forecast_result.get('scenario_sensitivity_pct', 0):.0f}% per forecast quarter around the baseline. "
                    "These are deterministic sensitivity cases, not probabilities. For cost/expense measures, a higher numerical case is not necessarily favourable."
                )

            chart_frame = frame.copy()
            chart_frame["Quarter order"] = chart_frame["Quarter Period"].astype(str)
            figure = px.line(
                chart_frame,
                x="Quarter order",
                y="Value",
                color="Series",
                markers=True,
                title=f"{metric}: actuals and selected-model forecast",
                labels={"Quarter order": "Quarter", "Value": metric},
            )
            st.plotly_chart(
                figure,
                use_container_width=True,
                key=f"forecast_chart_{metric_index}_{metric}",
            )
            render_visual_ai_action(
                {
                    "title": f"{metric}: actuals and selected-model forecast",
                    "metric": metric,
                    "selected_model": model_name,
                    "forecast_confidence": confidence,
                    "wape_pct": wape,
                    "history_count": history_count,
                    "values": frame[["Quarter Period", "Series", "Value"]].to_dict(orient="records"),
                    "latest_actual": latest_actual,
                    "end_forecast": end_forecast,
                    "projected_change": delta,
                    "projected_change_pct": delta_pct,
                },
                key=f"forecast_chart_{metric_index}_{metric}_visual",
            )

            aggregation_label = _forecast_aggregation(metric)
            st.caption(
                f"Quarterly basis: {aggregation_label}. Forecast method: {model_name} selected from "
                f"the simple models supported by {history_count} historical quarter(s)."
            )
            st.markdown(forecast_result["commentary"].get(metric, ""))

            confidence_reason = backtest.get("confidence_reason")
            if confidence_reason:
                st.caption(f"**Validation:** {confidence_reason}")

            tests = backtest.get("tests")
            if isinstance(tests, pd.DataFrame) and not tests.empty:
                with st.expander("Forecast validation details", expanded=False):
                    v1, v2, v3, v4 = st.columns(4)
                    v1.metric("Backtest holdouts", int(backtest.get("backtest_points") or 0))
                    v2.metric(
                        "MAE",
                        f"{float(backtest['mae']):,.2f}" if backtest.get("mae") is not None else "—",
                    )
                    v3.metric(
                        "WAPE",
                        f"{float(backtest['wape_pct']):.1f}%" if backtest.get("wape_pct") is not None else "—",
                    )
                    v4.metric(
                        "sMAPE",
                        f"{float(backtest['smape_pct']):.1f}%" if backtest.get("smape_pct") is not None else "—",
                    )

                    candidate_scores = selection.get("candidate_scores", []) or []
                    if candidate_scores:
                        model_table = pd.DataFrame(candidate_scores).copy()
                        show_cols = [
                            col for col in ["model_name", "wape_pct", "mae", "backtest_points"]
                            if col in model_table.columns
                        ]
                        if show_cols:
                            model_table = model_table[show_cols].rename(columns={
                                "model_name": "Candidate model",
                                "wape_pct": "WAPE %",
                                "mae": "MAE",
                                "backtest_points": "Holdouts",
                            })
                            st.markdown("**Candidate model comparison**")
                            st.dataframe(model_table, use_container_width=True, hide_index=True)
                            st.caption(
                                "Selection uses walk-forward error with a simplicity preference when models perform similarly."
                            )

                    validation_chart = tests[["Quarter", "Actual", "Backtest Forecast"]].copy()
                    validation_long = validation_chart.melt(
                        id_vars="Quarter",
                        value_vars=["Actual", "Backtest Forecast"],
                        var_name="Series",
                        value_name="Value",
                    )
                    validation_figure = px.line(
                        validation_long,
                        x="Quarter",
                        y="Value",
                        color="Series",
                        markers=True,
                        title=f"{metric}: historical backtest",
                    )
                    st.plotly_chart(
                        validation_figure,
                        use_container_width=True,
                        key=f"backtest_chart_{metric_index}_{metric}",
                    )
                    render_visual_ai_action(
                        {
                            "title": f"{metric}: historical backtest",
                            "metric": metric,
                            "selected_model": model_name,
                            "confidence": confidence,
                            "mae": backtest.get("mae"),
                            "wape_pct": backtest.get("wape_pct"),
                            "smape_pct": backtest.get("smape_pct"),
                            "holdouts": backtest.get("backtest_points"),
                            "tests": tests.to_dict(orient="records"),
                            "interpretation_constraint": "Historical backtest performance does not guarantee future accuracy.",
                        },
                        key=f"backtest_chart_{metric_index}_{metric}_visual",
                    )

                    display_columns = [
                        "Quarter",
                        "Actual",
                        "Backtest Forecast",
                        "Error",
                        "Absolute % Error",
                        "Training Quarters",
                    ]
                    st.dataframe(
                        tests[display_columns],
                        use_container_width=True,
                        hide_index=True,
                    )
                    st.caption(
                        "Walk-forward validation predicts each known quarter using only earlier quarters. "
                        "Historical backtest accuracy does not guarantee future forecast accuracy."
                    )

        st.divider()
        st.markdown("#### Executive forecast commentary")
        st.caption(
            "The metric-level commentary above is calculated locally. You can also ask the AI analyst "
            "to synthesize the board without inventing causes or external assumptions."
        )

        commentary_key = abs(hash(json.dumps({
            "quarters": forecast_result.get("quarters", []),
            "metrics": forecast_result.get("metrics", []),
            "horizon": forecast_result.get("horizon"),
            "scenario_sensitivity_pct": forecast_result.get("scenario_sensitivity_pct"),
        }, sort_keys=True, default=str)))

        if st.session_state.get("forecast_ai_commentary_key") != commentary_key:
            st.session_state["forecast_ai_commentary"] = None
            st.session_state["forecast_ai_commentary_key"] = commentary_key

        if st.session_state.get("forecast_ai_commentary"):
            st.markdown(st.session_state["forecast_ai_commentary"])
            render_ai_evidence_indicator(
                summary,
                purpose="forecast commentary",
                forecast_result=forecast_result,
                key="forecast_commentary_evidence",
            )
        elif st.button("Generate AI forecast commentary", key="generate_forecast_commentary"):
            try:
                with st.spinner("Analyzing forecast trends..."):
                    commentary = generate_forecast_ai_commentary(forecast_result)
                st.session_state["forecast_ai_commentary"] = commentary
                st.markdown(commentary)
                render_ai_evidence_indicator(
                    summary,
                    purpose="forecast commentary",
                    forecast_result=forecast_result,
                    key="forecast_commentary_evidence_new",
                )
            except Exception as error:
                st.warning("AI forecast commentary unavailable: " + friendly_error(error))

        export_history = forecast_result["history"].copy()
        export_history["Quarter Period"] = export_history["Quarter Period"].astype(str)
        st.download_button(
            "Download combined quarterly history as CSV",
            data=dataframe_csv_bytes(export_history),
            file_name="workiva_forecast_history.csv",
            mime="text/csv",
            key="download_forecast_history",
        )


# ============================================================
# COMPARISON TAB
# ============================================================

if active_workspace == "Comparison":

    st.subheader(
        "Quarter comparison"
    )


    quarter_list = (
        quarter_options
    )


    col1, col2 = (
        st.columns(
            2
        )
    )


    first_index = (

        quarter_list.index(
            selected_quarter
        )

        if selected_quarter
        in quarter_list

        else 0
    )


    second_default = min(
        first_index + 1,
        len(
            quarter_list
        ) - 1,
    )


    comparison_first = (
        col1.selectbox(
            "First quarter",
            quarter_list,
            index=first_index,
            key="comparison_first",
        )
    )


    comparison_second = (
        col2.selectbox(
            "Second quarter",
            quarter_list,
            index=second_default,
            key="comparison_second",
        )
    )


    comparison_focus = (
        st.text_input(
            "Optional comparison focus",
            placeholder=(
                "Example: focus on costs "
                "and regional performance"
            ),
        )
    )


    if st.button(
        "Compare quarters",
        key="compare_quarters_button",
    ):

        if (
            comparison_first
            == comparison_second
        ):

            st.warning(
                (
                    "Choose two different "
                    "quarters."
                )
            )


        else:

            try:

                with st.spinner(
                    (
                        "Loading and comparing "
                        "both Workiva quarters..."
                    )
                ):

                    first_bundle = (
                        quarter_bundle(
                            comparison_first
                        )
                    )


                    second_bundle = (
                        quarter_bundle(
                            comparison_second
                        )
                    )


                    comparison = (
                        compare_quarter_bundles(
                            first_bundle,
                            second_bundle,
                            comparison_focus,
                        )
                    )


                st.session_state[
                    "comparison_result"
                ] = comparison


                log_event(
                    "Compare quarters",
                    (
                        f"{comparison_first} "
                        f"vs "
                        f"{comparison_second}"
                    ),
                )


            except Exception as error:

                st.error(
                    friendly_error(
                        error
                    )
                )


    if (
        st.session_state[
            "comparison_result"
        ]
    ):

        render_comparison(
            st.session_state[
                "comparison_result"
            ],
            key_prefix="comparison_tab",
        )


# ============================================================
# DATA QUALITY TAB
# ============================================================

if active_workspace == "Data quality":

    st.subheader(
        "Data Quality Dashboard"
    )

    st.caption(
        "Deterministic quality checks run before AI interpretation. The readiness score is an "
        "analytical-use heuristic, not an audit opinion or accounting control certification."
    )

    quality_intel = summary.get("data_quality_intelligence", {}) or {}
    q1, q2, q3, q4 = st.columns(4)
    q1.metric("Readiness", quality_intel.get("status", "n/a"))
    q2.metric("Quality score", f"{quality_intel.get('score', 0):.0f}/100" if quality_intel.get("score") is not None else "n/a")
    q3.metric("Rows", len(data))
    q4.metric("Missing cells", int(data.isna().sum().sum()))

    check_rows = [
        {
            "Check": item.get("check"),
            "Status": item.get("status"),
            "Finding": item.get("detail"),
        }
        for item in (quality_intel.get("checks", []) or [])
    ]
    if check_rows:
        st.dataframe(pd.DataFrame(check_rows), use_container_width=True, hide_index=True)

    st.markdown("#### Cross-quarter consistency")
    st.caption(
        "Optional deeper check across up to four quarters. It checks row-count drift, schema changes, "
        "detected-type changes and category membership drift. It runs only when requested so normal page loads stay fast."
    )
    quality_default = quarter_options[: min(4, len(quarter_options))]
    scan_quarters = st.multiselect(
        "Quarters to validate",
        quarter_options,
        default=quality_default,
        max_selections=4,
        key="quality_scan_quarters",
    )
    scan_key = json.dumps(scan_quarters, sort_keys=True)
    if st.session_state.get("cross_quarter_quality_key") != scan_key:
        st.session_state["cross_quarter_quality"] = None
        st.session_state["cross_quarter_quality_key"] = scan_key

    if st.button("Run cross-quarter quality scan", key="run_cross_quarter_quality"):
        if len(scan_quarters) < 2:
            st.info("Select at least two quarters for a cross-quarter consistency scan.")
        else:
            try:
                with st.spinner("Checking cross-quarter consistency..."):
                    scan_bundles = [quarter_bundle(item) for item in scan_quarters[:4]]
                    st.session_state["cross_quarter_quality"] = compare_quarter_quality(scan_bundles)
                log_event("Run cross-quarter quality scan", ", ".join(scan_quarters[:4]))
            except Exception as error:
                st.warning("Cross-quarter quality scan unavailable: " + friendly_error(error))

    cross_quality = st.session_state.get("cross_quarter_quality")
    if cross_quality:
        st.markdown(f"**Cross-quarter status: {cross_quality.get('status')} · {cross_quality.get('score', 0):.0f}/100**")
        cross_rows = [
            {"Check": item.get("check"), "Status": item.get("status"), "Finding": item.get("detail")}
            for item in (cross_quality.get("checks", []) or [])
        ]
        if cross_rows:
            st.dataframe(pd.DataFrame(cross_rows), use_container_width=True, hide_index=True)

    st.divider()
    st.markdown("#### Column profile")

    quality_rows = []


    for column in data.columns:

        detail = (
            profile[
                "details"
            ].get(
                column,
                {},
            )
        )


        quality_rows.append(
            {

                "Column":
                    column,

                "Detected type":
                    detail.get(
                        "type",
                        "unknown",
                    ),

                "Missing":
                    int(
                        data[
                            column
                        ]
                        .isna()
                        .sum()
                    ),

                "Unique":
                    int(
                        data[
                            column
                        ]
                        .nunique(
                            dropna=True
                        )
                    ),
            }
        )


    st.dataframe(
        pd.DataFrame(
            quality_rows
        ),
        use_container_width=True,
    )


    with st.expander(
        "What the agent detected"
    ):

        st.write(
            "Numeric:",
            profile[
                "numeric"
            ],
        )

        st.write(
            "Categories:",
            profile[
                "categories"
            ],
        )

        st.write(
            "Dates:",
            profile[
                "dates"
            ],
        )

        st.write(
            "Revenue:",
            business[
                "revenue"
            ],
        )

        st.write(
            "Budget:",
            business[
                "budget"
            ],
        )

        st.write(
            "Costs:",
            business[
                "costs"
            ],
        )

        st.write(
            "Main dimension:",
            business[
                "dimension"
            ],
        )


        st.markdown("**Financial intelligence**")
        finance_labels = {
            "actual": "Actual",
            "forecast": "Forecast",
            "profit": "Profit",
            "ebitda": "EBITDA",
            "opex": "Opex",
            "capex": "Capex",
            "headcount": "Headcount / FTE",
            "volume": "Volume",
            "price": "Price",
            "prior_year": "Prior year",
            "margin": "Margin",
        }
        for key, label in finance_labels.items():
            if business.get(key):
                st.write(f"{label}:", business.get(key))

        derived_measures = business.get("derived_measures", {})
        if derived_measures:
            st.write("Derived measures:")
            for measure, formula in derived_measures.items():
                st.caption(f"{measure} = {formula}")
        else:
            st.caption(
                "No additional finance ratios were derived because the required source columns were not all present."
            )


    with st.expander(
        "Raw data explorer"
    ):

        st.dataframe(
            data,
            use_container_width=True,
        )


# ============================================================
# ASK AI TAB
# ============================================================

if active_workspace == "Ask AI":

    st.subheader(
        "Ask the Workiva Agent"
    )


    st.caption(
        (
            "Ask about the current dashboard, trends, forecasting, or quarter comparisons. "
            "The assistant receives the active calculated evidence and distinguishes actuals "
            "from forecasts."
        )
    )


    for message in (
        st.session_state[
            "chat_messages"
        ]
    ):

        with st.chat_message(
            message[
                "role"
            ]
        ):

            st.markdown(
                message[
                    "content"
                ]
            )


    # Context-aware shortcuts are generated from the current calculated state.
    # They remain local UI actions and do not call Copilot until clicked.
    previous_quarter = None
    try:
        current_index = quarter_options.index(selected_quarter)
        if current_index + 1 < len(quarter_options):
            previous_quarter = quarter_options[current_index + 1]
    except (ValueError, AttributeError):
        previous_quarter = None

    contextual_questions = build_contextual_ai_questions(
        summary,
        st.session_state.get("forecast_result"),
        previous_quarter=previous_quarter,
        selected_quarter=selected_quarter,
        comparison_result=st.session_state.get("comparison_result"),
        chat_history=st.session_state.get("chat_messages", []),
    )

    if contextual_questions:
        st.caption("Suggested next questions · ranked from the current evidence")
        shortcut_columns = st.columns(len(contextual_questions))
    else:
        shortcut_columns = []

    suggested_question = None
    for shortcut_index, item in enumerate(contextual_questions):
        if shortcut_columns[shortcut_index].button(
            item["label"],
            use_container_width=True,
            help=item.get("reason"),
            key=f"chat_contextual_{shortcut_index}_{selected_quarter}_{item['topic']}",
        ):
            suggested_question = item["prompt"]

    typed_question = st.chat_input(
        "Ask about Workiva data..."
    )

    # A button click and a typed question share exactly the same execution
    # path, avoiding duplicate AI calls and keeping interaction fast.
    question = typed_question or suggested_question


    if question:

        st.session_state[
            "chat_messages"
        ].append(
            {
                "role":
                    "user",

                "content":
                    question,
            }
        )


        with st.chat_message(
            "user"
        ):

            st.markdown(
                question
            )


        with st.chat_message(
            "assistant"
        ):

            requested_quarters = (
                extract_quarters(
                    question
                )
            )


            if (
                len(
                    requested_quarters
                )
                >= 2
            ):

                first_quarter = (
                    requested_quarters[
                        0
                    ]
                )

                second_quarter = (
                    requested_quarters[
                        1
                    ]
                )


                try:

                    with st.spinner(
                        (
                            "Finding both quarters "
                            "in Workiva and comparing them..."
                        )
                    ):

                        first_bundle = (
                            quarter_bundle(
                                first_quarter
                            )
                        )


                        second_bundle = (
                            quarter_bundle(
                                second_quarter
                            )
                        )


                        comparison = (
                            compare_quarter_bundles(
                                first_bundle,
                                second_bundle,
                                question,
                            )
                        )


                    st.session_state[
                        "comparison_result"
                    ] = comparison


                    render_comparison(
                        comparison,
                        key_prefix="chat_comparison",
                    )


                    answer = (
                        f"I compared "
                        f"{first_quarter} "
                        f"and "
                        f"{second_quarter} "
                        "using Workiva data."
                    )


                    log_event(
                        "AI comparison request",
                        question,
                    )


                except Exception as error:

                    answer = (
                        friendly_error(
                            error
                        )
                    )


                    st.error(
                        answer
                    )


            else:

                try:

                    with st.spinner(
                        "Analysing..."
                    ):

                        answer = (
                            chat_answer(
                                question,
                                summary,
                                (display_plan or active_plan or summary_plan),
                                data,
                                profile,
                                business,
                                st.session_state["chat_messages"],
                                st.session_state.get("forecast_result"),
                                st.session_state.get("comparison_result"),
                            )
                        )


                    st.markdown(
                        answer
                    )
                    chat_purpose = (
                        "forecast chat"
                        if any(word in question.lower() for word in ["forecast", "scenario", "projection", "backtest", "model"])
                        else "analytical chat"
                    )
                    render_ai_evidence_indicator(
                        summary,
                        purpose=chat_purpose,
                        forecast_result=st.session_state.get("forecast_result"),
                        comparison_result=st.session_state.get("comparison_result"),
                        key="ask_ai_evidence",
                        compact=True,
                    )


                    log_event(
                        "Ask AI",
                        question,
                    )


                except Exception as error:

                    answer = (
                        friendly_error(
                            error
                        )
                    )


                    st.error(
                        answer
                    )


        st.session_state[
            "chat_messages"
        ].append(
            {

                "role":
                    "assistant",

                "content":
                    answer,
            }
        )


        # Refresh after storing the response so the user can immediately
        # ask another question while retaining the full conversation history.
        st.rerun()


# ============================================================
# AUDIT LOG
# ============================================================

# Always available as a collapsed control so governance evidence is easy to
# reach without adding visual clutter to the main analytical workspace.
st.divider()

with st.expander(
    "Session audit log",
    expanded=False,
):

    st.caption(
        (
            "Session-level activity log. "
            "Workiva write actions are disabled."
        )
    )


    if not st.session_state[
        "audit_log"
    ]:

        st.write(
            "No activity recorded yet."
        )


    else:

        st.dataframe(
            pd.DataFrame(
                st.session_state[
                    "audit_log"
                ]
            ),
            use_container_width=True,
        )
